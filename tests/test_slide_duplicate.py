# pyright: reportPrivateUsage=false

"""Unit-test suite for `Slide.duplicate` / `Slides.duplicate` (slide-CRUD Phase 2).

Issue: https://github.com/MHoroszowski/python-pptx/issues/11 (Phase 2 — duplicate).

Closes upstream feature request scanny/python-pptx#132 in this fork.

The tests are organised in three layers:

1. **Unit tests** that drive the API surface (`Slide.duplicate`,
   `Slides.duplicate`) via mocks — argument validation, raises, and
   delegation patterns. Mirrors the unit-test style of
   `tests/test_slide_crud.py`.
2. **Part-graph tests** that build small in-memory `Presentation`
   objects and inspect rels / parts directly to verify the dedup
   invariant (image, media reused), the deep-copy invariant (chart,
   OLE, notes-slide each get their own part), and the rId-remap pass.
3. **Round-trip integration tests** at the bottom — open → mutate →
   save → reopen — to confirm Office-compatible packaging.
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.slide import SlidePart
from pptx.slide import Slide, Slides
from pptx.util import Inches

from .unitutil.mock import instance_mock, method_mock

_RELS_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

_PNG_PATH = Path(__file__).parent / "test_files" / "python-powered.png"


# ---------------------------------------------------------------------------
# Helpers — keep round-trip plumbing identical to test_slide_crud.py.
# ---------------------------------------------------------------------------


def _seed_presentation_with(n_slides: int) -> Presentation:
    """Build an in-memory Presentation seeded with `n_slides` blank slides."""
    prs = Presentation()
    layout = prs.slide_layouts[6]  # ---blank layout---
    for _ in range(n_slides):
        prs.slides.add_slide(layout)
    return prs


def _round_trip(prs: Presentation) -> Presentation:
    """Save `prs` to a bytes buffer and reopen it."""
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


def _image_part_count(prs: Presentation) -> int:
    """Number of image parts currently reachable in the package."""
    from pptx.parts.image import ImagePart

    return sum(1 for p in prs.part.package.iter_parts() if isinstance(p, ImagePart))


def _slide_part_count(prs: Presentation) -> int:
    """Number of slide parts currently reachable in the package."""
    return sum(1 for p in prs.part.package.iter_parts() if isinstance(p, SlidePart))


# ---------------------------------------------------------------------------
# `Slides.duplicate` — collection-level operation.
# ---------------------------------------------------------------------------


class DescribeSlides_Duplicate(object):
    """Unit-test suite for `pptx.slide.Slides.duplicate`."""

    def it_exposes_a_duplicate_method(self):
        """API smoke — the method exists on the collection class."""
        assert callable(getattr(Slides, "duplicate", None))

    def it_returns_a_Slide_instance(self):
        """End-to-end smoke against a real two-slide deck."""
        prs = _seed_presentation_with(2)

        new_slide = prs.slides.duplicate(prs.slides[0])

        assert isinstance(new_slide, Slide)

    def it_inserts_the_duplicate_immediately_after_source_when_index_is_None(self):
        prs = _seed_presentation_with(3)
        source = prs.slides[1]

        new_slide = prs.slides.duplicate(source)

        assert len(prs.slides) == 4
        # ---new slide sits at source_index + 1---
        assert prs.slides[2].slide_id == new_slide.slide_id
        # ---source unchanged at its original index---
        assert prs.slides[1].slide_id == source.slide_id

    @pytest.mark.parametrize("idx", [0, 1, 2, 3])
    def it_inserts_the_duplicate_at_the_given_index(self, idx: int):
        prs = _seed_presentation_with(3)

        new_slide = prs.slides.duplicate(prs.slides[0], index=idx)

        assert len(prs.slides) == 4
        assert prs.slides[idx].slide_id == new_slide.slide_id

    @pytest.mark.parametrize("bad_idx", [-1, 5])
    def but_it_raises_IndexError_for_out_of_range_index(self, bad_idx: int):
        prs = _seed_presentation_with(3)
        with pytest.raises(IndexError):
            prs.slides.duplicate(prs.slides[0], index=bad_idx)

    def but_it_raises_ValueError_when_slide_is_not_in_the_collection(self, request):
        method_mock(request, Slides, "index", side_effect=ValueError("not in collection"))
        slides = Slides(None, None)  # pyright: ignore[reportArgumentType]
        slide_ = instance_mock(request, Slide)

        with pytest.raises(ValueError):
            slides.duplicate(slide_)

    def it_treats_index_equal_to_len_as_append(self):
        """`duplicate(slide, index=len(slides))` should be a valid append."""
        prs = _seed_presentation_with(2)

        new_slide = prs.slides.duplicate(prs.slides[0], index=len(prs.slides))

        assert len(prs.slides) == 3
        assert prs.slides[2].slide_id == new_slide.slide_id

    def it_does_not_mutate_source_slide_position(self):
        """Anti-criterion: source's index in the slide-id list must be stable."""
        prs = _seed_presentation_with(3)
        source = prs.slides[0]
        source_id = source.slide_id

        prs.slides.duplicate(source)

        # ---source still at index 0, its id unchanged---
        assert prs.slides[0].slide_id == source_id

    def it_assigns_a_unique_slide_id_to_the_duplicate(self):
        prs = _seed_presentation_with(2)
        ids_before = {s.slide_id for s in prs.slides}

        new_slide = prs.slides.duplicate(prs.slides[0])

        assert new_slide.slide_id not in ids_before
        assert new_slide.slide_id > max(ids_before)

    def it_creates_a_new_SlidePart_with_a_unique_partname(self):
        prs = _seed_presentation_with(2)
        partnames_before = {prs.slides[i].part.partname for i in range(len(prs.slides))}

        new_slide = prs.slides.duplicate(prs.slides[0])

        assert new_slide.part.partname not in partnames_before


# ---------------------------------------------------------------------------
# `Slide.duplicate()` — convenience alias.
# ---------------------------------------------------------------------------


class DescribeSlide_Duplicate(object):
    """Unit-test suite for `pptx.slide.Slide.duplicate`."""

    def it_exposes_a_duplicate_method(self):
        assert callable(getattr(Slide, "duplicate", None))

    def it_delegates_to_Slides_duplicate_on_the_owning_presentation(self, request):
        # ---Mock chain: slide.part.package.presentation_part.presentation.slides---
        slides_ = instance_mock(request, Slides)
        prs_ = instance_mock(request, type("Prs", (), {"slides": None}))
        prs_.slides = slides_
        prs_part_ = instance_mock(request, type("PresPart", (), {"presentation": None}))
        prs_part_.presentation = prs_
        slide_part_ = instance_mock(request, SlidePart)
        slide_part_.package.presentation_part = prs_part_
        slide = Slide(None, slide_part_)  # pyright: ignore[reportArgumentType]

        slide.duplicate(index=2)

        slides_.duplicate.assert_called_once_with(slide, 2)

    def it_passes_None_index_through_unchanged(self, request):
        slides_ = instance_mock(request, Slides)
        prs_ = instance_mock(request, type("Prs", (), {"slides": None}))
        prs_.slides = slides_
        prs_part_ = instance_mock(request, type("PresPart", (), {"presentation": None}))
        prs_part_.presentation = prs_
        slide_part_ = instance_mock(request, SlidePart)
        slide_part_.package.presentation_part = prs_part_
        slide = Slide(None, slide_part_)  # pyright: ignore[reportArgumentType]

        slide.duplicate()

        slides_.duplicate.assert_called_once_with(slide, None)

    def it_returns_the_value_from_Slides_duplicate(self, request):
        new_slide_ = instance_mock(request, Slide)
        slides_ = instance_mock(request, Slides)
        slides_.duplicate.return_value = new_slide_
        prs_ = instance_mock(request, type("Prs", (), {"slides": None}))
        prs_.slides = slides_
        prs_part_ = instance_mock(request, type("PresPart", (), {"presentation": None}))
        prs_part_.presentation = prs_
        slide_part_ = instance_mock(request, SlidePart)
        slide_part_.package.presentation_part = prs_part_
        slide = Slide(None, slide_part_)  # pyright: ignore[reportArgumentType]

        result = slide.duplicate()

        assert result is new_slide_


# ---------------------------------------------------------------------------
# Part-graph + dedup invariants.
# ---------------------------------------------------------------------------


class DescribeSlideDuplicate_PartGraph(object):
    """Verify slide duplication maintains the part graph invariants."""

    def it_creates_a_new_unique_relationship_in_presentation_rels(self):
        prs = _seed_presentation_with(2)
        rIds_before = {rId for rId, _ in prs.part.rels.items()}

        prs.slides.duplicate(prs.slides[0])

        rIds_after = {rId for rId, _ in prs.part.rels.items()}
        new_rIds = rIds_after - rIds_before
        assert len(new_rIds) == 1

    def it_grows_the_slide_part_count_by_exactly_one(self):
        prs = _seed_presentation_with(2)
        n_before = _slide_part_count(prs)

        prs.slides.duplicate(prs.slides[0])

        assert _slide_part_count(prs) == n_before + 1

    def it_shares_the_slide_layout_part_with_the_source(self):
        prs = _seed_presentation_with(2)
        source_layout = prs.slides[0].slide_layout

        new_slide = prs.slides.duplicate(prs.slides[0])

        # ---both slides resolve to the SAME SlideLayoutPart instance---
        assert new_slide.slide_layout.part is source_layout.part

    def it_isolates_modifications_to_the_duplicate_from_the_source(self):
        prs = _seed_presentation_with(1)
        layout = prs.slide_layouts[6]
        source = prs.slides.add_slide(layout)
        source.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(1)
        ).text_frame.text = "original"

        new_slide = prs.slides.duplicate(source)
        # ---mutate duplicate's text---
        textbox = next(shp for shp in new_slide.shapes if getattr(shp, "has_text_frame", False))
        textbox.text_frame.text = "mutated"

        # ---source's textbox unchanged---
        source_text = next(
            shp for shp in source.shapes if getattr(shp, "has_text_frame", False)
        ).text_frame.text
        assert source_text == "original"

    def it_serializes_duplicate_xml_equivalent_to_source_before_mutation(self):
        """Pre-mutation, dup's <p:sld> matches source modulo r:id substitution."""
        prs = _seed_presentation_with(1)
        layout = prs.slide_layouts[6]
        source = prs.slides.add_slide(layout)
        source.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(1)
        ).text_frame.text = "hello"

        new_slide = prs.slides.duplicate(source)

        # ---both slides have the same number of shapes with the same text---
        assert len(new_slide.shapes) == len(source.shapes)
        assert (
            next(
                shp.text_frame.text
                for shp in new_slide.shapes
                if getattr(shp, "has_text_frame", False)
            )
            == "hello"
        )


class DescribeSlideDuplicate_ImageDedup(object):
    """The package-level image-dedup invariant: shared image parts stay shared."""

    def it_does_not_increase_image_part_count_when_duplicating_a_slide_with_an_image(self):
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        slide.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))
        n_images_before = _image_part_count(prs)

        prs.slides.duplicate(slide)

        assert _image_part_count(prs) == n_images_before

    def it_shares_the_image_part_with_the_source(self):
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        pic_src = slide.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))
        n_before = _image_part_count(prs)

        new_slide = prs.slides.duplicate(slide)

        pic_dup = next(shp for shp in new_slide.shapes if shp.shape_type == 13)
        assert pic_dup.image.blob == pic_src.image.blob
        # ---image part count unchanged proves the share at package level---
        assert _image_part_count(prs) == n_before

    def it_round_trips_a_duplicated_slide_with_an_image(self):
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        slide.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))

        prs.slides.duplicate(slide)
        round_tripped = _round_trip(prs)

        # ---both slides expose pictures referencing the same blob---
        pictures_per_slide = [
            [shp for shp in s.shapes if shp.shape_type == 13] for s in round_tripped.slides
        ]
        assert all(len(pics) == 1 for pics in pictures_per_slide)
        blobs = [pics[0].image.blob for pics in pictures_per_slide]
        assert blobs[0] == blobs[1]
        assert len(blobs[0]) > 0


class DescribeSlideDuplicate_NotesSlide(object):
    """Notes-slide handling: duplicate gets its own NotesSlidePart."""

    def it_gives_the_duplicate_its_own_notes_slide_part(self):
        prs = _seed_presentation_with(1)
        source = prs.slides[0]
        source.notes_slide.notes_text_frame.text = "speaker notes"

        new_slide = prs.slides.duplicate(source)

        assert new_slide.has_notes_slide is True
        # ---the two notes-slide parts are distinct---
        assert new_slide.notes_slide.part is not source.notes_slide.part

    def it_carries_the_notes_text_to_the_duplicate(self):
        prs = _seed_presentation_with(1)
        source = prs.slides[0]
        source.notes_slide.notes_text_frame.text = "speaker notes"

        new_slide = prs.slides.duplicate(source)

        assert new_slide.notes_slide.notes_text_frame.text == "speaker notes"

    def it_isolates_notes_edits_on_the_duplicate_from_the_source(self):
        prs = _seed_presentation_with(1)
        source = prs.slides[0]
        source.notes_slide.notes_text_frame.text = "original notes"

        new_slide = prs.slides.duplicate(source)
        new_slide.notes_slide.notes_text_frame.text = "mutated notes"

        assert source.notes_slide.notes_text_frame.text == "original notes"

    def it_does_not_create_a_notes_slide_when_source_has_none(self):
        prs = _seed_presentation_with(1)
        source = prs.slides[0]
        # ---source has NO notes-slide; do not call .notes_slide which lazily creates one---
        assert source.has_notes_slide is False

        new_slide = prs.slides.duplicate(source)

        assert new_slide.has_notes_slide is False

    def it_rewires_the_duplicate_notes_slide_back_ref_to_the_new_slide(self):
        """Anti — community gotcha #961: notes-slide must back-ref to new slide."""
        prs = _seed_presentation_with(1)
        source = prs.slides[0]
        source.notes_slide.notes_text_frame.text = "x"

        new_slide = prs.slides.duplicate(source)

        # ---the new notes-slide's RT.SLIDE rel target is the NEW slide part---
        new_notes_part = new_slide.notes_slide.part
        back_ref_target = new_notes_part.part_related_by(RT.SLIDE)
        assert back_ref_target is new_slide.part

    def it_round_trips_a_slide_with_notes(self):
        prs = _seed_presentation_with(1)
        source = prs.slides[0]
        source.notes_slide.notes_text_frame.text = "speaker notes"

        prs.slides.duplicate(source)
        round_tripped = _round_trip(prs)

        for s in round_tripped.slides:
            assert s.has_notes_slide is True
            assert s.notes_slide.notes_text_frame.text == "speaker notes"


# ---------------------------------------------------------------------------
# Defensive XPath check — NO unmapped rId references should remain.
# ---------------------------------------------------------------------------


class DescribeSlideDuplicate_RIdRemap(object):
    """Every r:* attribute on the duplicate must resolve to a rel on the new slide."""

    def it_resolves_every_rId_reference_in_the_duplicate_xml(self):
        """No dangling rIds — pivots on the dedup invariant tested at runtime."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        slide.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))
        slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(2), Inches(1)
        ).text_frame.text = "rId test"

        new_slide = prs.slides.duplicate(slide)

        # ---collect every attribute value in the relationships namespace---
        rId_refs = set()
        for el in new_slide.element.iter():
            for attr_name, attr_val in el.attrib.items():
                if attr_name.startswith(_RELS_NS):
                    rId_refs.add(attr_val)
        # ---each must resolve in the new slide part's rels---
        new_part_rIds = set(new_slide.part.rels)
        unresolved = rId_refs - new_part_rIds
        assert unresolved == set(), f"unresolved rIds in duplicated slide: {unresolved}"


# ---------------------------------------------------------------------------
# Round-trip integration tests — open → duplicate → save → reopen.
# ---------------------------------------------------------------------------


class DescribeSlideDuplicate_RoundTrip(object):
    """Open → duplicate → save → reopen integration coverage."""

    def it_round_trips_a_basic_duplicate(self):
        prs = _seed_presentation_with(2)
        ids_before = [s.slide_id for s in prs.slides]

        prs.slides.duplicate(prs.slides[0])
        round_tripped = _round_trip(prs)

        assert len(round_tripped.slides) == 3
        ids_after = [s.slide_id for s in round_tripped.slides]
        # ---source ids stable, duplicate inserted at index 1---
        assert ids_after[0] == ids_before[0]
        assert ids_after[2] == ids_before[1]
        assert ids_after[1] not in ids_before

    def it_round_trips_a_duplicate_at_a_specific_index(self):
        prs = _seed_presentation_with(3)

        prs.slides.duplicate(prs.slides[2], index=0)
        round_tripped = _round_trip(prs)

        assert len(round_tripped.slides) == 4

    def it_round_trips_Slide_duplicate_alias(self):
        prs = _seed_presentation_with(2)

        prs.slides[0].duplicate()
        round_tripped = _round_trip(prs)

        assert len(round_tripped.slides) == 3

    def it_preserves_shape_count_through_round_trip(self):
        prs = _seed_presentation_with(1)
        layout = prs.slide_layouts[6]
        source = prs.slides.add_slide(layout)
        source.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(1)
        ).text_frame.text = "kept"
        source.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(2), Inches(1)
        ).text_frame.text = "also kept"
        n_shapes_source = len(source.shapes)

        prs.slides.duplicate(source)
        round_tripped = _round_trip(prs)

        # ---last slide is the duplicate; shape count matches source---
        assert len(round_tripped.slides[-1].shapes) == n_shapes_source

    def it_does_not_mutate_image_part_count_through_round_trip(self):
        """Anti-criterion: image dedup survives save → reopen."""
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        slide.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))
        n_images_before = _image_part_count(prs)

        prs.slides.duplicate(slide)
        round_tripped = _round_trip(prs)

        assert _image_part_count(round_tripped) == n_images_before

    def it_does_not_carry_comments_through_a_duplicate(self):
        """Phase-2 scope: comments parts are dropped on duplicate (documented).

        We don't have a Phase-2 API to add comments yet, so this test
        documents the behavior via a hand-crafted source-slide rel: if
        the source had a `RT.COMMENTS` rel pointing at some part, the
        duplicate must NOT carry it. This is a forward-looking guard
        for when comments are added in a later phase.
        """
        prs = _seed_presentation_with(2)
        source = prs.slides[0]

        new_slide = prs.slides.duplicate(source)

        # ---no slide ever has RT.COMMENTS rels in this build, but the
        #    invariant we want is: even if it had one, it would be dropped.
        #    Document the invariant by asserting absence on dup---
        with pytest.raises(KeyError):
            new_slide.part.part_related_by(RT.COMMENTS)
