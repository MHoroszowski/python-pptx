# pyright: reportPrivateUsage=false

"""Unit-test suite for Tables 2.0 Phase 4 — sizing & ergonomics, closing the epic.

Covers:

- Read-only count properties on |Table|: ``row_count``, ``column_count``,
  ``dimensions`` — convenience accessors that don't instantiate the
  full |_RowCollection| / |_ColumnCollection|.
- Per-row height and per-column width round-trip preservation through
  save/reload (regression-lock — the underlying setters already work,
  but no test pinned that down before now).
- Anti-criteria: Phase 1/2/3 surfaces unaffected.

Issue: https://github.com/MHoroszowski/python-pptx/issues/12 (Phase 4, closing PR).
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.util import Emu, Inches

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


def _make_table(rows: int, cols: int):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    gf = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(2))
    return prs, gf.table


@pytest.fixture
def t3x4():
    _, t = _make_table(3, 4)
    return t


# ---------------------------------------------------------------------------
# Table.row_count / column_count / dimensions
# ---------------------------------------------------------------------------


class DescribeTable_CountProperties(object):
    """Unit-test suite for `row_count`, `column_count`, `dimensions`."""

    def it_reports_row_count_for_a_3x4_table(self, t3x4):
        assert t3x4.row_count == 3

    def it_reports_column_count_for_a_3x4_table(self, t3x4):
        assert t3x4.column_count == 4

    def it_reports_dimensions_as_rows_cols_tuple(self, t3x4):
        assert t3x4.dimensions == (3, 4)

    def it_matches_len_of_rows_and_columns_collections(self, t3x4):
        assert t3x4.row_count == len(t3x4.rows)
        assert t3x4.column_count == len(t3x4.columns)

    def it_increments_row_count_after_rows_add(self, t3x4):
        t3x4.rows.add()
        assert t3x4.row_count == 4
        assert t3x4.dimensions == (4, 4)

    def it_decrements_row_count_after_rows_remove(self, t3x4):
        t3x4.rows.remove(1)
        assert t3x4.row_count == 2
        assert t3x4.dimensions == (2, 4)

    def it_increments_column_count_after_columns_add(self, t3x4):
        t3x4.columns.add()
        assert t3x4.column_count == 5
        assert t3x4.dimensions == (3, 5)

    def it_decrements_column_count_after_columns_remove(self, t3x4):
        t3x4.columns.remove(0)
        assert t3x4.column_count == 3
        assert t3x4.dimensions == (3, 3)

    @pytest.mark.parametrize("attr", ["row_count", "column_count", "dimensions"])
    def it_is_read_only_no_setter(self, t3x4, attr):
        with pytest.raises(AttributeError):
            setattr(t3x4, attr, 99)


# ---------------------------------------------------------------------------
# Per-row height / per-column width — round-trip regression-lock
# ---------------------------------------------------------------------------


class DescribeSizing_RoundTrip(object):
    """Save → reload preserves explicit row heights and column widths."""

    def it_preserves_a_single_row_height_through_save_and_reload(self):
        prs, t = _make_table(3, 3)
        t.rows[0].height = Emu(500000)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        prs2 = Presentation(buf)
        t2 = next(shp for shp in prs2.slides[0].shapes if shp.has_table).table
        assert t2.rows[0].height == 500000

    def it_preserves_a_single_column_width_through_save_and_reload(self):
        prs, t = _make_table(3, 3)
        t.columns[1].width = Emu(1000000)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        prs2 = Presentation(buf)
        t2 = next(shp for shp in prs2.slides[0].shapes if shp.has_table).table
        assert t2.columns[1].width == 1000000

    def it_preserves_mixed_row_heights(self):
        prs, t = _make_table(3, 3)
        heights = [Emu(300000), Emu(600000), Emu(900000)]
        for idx, h in enumerate(heights):
            t.rows[idx].height = h

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        prs2 = Presentation(buf)
        t2 = next(shp for shp in prs2.slides[0].shapes if shp.has_table).table
        assert [t2.rows[i].height for i in range(3)] == [300000, 600000, 900000]

    def it_preserves_mixed_column_widths(self):
        prs, t = _make_table(2, 4)
        widths = [Emu(400000), Emu(800000), Emu(1200000), Emu(1600000)]
        for idx, w in enumerate(widths):
            t.columns[idx].width = w

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        prs2 = Presentation(buf)
        t2 = next(shp for shp in prs2.slides[0].shapes if shp.has_table).table
        assert [t2.columns[i].width for i in range(4)] == [400000, 800000, 1200000, 1600000]

    def it_preserves_count_properties_through_round_trip(self):
        prs, t = _make_table(4, 5)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        prs2 = Presentation(buf)
        t2 = next(shp for shp in prs2.slides[0].shapes if shp.has_table).table
        assert t2.row_count == 4
        assert t2.column_count == 5
        assert t2.dimensions == (4, 5)


# ---------------------------------------------------------------------------
# Anti / Regression
# ---------------------------------------------------------------------------


class DescribePhase4_Regression(object):
    """Anti-criteria: existing surfaces unaffected by Phase 4 additions."""

    def it_keeps_phase1_rows_add_remove_working(self, t3x4):
        t3x4.rows.add()
        assert t3x4.row_count == 4
        t3x4.rows.remove(0)
        assert t3x4.row_count == 3

    def it_keeps_phase1_columns_add_remove_working(self, t3x4):
        t3x4.columns.add()
        assert t3x4.column_count == 5
        t3x4.columns.remove(0)
        assert t3x4.column_count == 4

    def it_keeps_phase2_style_api_working(self, t3x4):
        assert t3x4.style_id == "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"
        t3x4.apply_style("No Style, No Grid")
        assert t3x4.style_name == "No Style, No Grid"

    def it_keeps_phase3_merge_api_working(self, t3x4):
        t3x4.merge_cells((0, 1), (0, 2))
        assert t3x4.cell(0, 0).grid_span == 3
        assert t3x4.cell(0, 0).row_span == 2
        t3x4.split_cells((0, 1), (0, 2))
        assert t3x4.cell(0, 0).grid_span == 1

    def it_keeps_existing_Row_height_setter_working(self, t3x4):
        t3x4.rows[0].height = Emu(500000)
        assert t3x4.rows[0].height == 500000

    def it_keeps_existing_Column_width_setter_working(self, t3x4):
        t3x4.columns[0].width = Emu(800000)
        assert t3x4.columns[0].width == 800000
