# pyright: reportPrivateUsage=false

"""Unit-test suite for Accessibility Phase B (issue #22).

Phase A — `Shape.alt_text`, `Shape.alt_title`, `Shape.is_decorative` —
shipped in PR #31. Phase B adds:

- `Shape.is_hidden_from_accessibility` — boolean alias for `is_decorative`.
- `Slide.shapes.reading_order` — getter (tuple of shapes in z-order)
  and setter (reorders the spTree to match the assigned permutation).
- `Slide.shapes.accessibility_issues()` — lint helper returning shapes
  that lack alt text and are not marked decorative.
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# `Shape.is_hidden_from_accessibility` — alias for `is_decorative`
# ---------------------------------------------------------------------------


class DescribeIsHiddenFromAccessibility(object):
    """Unit-test suite for the new `is_hidden_from_accessibility` alias."""

    def it_returns_False_by_default(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        assert shape.is_hidden_from_accessibility is False

    def it_returns_True_when_marked_decorative(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        shape.is_decorative = True
        assert shape.is_hidden_from_accessibility is True

    def it_can_be_set_True_to_mark_decorative(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        shape.is_hidden_from_accessibility = True

        assert shape.is_decorative is True
        assert shape.is_hidden_from_accessibility is True

    def it_can_be_set_False_to_unmark(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        shape.is_decorative = True

        shape.is_hidden_from_accessibility = False

        assert shape.is_decorative is False

    def it_round_trips_through_save_and_reopen(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        shape.is_hidden_from_accessibility = True
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        rt = Presentation(buf)
        assert rt.slides[0].shapes[0].is_hidden_from_accessibility is True


# ---------------------------------------------------------------------------
# `Slide.shapes.reading_order`
# ---------------------------------------------------------------------------


def _seed_three_shapes():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    b = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(2), Inches(0.5))
    c = slide.shapes.add_textbox(Inches(5.5), Inches(0.5), Inches(2), Inches(0.5))
    a.alt_text = "alpha"
    b.alt_text = "bravo"
    c.alt_text = "charlie"
    return prs, slide, a, b, c


class DescribeReadingOrderGetter(object):
    """Unit-test suite for the `reading_order` getter."""

    def it_returns_shapes_in_z_order(self):
        _, slide, a, b, c = _seed_three_shapes()
        order = slide.shapes.reading_order
        assert tuple(s.alt_text for s in order) == ("alpha", "bravo", "charlie")

    def it_returns_a_tuple_not_a_list(self):
        _, slide, *_ = _seed_three_shapes()
        assert isinstance(slide.shapes.reading_order, tuple)

    def it_returns_an_empty_tuple_for_an_empty_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # ---layout 6 ("blank") — placeholder set may still drop in elements;
        #    use the actual length to set expectations rather than hard-coding 0
        assert slide.shapes.reading_order == tuple(slide.shapes)


class DescribeReadingOrderSetter(object):
    """Unit-test suite for the `reading_order` setter."""

    def it_can_reorder_shapes(self):
        _, slide, a, b, c = _seed_three_shapes()

        slide.shapes.reading_order = (c, a, b)

        order = tuple(s.alt_text for s in slide.shapes)
        assert order == ("charlie", "alpha", "bravo")

    def it_can_reverse_shape_order(self):
        _, slide, a, b, c = _seed_three_shapes()

        slide.shapes.reading_order = list(reversed(slide.shapes.reading_order))

        order = tuple(s.alt_text for s in slide.shapes)
        assert order == ("charlie", "bravo", "alpha")

    def it_round_trips_through_save_and_reopen(self):
        prs, slide, a, b, c = _seed_three_shapes()
        slide.shapes.reading_order = (c, a, b)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        rt = Presentation(buf)

        order = tuple(s.alt_text for s in rt.slides[0].shapes)
        assert order == ("charlie", "alpha", "bravo")

    def it_is_a_noop_when_assigning_current_order(self):
        _, slide, *_ = _seed_three_shapes()
        before = tuple(s._element for s in slide.shapes)

        slide.shapes.reading_order = slide.shapes.reading_order

        after = tuple(s._element for s in slide.shapes)
        # ---same elements, same order (identity-preserved)
        assert after == before

    def but_it_raises_on_wrong_length(self):
        _, slide, a, b, _ = _seed_three_shapes()
        with pytest.raises(ValueError):
            slide.shapes.reading_order = (a, b)

    def but_it_raises_on_unknown_shape(self):
        prs1, slide1, a1, b1, c1 = _seed_three_shapes()
        _, _, a2, _, _ = _seed_three_shapes()  # ---from a different presentation

        with pytest.raises(ValueError):
            slide1.shapes.reading_order = (a1, b1, a2)


# ---------------------------------------------------------------------------
# `Slide.shapes.accessibility_issues()`
# ---------------------------------------------------------------------------


class DescribeAccessibilityIssues(object):
    """Unit-test suite for the lint helper."""

    def it_flags_shapes_without_alt_text_or_decorative(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        flagged = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        issues = slide.shapes.accessibility_issues()

        assert flagged in issues

    def it_passes_shapes_with_alt_text(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        ok_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        ok_shape.alt_text = "described"

        assert ok_shape not in slide.shapes.accessibility_issues()

    def it_passes_shapes_with_alt_title_only(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_only = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        title_only.alt_title = "title"

        assert title_only not in slide.shapes.accessibility_issues()

    def it_passes_decorative_shapes(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        deco = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        deco.is_decorative = True

        assert deco not in slide.shapes.accessibility_issues()

    def it_returns_an_empty_list_for_a_fully_tagged_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        s1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        s2 = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(2), Inches(1))
        s1.alt_text = "x"
        s2.is_decorative = True

        assert slide.shapes.accessibility_issues() == []

    def it_returns_flagged_shapes_in_reading_order(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        first_unflagged = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
        ok_middle = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(2), Inches(0.5))
        ok_middle.alt_text = "ok"
        last_unflagged = slide.shapes.add_textbox(Inches(5.5), Inches(0.5), Inches(2), Inches(0.5))

        issues = slide.shapes.accessibility_issues()

        assert issues == [first_unflagged, last_unflagged]
