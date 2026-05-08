# pyright: reportPrivateUsage=false

"""Unit-test suite for slide-CRUD Phase 4 — Presentation.sections.

Covers:

- The new oxml types in `pptx.oxml.presentation` (extLst plumbing,
  sectionLst, section, section-sldIdLst, section-sldId).
- The `pptx.sections` proxy classes — `Section` and `_Sections`.
- Round-trip preservation: open → mutate → save → reopen.
- Anti-criteria: section membership survives `Slides.move` and
  `Slides.remove` (because sections key on slide_id, not position).

Issue: https://github.com/MHoroszowski/python-pptx/issues/11 (Phase 4).
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.oxml.presentation import (
    SECTION_LIST_EXT_URI,
    CT_Presentation,
    CT_Section,
    CT_SectionList,
    CT_SectionSlideIdList,
)
from pptx.sections import Section, _new_section_id, _Sections

from .unitutil.cxml import element

# ---------------------------------------------------------------------------
# Section ID generator
# ---------------------------------------------------------------------------


class Describe_new_section_id(object):
    """Unit-test suite for `pptx.sections._new_section_id`."""

    def it_returns_a_GUID_in_braces(self):
        sid = _new_section_id()
        assert sid.startswith("{")
        assert sid.endswith("}")
        assert len(sid) == 38  # ---{ + 36 GUID chars + }

    def it_returns_unique_ids_on_each_call(self):
        ids = {_new_section_id() for _ in range(50)}
        assert len(ids) == 50


# ---------------------------------------------------------------------------
# OXML LAYER — sectionLst, section, section-sldIdLst.
# ---------------------------------------------------------------------------


class DescribeCT_SectionList(object):
    """Unit-test suite for `CT_SectionList` add/insert helpers."""

    def it_can_add_a_section(self):
        sectionLst = element("p14:sectionLst")
        assert isinstance(sectionLst, CT_SectionList)

        sec = sectionLst.add_section(name="Intro", section_id="abc")

        assert isinstance(sec, CT_Section)
        assert sec.name == "Intro"
        assert sec.id == "abc"
        assert len(sectionLst.section_lst) == 1

    @pytest.mark.parametrize(
        ("idx", "expected_position"),
        [(0, 0), (1, 1), (2, 2)],  # ---append==len, head, middle
    )
    def it_can_insert_a_section_at_a_specific_index(self, idx, expected_position):
        sectionLst = element("p14:sectionLst/(p14:section{name=A,id=a},p14:section{name=B,id=b})")

        sectionLst.insert_section_at(name="X", section_id="x", idx=idx)

        assert sectionLst.section_lst[expected_position].id == "x"

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_insert_out_of_range(self, bad_idx):
        sectionLst = element("p14:sectionLst/p14:section{name=A,id=a}")
        with pytest.raises(IndexError):
            sectionLst.insert_section_at(name="X", section_id="x", idx=bad_idx)


class DescribeCT_SectionSlideIdList(object):
    """Unit-test suite for `CT_SectionSlideIdList` add/remove helpers."""

    def it_can_add_a_sldId(self):
        sldIdLst = element("p14:sldIdLst")
        assert isinstance(sldIdLst, CT_SectionSlideIdList)

        sldIdLst.add_sldId(256)

        assert len(sldIdLst.sldId_lst) == 1
        assert sldIdLst.sldId_lst[0].id == 256

    def it_can_remove_a_sldId_for_a_slide_id(self):
        sldIdLst = element("p14:sldIdLst/(p14:sldId{id=256},p14:sldId{id=257})")

        removed = sldIdLst.remove_sldId_for(256)

        assert removed is True
        assert [s.id for s in sldIdLst.sldId_lst] == [257]

    def it_returns_False_when_sldId_is_absent(self):
        sldIdLst = element("p14:sldIdLst/p14:sldId{id=256}")
        assert sldIdLst.remove_sldId_for(999) is False


class DescribeCT_Presentation_section_helpers(object):
    """Unit-test suite for `CT_Presentation` section-list traversal helpers."""

    def it_returns_None_for_section_list_when_no_extLst(self):
        prs = element("p:presentation/p:sldIdLst/p:sldId{id=256,r:id=rId1}")
        assert isinstance(prs, CT_Presentation)
        assert prs.section_list is None

    def it_returns_None_for_section_list_when_no_section_ext(self):
        # ---build programmatically because cxml attr_val cannot embed the
        #    GUID URI with literal `{` and `}`---
        prs = element("p:presentation")
        assert isinstance(prs, CT_Presentation)
        extLst = prs.get_or_add_extLst()
        extLst._add_ext(uri="some-other-uri-not-the-section-one")

        assert prs.section_list is None

    def it_returns_the_sectionLst_when_present(self):
        prs = element("p:presentation")
        assert isinstance(prs, CT_Presentation)
        prs.get_or_add_section_list()  # ---creates extLst/ext/sectionLst chain

        sectionLst = prs.section_list
        assert isinstance(sectionLst, CT_SectionList)

    def it_creates_extLst_ext_and_sectionLst_on_get_or_add(self):
        prs = element("p:presentation")

        sectionLst = prs.get_or_add_section_list()

        assert isinstance(sectionLst, CT_SectionList)
        assert prs.extLst is not None
        ext = prs.extLst.ext_by_uri(SECTION_LIST_EXT_URI)
        assert ext is not None
        assert ext.sectionLst is sectionLst

    def it_is_a_noop_to_remove_section_list_when_already_absent(self):
        prs = element("p:presentation")
        prs.remove_section_list()  # ---should not raise---
        assert prs.section_list is None

    def it_drops_only_the_section_ext_keeping_other_extensions(self):
        prs = element("p:presentation")
        extLst = prs.get_or_add_extLst()
        extLst._add_ext(uri="other-uri-keep-me")
        prs.get_or_add_section_list()  # ---adds the section ext alongside

        prs.remove_section_list()

        assert prs.section_list is None
        # ---other ext is preserved, extLst is preserved---
        assert prs.extLst is not None
        assert prs.extLst.ext_by_uri("other-uri-keep-me") is not None

    def it_removes_extLst_entirely_when_only_section_ext_was_present(self):
        prs = element("p:presentation")
        prs.get_or_add_section_list()
        assert prs.extLst is not None  # ---sanity check before remove

        prs.remove_section_list()

        assert prs.extLst is None


# ---------------------------------------------------------------------------
# `Section` proxy — name, id, equality, slides view.
# ---------------------------------------------------------------------------


def _empty_sections() -> _Sections:
    """Build a `_Sections` proxy bound to a fresh empty presentation."""
    return _Sections(Presentation())


class DescribeSection(object):
    """Unit-test suite for `pptx.sections.Section`."""

    def it_knows_its_name(self):
        ct = element("p14:section{name=Intro,id=abc}")
        section = Section(ct, _empty_sections())
        assert section.name == "Intro"

    def it_can_change_its_name(self):
        ct = element("p14:section{name=Old,id=abc}")
        section = Section(ct, _empty_sections())

        section.name = "New"

        assert ct.name == "New"

    def it_knows_its_id(self):
        ct = element("p14:section{name=Intro,id=abc}")
        section = Section(ct, _empty_sections())
        assert section.id == "abc"

    def it_returns_an_empty_tuple_for_slides_when_sldIdLst_is_absent(self):
        ct = element("p14:section{name=Intro,id=abc}")
        section = Section(ct, _empty_sections())
        assert section.slides == ()

    def it_compares_equal_to_another_proxy_with_the_same_id(self):
        ct1 = element("p14:section{name=A,id=same-id}")
        ct2 = element("p14:section{name=B,id=same-id}")
        sections = _empty_sections()
        assert Section(ct1, sections) == Section(ct2, sections)

    def it_compares_unequal_to_another_proxy_with_a_different_id(self):
        ct1 = element("p14:section{name=A,id=id1}")
        ct2 = element("p14:section{name=A,id=id2}")
        sections = _empty_sections()
        assert Section(ct1, sections) != Section(ct2, sections)

    def it_is_hashable_by_id(self):
        ct = element("p14:section{name=A,id=abc}")
        section = Section(ct, _empty_sections())
        assert hash(section) == hash("abc")


# ---------------------------------------------------------------------------
# `_Sections` proxy — read-only collection semantics on empty deck.
# ---------------------------------------------------------------------------


class DescribeSections_EmptyPresentation(object):
    """Unit-test suite for `_Sections` over a presentation without sections."""

    def it_reports_zero_length(self):
        prs = Presentation()
        assert len(prs.sections) == 0

    def it_returns_an_empty_iterator(self):
        prs = Presentation()
        assert list(prs.sections) == []

    def it_does_not_create_extLst_just_to_report_zero(self):
        prs = Presentation()
        _ = len(prs.sections)
        assert prs._element.extLst is None

    def it_raises_IndexError_on_indexed_access(self):
        prs = Presentation()
        with pytest.raises(IndexError):
            prs.sections[0]


# ---------------------------------------------------------------------------
# `_Sections` — add / remove / index / membership.
# ---------------------------------------------------------------------------


class DescribeSections_AddRemove(object):
    """Unit-test suite for `_Sections` mutation surface."""

    def it_can_add_a_section(self):
        prs = Presentation()
        section = prs.sections.add_section("Intro")

        assert len(prs.sections) == 1
        assert section.name == "Intro"
        assert prs.sections[0] == section

    def it_can_add_multiple_sections_in_order(self):
        prs = Presentation()
        s1 = prs.sections.add_section("First")
        s2 = prs.sections.add_section("Second")
        s3 = prs.sections.add_section("Third")

        assert [s.name for s in prs.sections] == ["First", "Second", "Third"]
        assert prs.sections.index(s1) == 0
        assert prs.sections.index(s2) == 1
        assert prs.sections.index(s3) == 2

    def it_can_insert_a_section_after_another(self):
        prs = Presentation()
        first = prs.sections.add_section("First")
        last = prs.sections.add_section("Last")

        middle = prs.sections.add_section("Middle", after=first)

        assert [s.name for s in prs.sections] == ["First", "Middle", "Last"]
        assert prs.sections.index(middle) == 1
        assert prs.sections.index(last) == 2

    def it_assigns_unique_ids_to_each_section(self):
        prs = Presentation()
        s1 = prs.sections.add_section("A")
        s2 = prs.sections.add_section("B")
        assert s1.id != s2.id

    def but_it_raises_on_add_after_a_section_not_in_collection(self):
        prs1 = Presentation()
        prs2 = Presentation()
        foreign = prs2.sections.add_section("Foreign")

        with pytest.raises(ValueError):
            prs1.sections.add_section("X", after=foreign)

    def it_can_remove_a_section(self):
        prs = Presentation()
        s1 = prs.sections.add_section("Keep")
        s2 = prs.sections.add_section("Drop")

        prs.sections.remove(s2)

        assert len(prs.sections) == 1
        assert prs.sections[0] == s1

    def it_cleans_up_extLst_when_last_section_removed(self):
        prs = Presentation()
        only = prs.sections.add_section("Lonely")
        assert prs._element.extLst is not None  # ---created during add

        prs.sections.remove(only)

        assert len(prs.sections) == 0
        assert prs._element.extLst is None

    def but_it_raises_on_remove_of_section_not_in_collection(self):
        prs1 = Presentation()
        prs2 = Presentation()
        foreign = prs2.sections.add_section("Foreign")

        with pytest.raises(ValueError):
            prs1.sections.remove(foreign)

    def but_index_raises_on_section_not_in_collection(self):
        prs1 = Presentation()
        prs2 = Presentation()
        foreign = prs2.sections.add_section("Foreign")

        with pytest.raises(ValueError):
            prs1.sections.index(foreign)


# ---------------------------------------------------------------------------
# Slide assignment — Section.add_slide / remove_slide invariants.
# ---------------------------------------------------------------------------


def _seed_with_slides(n: int) -> Presentation:
    prs = Presentation()
    layout = prs.slide_layouts[6]
    for _ in range(n):
        prs.slides.add_slide(layout)
    return prs


class DescribeSection_SlideAssignment(object):
    """Unit-test suite for slide membership operations on Section."""

    def it_can_add_a_slide_to_a_section(self):
        prs = _seed_with_slides(2)
        section = prs.sections.add_section("Intro")

        section.add_slide(prs.slides[0])

        assert len(section.slides) == 1
        assert section.slides[0].slide_id == prs.slides[0].slide_id

    def it_unassigns_slide_from_other_section_on_add(self):
        prs = _seed_with_slides(1)
        first = prs.sections.add_section("First")
        second = prs.sections.add_section("Second")
        slide = prs.slides[0]

        first.add_slide(slide)
        second.add_slide(slide)

        assert first.slides == ()
        assert len(second.slides) == 1
        assert second.slides[0].slide_id == slide.slide_id

    def it_is_a_noop_to_add_a_slide_already_in_this_section(self):
        prs = _seed_with_slides(1)
        section = prs.sections.add_section("Intro")
        slide = prs.slides[0]

        section.add_slide(slide)
        section.add_slide(slide)  # ---should not duplicate

        assert len(section.slides) == 1

    def but_it_raises_on_add_of_slide_not_in_presentation(self):
        prs1 = _seed_with_slides(1)
        prs2 = _seed_with_slides(1)
        section = prs1.sections.add_section("Intro")

        with pytest.raises(ValueError):
            section.add_slide(prs2.slides[0])

    def it_can_remove_a_slide(self):
        prs = _seed_with_slides(2)
        section = prs.sections.add_section("Intro")
        slide = prs.slides[0]
        section.add_slide(slide)

        section.remove_slide(slide)

        assert section.slides == ()

    def but_it_raises_on_remove_of_slide_not_in_section(self):
        prs = _seed_with_slides(2)
        section = prs.sections.add_section("Intro")

        with pytest.raises(ValueError):
            section.remove_slide(prs.slides[0])

    def but_it_raises_on_remove_when_section_has_no_sldIdLst(self):
        # ---freshly-added section starts without a sldIdLst child
        prs = _seed_with_slides(1)
        section = prs.sections.add_section("Empty")

        with pytest.raises(ValueError):
            section.remove_slide(prs.slides[0])

    def it_returns_slides_in_section_order_not_presentation_order(self):
        prs = _seed_with_slides(3)
        section = prs.sections.add_section("Intro")

        section.add_slide(prs.slides[2])
        section.add_slide(prs.slides[0])
        section.add_slide(prs.slides[1])

        section_ids = [s.slide_id for s in section.slides]
        prs_ids = [s.slide_id for s in prs.slides]
        assert section_ids == [prs_ids[2], prs_ids[0], prs_ids[1]]


# ---------------------------------------------------------------------------
# Round-trip integration tests — open → mutate → save → reopen.
# ---------------------------------------------------------------------------


def _round_trip(prs: Presentation) -> Presentation:
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


class DescribeSectionsRoundTrip(object):
    """Save → reopen preservation of section structure."""

    def it_round_trips_a_presentation_with_no_sections(self):
        prs = _seed_with_slides(2)

        round_tripped = _round_trip(prs)

        assert len(round_tripped.sections) == 0

    def it_round_trips_sections_with_their_names_and_membership(self):
        prs = _seed_with_slides(3)
        intro = prs.sections.add_section("Intro")
        body = prs.sections.add_section("Body")
        intro.add_slide(prs.slides[0])
        body.add_slide(prs.slides[1])
        body.add_slide(prs.slides[2])

        rt = _round_trip(prs)

        assert [s.name for s in rt.sections] == ["Intro", "Body"]
        assert len(rt.sections[0].slides) == 1
        assert len(rt.sections[1].slides) == 2

    def it_preserves_section_id_through_round_trip(self):
        prs = _seed_with_slides(1)
        section = prs.sections.add_section("Intro")
        original_id = section.id

        rt = _round_trip(prs)

        assert rt.sections[0].id == original_id

    def it_preserves_section_membership_through_a_slides_move(self):
        """The whole point of keying sections by slide_id (not position)."""
        prs = _seed_with_slides(3)
        section = prs.sections.add_section("Intro")
        section.add_slide(prs.slides[0])
        target_slide_id = prs.slides[0].slide_id

        prs.slides.move(prs.slides[0], 2)  # ---move to last position
        rt = _round_trip(prs)

        section_slide_ids = [s.slide_id for s in rt.sections[0].slides]
        assert section_slide_ids == [target_slide_id]
        assert rt.slides[2].slide_id == target_slide_id

    def it_preserves_other_assignments_when_a_slide_is_removed(self):
        prs = _seed_with_slides(3)
        section = prs.sections.add_section("Body")
        section.add_slide(prs.slides[0])
        section.add_slide(prs.slides[1])
        section.add_slide(prs.slides[2])
        keep_id = prs.slides[0].slide_id
        also_keep_id = prs.slides[2].slide_id

        prs.slides.remove(prs.slides[1])  # ---middle slide goes
        rt = _round_trip(prs)

        section_ids = {s.slide_id for s in rt.sections[0].slides}
        # ---removed slide's id is now stale and silently skipped
        assert keep_id in section_ids
        assert also_keep_id in section_ids
        assert len(rt.sections[0].slides) == 2

    def it_preserves_sibling_extLst_extensions_when_adding_a_section(self):
        """User-authored ext alongside our section ext must round-trip intact."""
        prs = Presentation()
        # ---author a foreign extension before any section work---
        extLst = prs._element.get_or_add_extLst()
        extLst._add_ext(uri="user-foreign-ext-uri-keep-me")

        prs.sections.add_section("Intro")
        rt = _round_trip(prs)

        # ---both the section ext and the foreign ext survive---
        assert len(rt.sections) == 1
        rt_extLst = rt._element.extLst
        assert rt_extLst is not None
        assert rt_extLst.ext_by_uri("user-foreign-ext-uri-keep-me") is not None
        assert rt_extLst.ext_by_uri(SECTION_LIST_EXT_URI) is not None

    def it_preserves_orphan_section_sldId_entries_on_round_trip(self):
        """`Slides.remove` does NOT auto-prune section sldIdLst entries.

        Per PowerPoint compatibility doctrine: leave foreign data alone.
        The orphan id stays in the section XML; `Section.slides`
        silently skips it on read but the file content is preserved.
        """
        prs = _seed_with_slides(3)
        section = prs.sections.add_section("Body")
        section.add_slide(prs.slides[0])
        section.add_slide(prs.slides[1])
        section.add_slide(prs.slides[2])
        orphan_id = prs.slides[1].slide_id

        prs.slides.remove(prs.slides[1])
        rt = _round_trip(prs)

        # ---visible API skips the orphan---
        assert len(rt.sections[0].slides) == 2
        # ---raw XML still contains the orphan id (no auto-prune)---
        rt_section_lst = rt._element.section_list
        assert rt_section_lst is not None
        rt_section = rt_section_lst.section_lst[0]
        rt_sldIdLst = rt_section.sldIdLst
        assert rt_sldIdLst is not None
        raw_ids = [s.id for s in rt_sldIdLst.sldId_lst]
        assert orphan_id in raw_ids

    def it_round_trips_an_empty_section_with_zero_slides(self):
        """Empty sections must round-trip with their `<p14:sldIdLst/>` intact."""
        prs = _seed_with_slides(1)
        section = prs.sections.add_section("Empty")
        original_id = section.id

        rt = _round_trip(prs)

        assert len(rt.sections) == 1
        assert rt.sections[0].name == "Empty"
        assert rt.sections[0].id == original_id
        assert rt.sections[0].slides == ()

    def it_round_trips_unicode_and_xml_special_chars_in_section_name(self):
        """Section name must survive round-trip for non-ASCII and XML-reserved chars."""
        prs = _seed_with_slides(1)
        tricky_name = '§ec†ion <1> & "intro" 🎯'
        prs.sections.add_section(tricky_name)

        rt = _round_trip(prs)

        assert rt.sections[0].name == tricky_name

    def it_makes_slides_unsectioned_when_their_section_is_removed(self):
        """Slides survive section removal; they just lose their assignment."""
        prs = _seed_with_slides(2)
        section = prs.sections.add_section("Body")
        section.add_slide(prs.slides[0])
        section.add_slide(prs.slides[1])
        slide_ids_before = {s.slide_id for s in prs.slides}

        prs.sections.remove(section)

        # ---slides themselves remain in the presentation---
        assert {s.slide_id for s in prs.slides} == slide_ids_before
        # ---no section claims them now---
        assert len(prs.sections) == 0
