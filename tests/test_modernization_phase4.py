# pyright: reportPrivateUsage=false

"""Unit-test suite for Modernization Phase 4 — shape-tree ergonomics.

Covers:

- ``_BaseShapes.iter_leaf_shapes()`` (closes scanny/python-pptx#435):
  recursive traversal that descends into ``GroupShape`` children, yielding
  only non-group leaf shapes.
- ``_BaseShapes.__getitem__(str)`` and ``SlidePlaceholders.__getitem__(str)``
  (closes scanny/python-pptx#800): Mapping-like name access. Plus
  ``__contains__`` and ``keys()`` helpers on both collections.
- ``BaseShape.find_by_xpath(xpath, namespaces=None)``: power-user XPath
  escape hatch over a shape's element subtree.
- ``_BaseShapes.in_selection_pane_order()`` (closes scanny/python-pptx#532):
  return shapes in PowerPoint's Selection Pane order (reverse XML / z-order).
- Anti-criteria: integer ``__getitem__`` keeps existing behavior; iter
  semantics unchanged.

Issue: https://github.com/MHoroszowski/python-pptx/issues/29 (Phase 4).
"""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


def _make_slide_with_two_placeholders():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    return prs, slide


@pytest.fixture
def slide_fixture():
    _, slide = _make_slide_with_two_placeholders()
    return slide


@pytest.fixture
def slide_with_extra_textboxes():
    _, slide = _make_slide_with_two_placeholders()
    # ---add three textboxes so we have 5 shapes total---
    for i in range(3):
        tb = slide.shapes.add_textbox(Inches(1 + i), Inches(3), Inches(1), Inches(0.5))
        tb.name = "Box %d" % (i + 1)
    return slide


# ---------------------------------------------------------------------------
# iter_leaf_shapes (scanny#435)
# ---------------------------------------------------------------------------


class DescribeShapes_iter_leaf_shapes(object):
    """Unit-test suite for `_BaseShapes.iter_leaf_shapes`."""

    def it_returns_an_iterator(self, slide_fixture):
        result = slide_fixture.shapes.iter_leaf_shapes()
        # ---generator object satisfies the Iterator protocol---
        assert hasattr(result, "__next__")
        assert hasattr(result, "__iter__")

    def it_yields_top_level_shapes_when_no_groups_present(self, slide_fixture):
        leaves = list(slide_fixture.shapes.iter_leaf_shapes())
        top_level = list(slide_fixture.shapes)
        assert len(leaves) == len(top_level)
        assert [s.name for s in leaves] == [s.name for s in top_level]

    def it_excludes_GroupShape_instances(self, slide_with_extra_textboxes):
        for s in slide_with_extra_textboxes.shapes.iter_leaf_shapes():
            assert not isinstance(s, GroupShape)

    def it_yields_5_shapes_for_a_slide_with_5_top_level_shapes(self, slide_with_extra_textboxes):
        leaves = list(slide_with_extra_textboxes.shapes.iter_leaf_shapes())
        assert len(leaves) == 5


# ---------------------------------------------------------------------------
# Mapping-like access (scanny#800)
# ---------------------------------------------------------------------------


class DescribeShapes_MappingAccess(object):
    """Unit-test suite for `__getitem__(str)`, `__contains__`, `keys()`."""

    def it_returns_shape_by_name_via_string_key(self, slide_fixture):
        title = slide_fixture.shapes["Title 1"]
        assert title.name == "Title 1"

    def it_raises_KeyError_on_unknown_name(self, slide_fixture):
        with pytest.raises(KeyError):
            slide_fixture.shapes["Bogus"]

    def it_keeps_integer_indexing_unchanged(self, slide_fixture):
        first = slide_fixture.shapes[0]
        # ---first shape on the Title+Content layout is the title placeholder---
        assert first.name == "Title 1"

    def it_raises_IndexError_on_int_out_of_range(self, slide_fixture):
        with pytest.raises(IndexError):
            slide_fixture.shapes[99]

    def it_supports_string_membership_check(self, slide_fixture):
        assert "Title 1" in slide_fixture.shapes
        assert "Bogus" not in slide_fixture.shapes

    def it_supports_integer_index_range_membership(self, slide_fixture):
        assert 0 in slide_fixture.shapes
        assert 99 not in slide_fixture.shapes

    def it_returns_False_for_other_key_types(self, slide_fixture):
        assert (1.5 in slide_fixture.shapes) is False
        assert (None in slide_fixture.shapes) is False

    def it_rejects_bool_keys_explicitly(self, slide_fixture):
        # ---bool subclasses int; without this guard `shapes[True]` would
        # ---silently resolve to index 1, almost always a bug.
        with pytest.raises(TypeError):
            slide_fixture.shapes[True]
        with pytest.raises(TypeError):
            slide_fixture.shapes[False]
        # ---and `__contains__` should not match either way---
        assert (True in slide_fixture.shapes) is False
        assert (False in slide_fixture.shapes) is False

    def it_lists_names_via_keys(self, slide_fixture):
        names = slide_fixture.shapes.keys()
        assert names == ["Title 1", "Content Placeholder 2"]


class DescribeSlidePlaceholders_MappingAccess(object):
    """Unit-test suite for `SlidePlaceholders.__getitem__(str)` and friends."""

    def it_returns_placeholder_by_name(self, slide_fixture):
        title = slide_fixture.placeholders["Title 1"]
        assert title.name == "Title 1"

    def it_raises_KeyError_on_unknown_name(self, slide_fixture):
        with pytest.raises(KeyError):
            slide_fixture.placeholders["Bogus"]

    def it_keeps_integer_idx_lookup_unchanged(self, slide_fixture):
        # ---placeholder idx 0 = Title---
        title = slide_fixture.placeholders[0]
        assert title.name == "Title 1"

    def it_raises_KeyError_on_unknown_idx(self, slide_fixture):
        with pytest.raises(KeyError):
            slide_fixture.placeholders[999]

    def it_supports_string_and_int_membership(self, slide_fixture):
        assert "Title 1" in slide_fixture.placeholders
        assert "Bogus" not in slide_fixture.placeholders
        assert 0 in slide_fixture.placeholders
        assert 999 not in slide_fixture.placeholders

    def it_lists_names_via_keys(self, slide_fixture):
        names = slide_fixture.placeholders.keys()
        assert "Title 1" in names


# ---------------------------------------------------------------------------
# find_by_xpath
# ---------------------------------------------------------------------------


class DescribeShape_find_by_xpath(object):
    """Unit-test suite for `BaseShape.find_by_xpath`."""

    def it_returns_matching_elements_with_default_nsmap(self, slide_fixture):
        title = slide_fixture.shapes.title
        results = title.find_by_xpath(".//p:nvSpPr")
        assert len(results) == 1

    def it_returns_empty_list_on_no_match(self, slide_fixture):
        title = slide_fixture.shapes.title
        assert title.find_by_xpath(".//a:nope_does_not_exist") == []

    def it_accepts_a_custom_namespace_dict(self, slide_fixture):
        title = slide_fixture.shapes.title
        custom = title.find_by_xpath(
            ".//foo:nvSpPr",
            namespaces={"foo": "http://schemas.openxmlformats.org/presentationml/2006/main"},
        )
        assert len(custom) == 1


# ---------------------------------------------------------------------------
# in_selection_pane_order (scanny#532)
# ---------------------------------------------------------------------------


class DescribeShapes_in_selection_pane_order(object):
    """Unit-test suite for `_BaseShapes.in_selection_pane_order`."""

    def it_returns_a_tuple(self, slide_fixture):
        result = slide_fixture.shapes.in_selection_pane_order()
        assert isinstance(result, tuple)

    def it_reverses_the_xml_document_order(self, slide_fixture):
        xml_order = [s.name for s in slide_fixture.shapes]
        sp_order = [s.name for s in slide_fixture.shapes.in_selection_pane_order()]
        assert sp_order == list(reversed(xml_order))

    def it_preserves_length(self, slide_with_extra_textboxes):
        assert len(slide_with_extra_textboxes.shapes.in_selection_pane_order()) == len(
            slide_with_extra_textboxes.shapes
        )

    def it_does_not_mutate_the_collection(self, slide_with_extra_textboxes):
        before = [s.name for s in slide_with_extra_textboxes.shapes]
        _ = slide_with_extra_textboxes.shapes.in_selection_pane_order()
        after = [s.name for s in slide_with_extra_textboxes.shapes]
        assert before == after


# ---------------------------------------------------------------------------
# Anti / Regression
# ---------------------------------------------------------------------------


class DescribePhase4_Regression(object):
    """Anti-criteria — existing surfaces unchanged."""

    def it_keeps_phase2_by_name_working(self, slide_fixture):
        title = slide_fixture.shapes.by_name("Title 1")
        assert title.name == "Title 1"

    def it_keeps_iteration_yielding_shape_values(self, slide_fixture):
        # ---if we'd switched to Mapping ABC, __iter__ would yield keys---
        for s in slide_fixture.shapes:
            assert hasattr(s, "name")
            assert hasattr(s, "shape_id")

    def it_keeps_len_returning_shape_count(self, slide_fixture):
        assert len(slide_fixture.shapes) == 2  # Title + Content placeholders
