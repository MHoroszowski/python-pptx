# pyright: reportPrivateUsage=false

"""Integration test suite for customXml round-trip.

Loads each synthetic fixture under ``tests/test_files/customxml/``, exercises
the public API against it, saves to a fresh BytesIO, reloads, and asserts the
state survived.

Real third-party fixtures (SharePoint-saved, Office.js-produced, VSTO-tooled)
will land later under ``sharepoint-saved.pptx`` etc. once captured during the
manual PowerPoint UI matrix in ``Plans/customxml-implementation-plan.md`` §5.4.
"""

from __future__ import annotations

import os
from io import BytesIO

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.custom_xml import CustomXmlPart

_FIXTURE_DIR = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    os.pardir,
    "test_files",
    "customxml",
)


def _fixture(name: str) -> str:
    return os.path.join(_FIXTURE_DIR, name)


def _roundtrip(prs):
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


class DescribePresentationScopedFixture:
    def it_loads_the_part(self):
        prs = Presentation(_fixture("presentation-scoped.pptx"))
        assert len(prs.custom_xml_parts) == 1

    def it_upgrades_loaded_part_to_CustomXmlPart_class(self):
        prs = Presentation(_fixture("presentation-scoped.pptx"))
        part = prs.custom_xml_parts[0]
        assert isinstance(part, CustomXmlPart)

    def it_preserves_the_payload(self):
        prs = Presentation(_fixture("presentation-scoped.pptx"))
        part = prs.custom_xml_parts.by_name("provenance")
        assert part is not None
        assert part.element.tag == "{urn:my:provenance}provenance"
        source = part.element.find("{urn:my:provenance}source")
        assert source is not None
        assert source.text == "integration-fixture"

    def it_preserves_the_pinned_guid(self):
        prs = Presentation(_fixture("presentation-scoped.pptx"))
        part = prs.custom_xml_parts[0]
        assert part.datastore_item_id == "{1A2B3C4D-5E6F-7890-ABCD-EF1234567890}"

    def it_preserves_the_schema_refs(self):
        prs = Presentation(_fixture("presentation-scoped.pptx"))
        part = prs.custom_xml_parts[0]
        assert part.schema_refs == ("urn:my:provenance",)

    def it_preserves_the_presentation_scope_through_save(self):
        prs = Presentation(_fixture("presentation-scoped.pptx"))
        reloaded = _roundtrip(prs)
        prs_rel_types = {r.reltype for r in reloaded.part.rels.values()}
        pkg_rel_types = {r.reltype for r in reloaded.part.package._rels.values()}
        assert RT.CUSTOM_XML in prs_rel_types
        assert RT.CUSTOM_XML not in pkg_rel_types


class DescribePackageScopedFixture:
    def it_loads_the_part(self):
        prs = Presentation(_fixture("package-scoped.pptx"))
        assert len(prs.custom_xml_parts) == 1

    def it_preserves_the_payload(self):
        prs = Presentation(_fixture("package-scoped.pptx"))
        part = prs.custom_xml_parts.by_name("vsto")
        assert part is not None
        assert part.element.tag == "{urn:my:vsto}vsto-config"

    def it_preserves_the_package_scope_through_save(self):
        prs = Presentation(_fixture("package-scoped.pptx"))
        reloaded = _roundtrip(prs)
        prs_rel_types = {r.reltype for r in reloaded.part.rels.values()}
        pkg_rel_types = {r.reltype for r in reloaded.part.package._rels.values()}
        assert RT.CUSTOM_XML in pkg_rel_types
        assert RT.CUSTOM_XML not in prs_rel_types

    def it_preserves_the_pinned_guid(self):
        prs = Presentation(_fixture("package-scoped.pptx"))
        part = prs.custom_xml_parts[0]
        assert part.datastore_item_id == "{ABCDEF12-3456-7890-ABCD-EF1234567890}"


class DescribeMultipartFixture:
    def it_loads_two_customxml_parts_at_mixed_scopes(self):
        prs = Presentation(_fixture("multipart.pptx"))
        assert len(prs.custom_xml_parts) == 2 + 1  # provenance + extra + readme blob

    def it_preserves_custom_document_properties(self):
        prs = Presentation(_fixture("multipart.pptx"))
        assert prs.custom_properties["Source"] == "deck-builder-cli@1.4.2"
        assert prs.custom_properties["BuildNumber"] == 42
        assert prs.custom_properties["IsDraft"] is True

    def it_finds_each_part_by_name(self):
        prs = Presentation(_fixture("multipart.pptx"))
        assert prs.custom_xml_parts.by_name("provenance") is not None
        assert prs.custom_xml_parts.by_name("extra") is not None
        assert prs.custom_xml_parts.by_name("readme") is not None

    def it_round_trips_through_save_load_with_mutations(self):
        prs = Presentation(_fixture("multipart.pptx"))
        # mutate something in each layer
        prs.custom_properties["NewKey"] = "added"
        prs.custom_xml_parts.by_name("provenance").add_item("added-by-test", "value")

        reloaded = _roundtrip(prs)

        assert reloaded.custom_properties["NewKey"] == "added"
        assert reloaded.custom_properties["Source"] == "deck-builder-cli@1.4.2"
        prov = reloaded.custom_xml_parts.by_name("provenance")
        assert prov is not None
        # The added child element survived the round-trip
        added = [c for c in prov.element if c.tag.endswith("added-by-test")]
        assert len(added) == 1
        assert added[0].text == "value"

    def it_round_trips_the_string_blob_helper(self):
        prs = Presentation(_fixture("multipart.pptx"))
        content = prs.custom_xml_parts.read_string_blob("readme")
        assert content is not None
        assert "# Hello" in content
        assert "markdown content" in content

    def it_remove_then_save_drops_the_part(self):
        prs = Presentation(_fixture("multipart.pptx"))
        provenance = prs.custom_xml_parts.by_name("provenance")
        prs.custom_xml_parts.remove(provenance)
        reloaded = _roundtrip(prs)
        assert reloaded.custom_xml_parts.by_name("provenance") is None
        # Other parts still present
        assert reloaded.custom_xml_parts.by_name("extra") is not None
        assert reloaded.custom_xml_parts.by_name("readme") is not None


class DescribeCleanFixture:
    """A presentation with no customXml at all should have no related rels."""

    def it_has_no_customxml_parts(self):
        prs = Presentation(_fixture("clean.pptx"))
        assert len(prs.custom_xml_parts) == 0

    def it_round_trips_with_no_rels_added(self):
        prs = Presentation(_fixture("clean.pptx"))
        # do nothing
        reloaded = _roundtrip(prs)
        prs_rel_types = {r.reltype for r in reloaded.part.rels.values()}
        pkg_rel_types = {r.reltype for r in reloaded.part.package._rels.values()}
        assert RT.CUSTOM_XML not in prs_rel_types
        assert RT.CUSTOM_XML not in pkg_rel_types
        assert RT.CUSTOM_PROPERTIES not in pkg_rel_types

    def it_can_have_customxml_added_after_loading(self):
        prs = Presentation(_fixture("clean.pptx"))
        prs.custom_xml_parts.add(
            b'<after-load xmlns="u:al"/>',
            name="after-load",
        )
        reloaded = _roundtrip(prs)
        part = reloaded.custom_xml_parts.by_name("after-load")
        assert part is not None
        assert part.element.tag == "{u:al}after-load"


class DescribeCoreAndCustomCoexistence:
    def it_preserves_core_properties_alongside_custom_ones(self):
        prs = Presentation(_fixture("multipart.pptx"))
        prs.core_properties.author = "Athena"
        prs.core_properties.subject = "Integration test"

        reloaded = _roundtrip(prs)

        assert reloaded.core_properties.author == "Athena"
        assert reloaded.core_properties.subject == "Integration test"
        # custom properties still intact
        assert reloaded.custom_properties["Source"] == "deck-builder-cli@1.4.2"
        # customXml parts still intact
        assert len(reloaded.custom_xml_parts) == 3
