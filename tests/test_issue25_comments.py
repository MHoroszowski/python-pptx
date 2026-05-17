"""Wave-1 tests for issue #25 (Threaded Comments & Review).

Covers the oxml element classes and part classes that form the foundation
(SF1 + SF2, ISA ISC-1..12). `Slide.comments` is deliberately NOT exercised
here — that is Wave 2.

TDD ordering per repo CLAUDE.md §3: these tests are written to fail first.
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import PartFactory
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml
from pptx.oxml.comments import (
    CT_Comment,
    CT_CommentAuthor,
    CT_CommentAuthorList,
    CT_CommentList,
    CT_ThreadedComment,
)
from pptx.oxml.ns import nsdecls, nsuri, qn
from pptx.parts.comments import (
    CommentAuthorsPart,
    CommentsPart,
    ModernCommentsPart,
)

# -- oxml: legacy author list --------------------------------------------------


class DescribeCT_CommentAuthorList:
    def it_parses_a_hand_built_cmAuthorLst(self):
        xml = (
            "<p:cmAuthorLst %s>"
            '<p:cmAuthor id="0" name="Ada" initials="AL" lastIdx="1" clrIdx="0"/>'
            "</p:cmAuthorLst>" % nsdecls("p")
        )
        elm = parse_xml(xml)
        assert isinstance(elm, CT_CommentAuthorList)
        authors = elm.cmAuthor_lst
        assert len(authors) == 1
        assert isinstance(authors[0], CT_CommentAuthor)
        assert authors[0].id == 0
        assert authors[0].name == "Ada"
        assert authors[0].initials == "AL"


# -- oxml: legacy comment list -------------------------------------------------


class DescribeCT_CommentList:
    def it_parses_a_hand_built_cmLst(self):
        xml = (
            "<p:cmLst %s>"
            '<p:cm authorId="0" dt="2026-05-16T00:00:00" idx="1">'
            '<p:pos x="100" y="200"/>'
            "<p:text>Hello legacy</p:text>"
            "</p:cm>"
            "</p:cmLst>" % nsdecls("p")
        )
        elm = parse_xml(xml)
        assert isinstance(elm, CT_CommentList)
        comments = elm.cm_lst
        assert len(comments) == 1
        cm = comments[0]
        assert isinstance(cm, CT_Comment)
        assert cm.authorId == 0
        assert cm.idx == 1
        assert cm.text == "Hello legacy"


# -- oxml: modern threaded comment ---------------------------------------------


class DescribeModernThreadedNamespace:
    def it_registers_the_2018_8_main_namespace(self):
        assert nsuri("p188") == "http://schemas.microsoft.com/office/powerpoint/2018/8/main"
        # qn() must resolve a p188-prefixed tag to its Clark name
        assert qn("p188:cm").startswith(
            "{http://schemas.microsoft.com/office/powerpoint/2018/8/main}"
        )


class DescribeCT_ThreadedComment:
    def it_parses_a_hand_built_p188_cm(self):
        xml = (
            '<p188:cm %s id="{2E5D6F1A-0000-0000-0000-000000000001}" '
            'authorId="{11111111-1111-1111-1111-111111111111}" '
            'created="2026-05-16T12:00:00.000">'
            "<p188:txBody>"
            '<a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            '<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            "<a:r><a:t>modern comment</a:t></a:r></a:p>"
            "</p188:txBody>"
            "</p188:cm>" % nsdecls("p188")
        )
        elm = parse_xml(xml)
        assert isinstance(elm, CT_ThreadedComment)
        assert elm.id == "{2E5D6F1A-0000-0000-0000-000000000001}"
        assert elm.authorId == "{11111111-1111-1111-1111-111111111111}"
        assert elm.created == "2026-05-16T12:00:00.000"

    def it_parses_a_p188_cm_with_reply(self):
        xml = (
            '<p188:cm %s id="{AAA}" authorId="{BBB}" created="2026-05-16T12:00:00.000">'
            "<p188:txBody>"
            '<a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            "</p188:txBody>"
            "<p188:replyLst>"
            '<p188:reply id="{CCC}" authorId="{DDD}" created="2026-05-16T13:00:00.000">'
            "<p188:txBody>"
            '<a:bodyPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            "</p188:txBody>"
            "</p188:reply>"
            "</p188:replyLst>"
            "</p188:cm>" % nsdecls("p188")
        )
        elm = parse_xml(xml)
        assert isinstance(elm, CT_ThreadedComment)
        assert elm.replyLst is not None
        assert len(elm.replyLst.reply_lst) == 1


# -- parts: comment authors ----------------------------------------------------


class DescribeCommentAuthorsPart:
    def it_creates_an_empty_part_via_new(self):
        prs = Presentation()
        part = CommentAuthorsPart.new(prs.part.package)
        assert part.content_type == CT.PML_COMMENT_AUTHORS
        assert str(part.partname) == "/ppt/commentAuthors.xml"
        assert list(part.iter_authors()) == []

    def it_allocates_the_next_free_id_on_add_author(self):
        prs = Presentation()
        part = CommentAuthorsPart.new(prs.part.package)
        a0 = part.add_author("Ada", "AL")
        a1 = part.add_author("Babbage", "CB")
        assert a0.id == 0
        assert a1.id == 1
        assert a0.name == "Ada"
        assert [a.name for a in part.iter_authors()] == ["Ada", "Babbage"]

    def it_dedupes_same_name_on_get_or_add_author(self):
        # ISC-5: get_or_add_author reuses an existing same-name author id.
        prs = Presentation()
        part = CommentAuthorsPart.new(prs.part.package)
        first = part.get_or_add_author("Ada")
        again = part.get_or_add_author("Ada")
        assert first.id == again.id
        assert len(list(part.iter_authors())) == 1

    def it_round_trips_through_a_presentation_save_reopen(self):
        prs = Presentation()
        pres_part = prs.part
        authors_part = CommentAuthorsPart.new(pres_part.package)
        authors_part.add_author("Ada", "AL")
        pres_part.relate_to(authors_part, RT.COMMENT_AUTHORS)

        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)

        prs2 = Presentation(stream)
        reloaded = prs2.part.part_related_by(RT.COMMENT_AUTHORS)
        assert reloaded.content_type == CT.PML_COMMENT_AUTHORS
        names = [a.name for a in reloaded.iter_authors()]
        assert names == ["Ada"]


# -- parts: legacy + modern comments parts -------------------------------------


class DescribeCommentsPart:
    def it_creates_an_empty_legacy_part_via_new(self):
        prs = Presentation()
        part = CommentsPart.new(prs.part.package, PackURI("/ppt/comments/comment1.xml"))
        assert part.content_type == CT.PML_COMMENTS
        # empty <p:cmLst>
        assert len(part._element.cm_lst) == 0


class DescribeModernCommentsPart:
    def it_creates_an_empty_modern_part_via_new(self):
        prs = Presentation()
        part = ModernCommentsPart.new(
            prs.part.package, PackURI("/ppt/comments/modernComment_slide1.xml")
        )
        assert part.content_type == CT.PML_THREADED_COMMENTS


# -- content-type registry -----------------------------------------------------


class DescribeContentTypePartClassMap:
    def it_maps_the_two_legacy_part_classes(self):
        assert PartFactory.part_type_for[CT.PML_COMMENT_AUTHORS] is CommentAuthorsPart
        assert PartFactory.part_type_for[CT.PML_COMMENTS] is CommentsPart

    def it_maps_the_modern_threaded_part_class(self):
        assert PartFactory.part_type_for[CT.PML_THREADED_COMMENTS] is ModernCommentsPart


# -- Wave 2: Slide.comments collection + Comment/replies proxy -----------------


def _blank_slide():
    """A fresh presentation + one blank slide (layout 6 = Blank)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


class DescribeSlideComments:
    def it_exposes_an_empty_comments_collection_without_mutating(self):
        _prs, slide = _blank_slide()
        assert hasattr(slide, "comments")
        assert len(slide.comments) == 0
        assert list(slide.comments) == []
        # reading must not have created the modern-comments part (ISC-21
        # spirit: no side effects on read)
        assert slide.part.has_modern_comments is False

    def it_add_returns_a_Comment_and_increments_len(self):
        from pptx.comments import Comment

        _prs, slide = _blank_slide()
        comment = slide.comments.add("Hello", "Ada")
        assert isinstance(comment, Comment)
        assert comment.text == "Hello"
        assert len(slide.comments) == 1

    def it_creates_the_part_and_rel_on_first_add(self):
        _prs, slide = _blank_slide()
        assert slide.part.has_modern_comments is False
        slide.comments.add("first", "Ada")
        assert slide.part.has_modern_comments is True
        # same single part returned on every access
        assert slide.part.modern_comments_part is slide.part.modern_comments_part

    def it_dedupes_the_author_on_the_modern_authors_part(self):
        # ISC-5/ISC-21: two comments by the same author name -> one entry on
        # the MODERN authors part (NOT the legacy commentAuthors.xml — that
        # was the issue-#25 repair-dialog bug; the old assertion encoded it).
        prs, slide = _blank_slide()
        slide.comments.add("a", "Ada")
        slide.comments.add("b", "Ada")
        authors = list(prs.part.package.presentation_part.authors_part.iter_authors())
        assert [a.name for a in authors] == ["Ada"]
        # and the legacy part was never created for a modern-only deck
        assert prs.part.package.presentation_part.has_comment_authors is False

    def it_iterates_comments_in_document_order(self):
        _prs, slide = _blank_slide()
        slide.comments.add("one", "Ada")
        slide.comments.add("two", "Babbage")
        slide.comments.add("three", "Ada")
        assert [c.text for c in slide.comments] == ["one", "two", "three"]

    def it_remove_decrements_len(self):
        _prs, slide = _blank_slide()
        slide.comments.add("keep", "Ada")
        target = slide.comments.add("drop", "Ada")
        assert len(slide.comments) == 2
        slide.comments.remove(target)
        assert len(slide.comments) == 1
        assert [c.text for c in slide.comments] == ["keep"]

    def it_leaves_valid_rels_when_the_last_comment_is_removed(self):
        # ISC-21 anti-criterion: removing the last comment must NOT drop
        # the part/relationship and the file must still reopen cleanly.
        prs, slide = _blank_slide()
        only = slide.comments.add("solo", "Ada")
        slide.comments.remove(only)
        assert len(slide.comments) == 0
        assert slide.part.has_modern_comments is True  # rel intact
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        prs2 = Presentation(stream)
        assert len(prs2.slides[0].comments) == 0

    def it_round_trips_text_and_author_through_save_reopen(self):
        # ISC-19: save -> BytesIO -> reopen preserves text + author.
        prs, slide = _blank_slide()
        slide.comments.add("first body", "Ada Lovelace")
        slide.comments.add("second body", "Charles Babbage")
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        reloaded = list(Presentation(stream).slides[0].comments)
        assert [c.text for c in reloaded] == ["first body", "second body"]
        assert [c.author for c in reloaded] == ["Ada Lovelace", "Charles Babbage"]


class DescribeCommentReplies:
    def it_add_threads_a_reply_and_preserves_order(self):
        from pptx.comments import CommentReply

        _prs, slide = _blank_slide()
        comment = slide.comments.add("question?", "Ada")
        r1 = comment.replies.add("answer 1", "Babbage")
        comment.replies.add("answer 2", "Ada")
        assert isinstance(r1, CommentReply)
        assert len(comment.replies) == 2
        assert [r.text for r in comment.replies] == ["answer 1", "answer 2"]

    def it_does_not_detach_the_parent_comment_on_reply_add(self):
        # ISC-27 anti-criterion: adding a reply must not re-parent or drop
        # the top-level comment.
        _prs, slide = _blank_slide()
        comment = slide.comments.add("parent", "Ada")
        comment.replies.add("child", "Babbage")
        assert len(slide.comments) == 1
        assert slide.comments[0].text == "parent"

    def it_round_trips_replies_through_save_reopen(self):
        prs, slide = _blank_slide()
        comment = slide.comments.add("top", "Ada")
        comment.replies.add("reply one", "Babbage")
        comment.replies.add("reply two", "Ada")
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        reloaded = list(Presentation(stream).slides[0].comments)
        assert len(reloaded) == 1
        replies = list(reloaded[0].replies)
        assert [r.text for r in replies] == ["reply one", "reply two"]
        assert [r.author for r in replies] == ["Babbage", "Ada"]


class DescribeCommentProxyFields:
    def it_resolves_author_text_and_tz_aware_created_at(self):
        from datetime import datetime

        _prs, slide = _blank_slide()
        comment = slide.comments.add("body text", "Grace Hopper")
        assert comment.author == "Grace Hopper"
        assert comment.text == "body text"
        assert isinstance(comment.created_at, datetime)
        assert comment.created_at.tzinfo is not None  # tz-aware

    def it_reports_None_anchor_position_when_unanchored(self):
        _prs, slide = _blank_slide()
        comment = slide.comments.add("floating", "Ada")
        assert comment.anchor_position is None

    def it_resolves_anchor_position_to_shape_left_top_when_anchored(self):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        prs, slide = _blank_slide()
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(914400), Emu(457200), Emu(100), Emu(200)
        )
        comment = slide.comments.add("on the box", "Ada", anchor=shape)
        assert comment.anchor_position == (914400, 457200)
        # round-trips
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        reloaded = list(Presentation(stream).slides[0].comments)[0]
        assert reloaded.anchor_position == (914400, 457200)

    def it_returns_None_author_for_an_unknown_author_guid(self):
        _prs, slide = _blank_slide()
        comment = slide.comments.add("x", "Ada")
        comment._cm.authorId = "{00000000-0000-0000-0000-000000000000}"
        assert comment.author is None


# -- Wave 3: SF6 resolve / SF7 Shape.comments / SF8 legacy coexistence ---------


class DescribeCommentResolve:
    """SF6 (ISC-33..37): Comment.resolve() / Comment.resolved."""

    def it_defaults_resolved_to_False(self):
        _prs, slide = _blank_slide()
        comment = slide.comments.add("open question", "Ada")
        assert comment.resolved is False

    def it_marks_a_thread_resolved(self):
        _prs, slide = _blank_slide()
        comment = slide.comments.add("please review", "Ada")
        comment.resolve()
        assert comment.resolved is True

    def it_can_reopen_a_resolved_thread(self):
        _prs, slide = _blank_slide()
        comment = slide.comments.add("review me", "Ada")
        comment.resolve()
        comment.reopen()
        assert comment.resolved is False

    def it_round_trips_the_resolved_flag(self):
        # ISC-35: save -> reopen preserves resolution state.
        prs, slide = _blank_slide()
        a = slide.comments.add("resolved one", "Ada")
        slide.comments.add("still open", "Babbage")
        a.resolve()
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        reloaded = list(Presentation(stream).slides[0].comments)
        assert [c.resolved for c in reloaded] == [True, False]

    def it_raises_when_resolving_a_legacy_backed_comment(self):
        # ISC-37 anti: the legacy <p:cm> schema has no resolution concept.
        # Decision: raise (not silent no-op) so callers can't believe a
        # legacy thread was resolved when the file format cannot record it.
        prs, slide = _legacy_plus_modern_deck()
        legacy = next(c for c in slide.comments if c.is_legacy)
        with pytest.raises(TypeError):
            legacy.resolve()
        assert legacy.resolved is False


class DescribeShapeComments:
    """SF7 (ISC-38..41): per-shape comment filter."""

    def it_returns_an_empty_collection_for_a_shape_with_no_comments(self):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        _prs, slide = _blank_slide()
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(1), Emu(2), Emu(3), Emu(4))
        assert hasattr(shape, "comments")
        assert list(shape.comments) == []
        assert len(shape.comments) == 0

    def it_yields_only_comments_anchored_to_that_shape(self):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        _prs, slide = _blank_slide()
        a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(10), Emu(20), Emu(30), Emu(40))
        b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(50), Emu(60), Emu(70), Emu(80))
        slide.comments.add("on A one", "Ada", anchor=a)
        slide.comments.add("on B", "Babbage", anchor=b)
        slide.comments.add("on A two", "Ada", anchor=a)
        slide.comments.add("floating", "Grace")
        assert [c.text for c in a.comments] == ["on A one", "on A two"]

    def it_does_not_leak_a_comment_across_shapes(self):
        # ISC-40: a comment anchored to shape A must not appear under B.
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        _prs, slide = _blank_slide()
        a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(10), Emu(20), Emu(30), Emu(40))
        b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(50), Emu(60), Emu(70), Emu(80))
        slide.comments.add("only on A", "Ada", anchor=a)
        assert [c.text for c in a.comments] == ["only on A"]
        assert list(b.comments) == []

    def it_round_trips_the_per_shape_filter(self):
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu

        prs, slide = _blank_slide()
        a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(11), Emu(22), Emu(33), Emu(44))
        slide.comments.add("anchored", "Ada", anchor=a)
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        slide2 = Presentation(stream).slides[0]
        shape2 = next(s for s in slide2.shapes if s.shape_id == a.shape_id)
        assert [c.text for c in shape2.comments] == ["anchored"]


def _legacy_plus_modern_deck():
    """Build a deck that already has a LEGACY <p:cm> comment, then add a
    MODERN threaded comment to the same slide.

    The legacy comment is hand-injected into a saved package by editing the
    zip directly (our public ``.add()`` only writes modern), mirroring how
    the rest of this module builds raw-xml fixtures. Returns
    ``(prs, slide)`` of the *reopened* package with both families present.
    """
    import zipfile

    prs0 = Presentation()
    prs0.slides.add_slide(prs0.slide_layouts[6])
    base = io.BytesIO()
    prs0.save(base)
    base.seek(0)

    legacy_authors = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<p:cmAuthorLst xmlns:p="http://schemas.openxmlformats.org/'
        'presentationml/2006/main">'
        '<p:cmAuthor id="0" name="Legacy Larry" initials="LL" '
        'lastIdx="1" clrIdx="0"/>'
        "</p:cmAuthorLst>"
    )
    legacy_comments = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<p:cmLst xmlns:p="http://schemas.openxmlformats.org/'
        'presentationml/2006/main">'
        '<p:cm authorId="0" dt="2026-05-16T00:00:00" idx="1">'
        '<p:pos x="100" y="200"/>'
        "<p:text>Legacy feedback</p:text>"
        "</p:cm>"
        "</p:cmLst>"
    )

    src = zipfile.ZipFile(base, "r")
    names = src.namelist()
    ct = src.read("[Content_Types].xml").decode("utf-8")
    ct = ct.replace(
        "</Types>",
        '<Override PartName="/ppt/commentAuthors.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.'
        'presentationml.commentAuthors+xml"/>'
        '<Override PartName="/ppt/comments/comment1.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.'
        'presentationml.comments+xml"/>'
        "</Types>",
    )
    pres_rels_name = "ppt/_rels/presentation.xml.rels"
    pres_rels = src.read(pres_rels_name).decode("utf-8")
    pres_rels = pres_rels.replace(
        "</Relationships>",
        '<Relationship Id="rIdLegacyAuth" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
        'relationships/commentAuthors" Target="commentAuthors.xml"/>'
        "</Relationships>",
    )
    slide_rels_name = "ppt/slides/_rels/slide1.xml.rels"
    slide_rels = src.read(slide_rels_name).decode("utf-8")
    slide_rels = slide_rels.replace(
        "</Relationships>",
        '<Relationship Id="rIdLegacyCmt" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
        'relationships/comments" Target="../comments/comment1.xml"/>'
        "</Relationships>",
    )

    patched = {
        "[Content_Types].xml": ct,
        pres_rels_name: pres_rels,
        slide_rels_name: slide_rels,
        "ppt/commentAuthors.xml": legacy_authors,
        "ppt/comments/comment1.xml": legacy_comments,
    }
    out = io.BytesIO()
    dst = zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED)
    for name in names:
        dst.writestr(name, patched.get(name, src.read(name)))
    for name, data in patched.items():
        if name not in names:
            dst.writestr(name, data)
    src.close()
    dst.close()
    out.seek(0)

    prs = Presentation(out)
    slide = prs.slides[0]
    slide.comments.add("Modern reply-era feedback", "Modern Mary")
    return prs, slide


class DescribeLegacyModernCoexistence:
    """SF8 (ISC-42..46, ISC-67): legacy <p:cm> + modern <p188:cm> coexist."""

    def it_keeps_the_legacy_part_when_a_modern_comment_is_added(self):
        # ISC-43: adding a modern comment must NOT delete commentsN.xml.
        prs, slide = _legacy_plus_modern_deck()
        partnames = {p.partname for p in prs.part.package.iter_parts()}
        assert PackURI("/ppt/comments/comment1.xml") in partnames

    def it_enumerates_both_families_after_save_and_reopen(self):
        # ISC-44: slide.comments must read legacy <p:cm> AND modern
        # <p188:cm> and present both. ISC-67: legacy never silently dropped.
        prs, _slide = _legacy_plus_modern_deck()
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        comments = list(Presentation(stream).slides[0].comments)
        texts = sorted(c.text for c in comments)
        assert texts == ["Legacy feedback", "Modern reply-era feedback"]

    def it_exposes_is_legacy_to_distinguish_the_families(self):
        prs, slide = _legacy_plus_modern_deck()
        by_text = {c.text: c for c in slide.comments}
        assert by_text["Legacy feedback"].is_legacy is True
        assert by_text["Modern reply-era feedback"].is_legacy is False

    def it_does_not_mangle_legacy_author_ids_on_modern_write(self):
        # ISC-46 anti: modern write must not rewrite legacy <p:cm>/author ids.
        prs, _slide = _legacy_plus_modern_deck()
        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)
        import zipfile

        zf = zipfile.ZipFile(stream, "r")
        legacy_xml = zf.read("ppt/comments/comment1.xml").decode("utf-8")
        authors_xml = zf.read("ppt/commentAuthors.xml").decode("utf-8")
        zf.close()
        assert 'authorId="0"' in legacy_xml
        assert 'idx="1"' in legacy_xml
        assert 'id="0"' in authors_xml
        assert 'name="Legacy Larry"' in authors_xml

    def it_resolves_the_legacy_author_name(self):
        prs, slide = _legacy_plus_modern_deck()
        legacy = next(c for c in slide.comments if c.is_legacy)
        assert legacy.author == "Legacy Larry"

    def it_treats_a_legacy_comment_replies_as_empty_read_only(self):
        # legacy schema has no reply thread — SF4 documented legacy behavior.
        prs, slide = _legacy_plus_modern_deck()
        legacy = next(c for c in slide.comments if c.is_legacy)
        assert len(legacy.replies) == 0
        assert list(legacy.replies) == []


class DescribeThreadedCommentBodyPrRegression:
    """Regression: <p188:txBody> MUST carry <a:bodyPr> (issue #25 silent-drop).

    The add-path builds txBody via get_or_add_txBody() (a bare element);
    a CT_TextBody without <a:bodyPr> is schema-malformed and PowerPoint
    SILENTLY drops the comment (no repair dialog, comment just absent —
    caught only by maintainer visual review). Trinity was green while
    every comment was invisible. This locks the fix.
    """

    def _modern_xml(self):
        import io
        import zipfile

        from pptx import Presentation

        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[1])
        c = s.comments.add("Body check", author="Rev")
        c.replies.add("reply body check", author="Rev2")
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        z = zipfile.ZipFile(buf)
        name = next(n for n in z.namelist() if "modernComment" in n)
        return z.read(name).decode()

    def it_emits_bodyPr_in_every_comment_txBody(self):
        import re

        xml = self._modern_xml()
        bodies = re.findall(r"<p188:txBody>(.*?)</p188:txBody>", xml, re.S)
        assert bodies, "expected at least one <p188:txBody>"
        assert all("bodyPr" in b for b in bodies), (
            "every threaded-comment/reply txBody must contain <a:bodyPr> "
            "or PowerPoint silently drops the comment"
        )

    def it_places_bodyPr_before_the_paragraph(self):
        import re

        xml = self._modern_xml()
        body = re.search(r"<p188:txBody>(.*?)</p188:txBody>", xml, re.S).group(1)
        assert body.index("bodyPr") < body.index("<a:p"), (
            "<a:bodyPr> must precede <a:p> (CT_TextBody child order)"
        )


class DescribeThreadedCommentPowerPointContract:
    """Regression: the modern-comment OOXML must match PowerPoint ground truth.

    Captured 2026-05-17 from a threaded comment PowerPoint for Mac itself
    authored+saved. Three string axes plus a slide-binding element were
    inferred wrong in Waves 1-2 (content type ``threadedComments+xml``,
    reltypes ``threadedComment``/``threadedCommentAuthors``, and a missing
    ``<pc:sldMkLst>``). With any of them wrong PowerPoint silently fails to
    load/bind the part → empty Comments pane while the test trinity is
    green. This locks every axis to the PowerPoint-emitted contract.
    """

    def it_uses_the_powerpoint_content_type_and_reltypes(self):
        assert CT.PML_THREADED_COMMENTS == "application/vnd.ms-powerpoint.comments+xml"
        assert RT.THREADED_COMMENT == (
            "http://schemas.microsoft.com/office/2018/10/relationships/comments"
        )
        assert RT.AUTHORS == "http://schemas.microsoft.com/office/2018/10/relationships/authors"

    def _round_trip_zip(self):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.comments.add("Looks great! Ship it.", author="Alex Reviewer")
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        import zipfile

        return zipfile.ZipFile(buf), s.slide_id

    def it_stamps_the_powerpoint_content_type_in_content_types(self):
        z, _ = self._round_trip_zip()
        ctypes = z.read("[Content_Types].xml").decode()
        assert "application/vnd.ms-powerpoint.comments+xml" in ctypes
        assert "threadedComments+xml" not in ctypes

    def it_relates_slide_to_part_and_presentation_to_authors_by_pp_reltype(self):
        z, _ = self._round_trip_zip()
        slide_rels = z.read("ppt/slides/_rels/slide1.xml.rels").decode()
        prs_rels = z.read("ppt/_rels/presentation.xml.rels").decode()
        assert "/2018/10/relationships/comments" in slide_rels
        assert "/2018/10/relationships/threadedComment" not in slide_rels
        assert "/2018/10/relationships/authors" in prs_rels
        assert "threadedCommentAuthors" not in prs_rels

    def it_binds_each_comment_to_its_slide_via_sldMkLst(self):
        import re

        z, slide_id = self._round_trip_zip()
        name = next(n for n in z.namelist() if "modernComment" in n)
        xml = z.read(name).decode()
        # <pc:sldMkLst> present with <pc:docMk/> and <pc:sldMk sldId=...>
        assert "sldMkLst" in xml, "every <p188:cm> needs a <pc:sldMkLst> slide binding"
        assert "docMk" in xml
        assert "sldMk" in xml
        m = re.search(r'<pc:sldMk[^>]*sldId="(\d+)"', xml)
        assert m is not None, "<pc:sldMk> must carry the slide's sldId"
        assert int(m.group(1)) == slide_id, (
            "sldMk/@sldId must equal the slide's <p:sldId>/@id (%d) so "
            "PowerPoint binds the comment to the right slide" % slide_id
        )

    def it_places_sldMkLst_before_txBody_in_child_order(self):
        z, _ = self._round_trip_zip()
        name = next(n for n in z.namelist() if "modernComment" in n)
        xml = z.read(name).decode()
        assert xml.index("sldMkLst") < xml.index("txBody"), (
            "<pc:sldMkLst> must precede <p188:txBody> (PowerPoint child order)"
        )

    def it_keeps_exactly_one_sldMkLst_when_set_twice(self):
        # Cato low-finding lock: set_slide_marker must REPLACE, not prepend a
        # second <pc:sldMkLst>, if invoked again on the same comment element.
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[1])
        c = s.comments.add("once", author="Rev")
        c._cm.set_slide_marker(s.slide_id)
        c._cm.set_slide_marker(s.slide_id)
        from pptx.oxml.ns import qn

        assert len(c._cm.findall(qn("pc:sldMkLst"))) == 1, (
            "set_slide_marker must be idempotent (exactly one <pc:sldMkLst>)"
        )

    def it_places_replyLst_before_txBody_per_ms_pptx_sequence(self):
        # [MS-PPTX] CT_Comment sequence: EG_CommentAnchor, pos?, replyLst?,
        # EG_CommentProperties(txBody, extLst). replyLst MUST precede txBody.
        # Emitting replyLst AFTER txBody makes PowerPoint render the parent
        # comment but SILENTLY DROP every reply (issue #25, the reply defect).
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[1])
        c = s.comments.add("parent body", author="Rev")
        c.replies.add("a reply body", author="Rev2")
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        import zipfile

        z = zipfile.ZipFile(buf)
        name = next(n for n in z.namelist() if "modernComment" in n)
        xml = z.read(name).decode()
        assert "replyLst" in xml, "expected a <p188:replyLst> for the reply"
        assert "a reply body" in xml, "reply text must serialize"
        assert xml.index("replyLst") < xml.index("txBody"), (
            "<p188:replyLst> must precede <p188:txBody> per [MS-PPTX] "
            "CT_Comment sequence, or PowerPoint drops every reply"
        )

    def it_keeps_reply_txBody_with_bodyPr_after_reorder(self):
        # Guard the reorder didn't regress the reply body shape.
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[1])
        c = s.comments.add("p", author="Rev")
        c.replies.add("r", author="Rev2")
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        import re
        import zipfile

        z = zipfile.ZipFile(buf)
        name = next(n for n in z.namelist() if "modernComment" in n)
        xml = z.read(name).decode()
        reply = re.search(r"<p188:reply\b.*?</p188:reply>", xml, re.S).group(0)
        assert "bodyPr" in reply
        assert reply.index("bodyPr") < reply.index("<a:p")


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
