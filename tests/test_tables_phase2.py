# pyright: reportPrivateUsage=false

"""Unit-test suite for Tables 2.0 Phase 2 — table style API.

Covers:

- New oxml: `CT_TableStyleId` element class plus `CT_TableProperties.style_id`
  property (read/write/clear) routed through a `ZeroOrOne` `<a:tableStyleId>`
  child element.
- Public API on |Table|: `style_id` (r/w GUID), `style_name` (reverse lookup),
  `apply_style(name_or_guid)` accepting either a friendly name (case-insensitive
  registry lookup) or a raw brace-wrapped GUID.
- Built-in registry in `pptx.enum.table`: `PP_TABLE_STYLE`,
  `lookup_table_style`, `style_name_for`, `register_table_style`.
- Round-trip: save a presentation with a custom style applied, reopen it,
  assert the GUID survived.
- Anti-criteria: existing toggles (`first_row` etc.) still round-trip; the
  default style remains `Medium Style 2 - Accent 1` for newly-added tables;
  the GUID emitted always uses the canonical brace-wrapped upper-case-hex
  shape.

Issue: https://github.com/MHoroszowski/python-pptx/issues/12 (Phase 2).
"""

from __future__ import annotations

import io
import re

import pytest

from pptx import Presentation
from pptx.enum.table import (
    PP_TABLE_STYLE,
    lookup_table_style,
    register_table_style,
    style_name_for,
)
from pptx.oxml.table import CT_TableProperties, CT_TableStyleId
from pptx.table import Table, _looks_like_guid
from pptx.util import Inches

from .unitutil.cxml import element

DEFAULT_GUID = "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"
ALT_GUID = "{F5AB1C69-6EDB-4FF4-983F-18BD219EF322}"  # Medium Style 2 - Accent 3
NO_GRID_GUID = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid

CANONICAL_GUID_RE = re.compile(
    r"^\{[0-9A-F]{8}-[0-9A-F]{4}-[0-9A-F]{4}-[0-9A-F]{4}-[0-9A-F]{12}\}$"
)


# ---------------------------------------------------------------------------
# OXML LAYER — CT_TableStyleId + CT_TableProperties.style_id
# ---------------------------------------------------------------------------


class DescribeCT_TableStyleId(object):
    """Unit-test suite for `CT_TableStyleId.value` round-trip."""

    def it_reads_value_from_element_text(self):
        elm = element("a:tableStyleId")
        assert isinstance(elm, CT_TableStyleId)
        elm.text = DEFAULT_GUID

        assert elm.value == DEFAULT_GUID

    def it_writes_value_to_element_text(self):
        elm = element("a:tableStyleId")
        elm.value = DEFAULT_GUID

        assert elm.text == DEFAULT_GUID

    def it_returns_empty_string_when_no_text(self):
        elm = element("a:tableStyleId")
        assert elm.value == ""


class DescribeCT_TableProperties_StyleId(object):
    """Unit-test suite for `CT_TableProperties.style_id` property."""

    def it_reads_None_when_tableStyleId_absent(self):
        tblPr = element("a:tblPr")
        assert isinstance(tblPr, CT_TableProperties)
        assert tblPr.style_id is None

    def it_reads_guid_when_tableStyleId_present(self):
        tblPr = element("a:tblPr")
        styleId_elm = tblPr.get_or_add_tableStyleId()
        styleId_elm.text = DEFAULT_GUID

        assert tblPr.style_id == DEFAULT_GUID

    def it_creates_tableStyleId_on_set(self):
        tblPr = element("a:tblPr")

        tblPr.style_id = DEFAULT_GUID

        assert tblPr.tableStyleId is not None
        assert tblPr.tableStyleId.text == DEFAULT_GUID

    def it_updates_existing_tableStyleId_on_set(self):
        tblPr = element("a:tblPr")
        tblPr.style_id = DEFAULT_GUID

        tblPr.style_id = ALT_GUID

        assert tblPr.style_id == ALT_GUID

    def it_removes_tableStyleId_on_set_None(self):
        tblPr = element("a:tblPr")
        tblPr.style_id = DEFAULT_GUID
        assert tblPr.tableStyleId is not None

        tblPr.style_id = None

        assert tblPr.tableStyleId is None

    def it_is_a_no_op_to_clear_when_already_absent(self):
        tblPr = element("a:tblPr")
        assert tblPr.tableStyleId is None

        tblPr.style_id = None  # ---should not raise---

        assert tblPr.tableStyleId is None

    def it_emits_tableStyleId_as_first_child_of_tblPr(self):
        # ---ECMA-376 §21.1.3.15 sequence rule---
        tblPr = element("a:tblPr")
        tblPr.style_id = DEFAULT_GUID

        children = list(tblPr)

        assert len(children) == 1
        assert children[0].tag.endswith("}tableStyleId")

    def it_inserts_tableStyleId_before_extLst_when_tblPr_has_extLst(self):
        # ---regression for ECMA-376 §21.1.3.15: `tableStyleId` must come
        # ---BEFORE `extLst`. PowerPoint-authored decks may carry an
        # ---existing `<a:extLst>` on `<a:tblPr>`; setting style_id must not
        # ---append after it.
        tblPr = element("a:tblPr/a:extLst")

        tblPr.style_id = DEFAULT_GUID

        children = list(tblPr)
        assert len(children) == 2
        assert children[0].tag.endswith("}tableStyleId")
        assert children[1].tag.endswith("}extLst")


# ---------------------------------------------------------------------------
# REGISTRY — pptx.enum.table
# ---------------------------------------------------------------------------


class DescribePP_TABLE_STYLE_Registry(object):
    """Unit-test suite for `PP_TABLE_STYLE` and helper functions."""

    def it_includes_the_python_pptx_default_style(self):
        assert PP_TABLE_STYLE["Medium Style 2 - Accent 1"] == DEFAULT_GUID

    def it_includes_all_six_medium_style_2_accents(self):
        # ---scanny#27's 37 commenters most often asked about Medium Style 2 ---
        # ---and its accents; verify all six are present and unique GUIDs---
        guids = [PP_TABLE_STYLE["Medium Style 2 - Accent %d" % n] for n in range(1, 7)]
        assert len(set(guids)) == 6
        for g in guids:
            assert CANONICAL_GUID_RE.match(g), g

    def it_uses_canonical_brace_upper_hex_shape_for_every_guid(self):
        for name, guid in PP_TABLE_STYLE.items():
            assert CANONICAL_GUID_RE.match(guid), "%s: %r" % (name, guid)

    def it_has_no_duplicate_guid_entries(self):
        guids = list(PP_TABLE_STYLE.values())
        assert len(guids) == len(set(guids))

    def it_includes_the_no_style_no_grid_entry(self):
        assert PP_TABLE_STYLE["No Style, No Grid"] == NO_GRID_GUID


class Describe_lookup_table_style(object):
    """Unit-test suite for `lookup_table_style`."""

    def it_returns_guid_for_canonical_name(self):
        assert lookup_table_style("Medium Style 2 - Accent 1") == DEFAULT_GUID

    def it_is_case_insensitive(self):
        assert lookup_table_style("medium style 2 - accent 1") == DEFAULT_GUID
        assert lookup_table_style("MEDIUM STYLE 2 - ACCENT 1") == DEFAULT_GUID

    def it_raises_ValueError_on_unknown_name(self):
        with pytest.raises(ValueError) as excinfo:
            lookup_table_style("Bogus Style")
        assert "Bogus Style" in str(excinfo.value)


class Describe_style_name_for(object):
    """Unit-test suite for `style_name_for`."""

    def it_returns_friendly_name_for_known_guid(self):
        assert style_name_for(DEFAULT_GUID) == "Medium Style 2 - Accent 1"

    def it_returns_None_for_unknown_guid(self):
        assert style_name_for("{00000000-0000-0000-0000-000000000000}") is None


class Describe_register_table_style(object):
    """Unit-test suite for `register_table_style` extensibility."""

    def it_allows_runtime_extension(self):
        custom_guid = "{ABCDEF01-2345-6789-ABCD-EF0123456789}"
        register_table_style("CorpStyle Custom", custom_guid)
        try:
            assert lookup_table_style("CorpStyle Custom") == custom_guid
            assert style_name_for(custom_guid) == "CorpStyle Custom"
            assert lookup_table_style("corpstyle custom") == custom_guid
        finally:
            # ---test isolation: clean up the registry---
            del PP_TABLE_STYLE["CorpStyle Custom"]  # type: ignore[attr-defined]


# ---------------------------------------------------------------------------
# PUBLIC API — Table.style_id / style_name / apply_style
# ---------------------------------------------------------------------------


@pytest.fixture
def table_fixture():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    gf = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    return gf.table


class DescribeTable_StyleAPI(object):
    """Unit-test suite for `Table.style_id`, `Table.style_name`, `Table.apply_style`."""

    def it_returns_default_style_id_for_a_new_table(self, table_fixture):
        # ---fork's `_tbl_tmpl` bakes Medium Style 2 - Accent 1 as default---
        assert table_fixture.style_id == DEFAULT_GUID

    def it_returns_default_style_name_for_a_new_table(self, table_fixture):
        assert table_fixture.style_name == "Medium Style 2 - Accent 1"

    def it_can_set_style_id_directly(self, table_fixture):
        table_fixture.style_id = ALT_GUID

        assert table_fixture.style_id == ALT_GUID

    def it_can_clear_style_id_with_None(self, table_fixture):
        table_fixture.style_id = None

        assert table_fixture.style_id is None
        assert table_fixture.style_name is None

    def it_returns_None_style_name_when_guid_is_unregistered(self, table_fixture):
        table_fixture.style_id = "{00000000-0000-0000-0000-000000000000}"

        assert table_fixture.style_id == "{00000000-0000-0000-0000-000000000000}"
        assert table_fixture.style_name is None

    def it_applies_a_style_by_friendly_name(self, table_fixture):
        table_fixture.apply_style("Medium Style 2 - Accent 3")

        assert table_fixture.style_id == ALT_GUID
        assert table_fixture.style_name == "Medium Style 2 - Accent 3"

    def it_applies_a_style_by_raw_guid(self, table_fixture):
        table_fixture.apply_style(NO_GRID_GUID)

        assert table_fixture.style_id == NO_GRID_GUID
        assert table_fixture.style_name == "No Style, No Grid"

    def it_resolves_friendly_name_case_insensitively(self, table_fixture):
        table_fixture.apply_style("medium style 2 - accent 3")

        assert table_fixture.style_id == ALT_GUID

    def it_raises_ValueError_on_unknown_friendly_name(self, table_fixture):
        with pytest.raises(ValueError):
            table_fixture.apply_style("not-a-real-style")

    def it_lets_a_user_round_trip_a_custom_guid_losslessly(self, table_fixture):
        # ---a GUID not in the registry passes through verbatim, name returns None---
        custom_guid = "{12345678-1234-1234-1234-123456789ABC}"
        table_fixture.apply_style(custom_guid)

        assert table_fixture.style_id == custom_guid
        assert table_fixture.style_name is None


# ---------------------------------------------------------------------------
# Helper — _looks_like_guid
# ---------------------------------------------------------------------------


class Describe_looks_like_guid(object):
    """Unit-test suite for the GUID-shape detector used by apply_style."""

    @pytest.mark.parametrize(
        "value",
        [
            "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}",
            "{abcdef01-2345-6789-abcd-ef0123456789}",  # ---lowercase OK---
            "{ABCDEF01-2345-6789-ABCD-EF0123456789}",
        ],
    )
    def it_accepts_canonical_guid_shapes(self, value):
        assert _looks_like_guid(value) is True

    @pytest.mark.parametrize(
        "value",
        [
            "Medium Style 2 - Accent 1",  # ---name, not GUID---
            "5C22544A-7EE6-4342-B048-85BDC9FD1C3A",  # ---no braces---
            "{5C22544A}",  # ---wrong length---
            "",
            "{NOT-A-GU-ID00-0000-000000000000}",  # ---invalid hex---
        ],
    )
    def it_rejects_non_guid_strings(self, value):
        assert _looks_like_guid(value) is False


# ---------------------------------------------------------------------------
# Round-trip — save/reload preserves style_id
# ---------------------------------------------------------------------------


class DescribeStyle_RoundTrip(object):
    """Save a presentation with a non-default style and reload — GUID survives."""

    def it_preserves_style_id_through_save_and_reload(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        gf = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        gf.table.apply_style("Light Style 2 - Accent 4")
        applied_guid = gf.table.style_id

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        prs2 = Presentation(buf)
        # ---first slide, first shape is the table---
        tbl2 = next(shp for shp in prs2.slides[0].shapes if shp.has_table).table
        assert tbl2.style_id == applied_guid
        assert tbl2.style_name == "Light Style 2 - Accent 4"

    def it_preserves_a_cleared_style_through_save_and_reload(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        gf = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        gf.table.style_id = None

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        prs2 = Presentation(buf)
        tbl2 = next(shp for shp in prs2.slides[0].shapes if shp.has_table).table
        assert tbl2.style_id is None


# ---------------------------------------------------------------------------
# Anti / Regression
# ---------------------------------------------------------------------------


class DescribePhase2_Regression(object):
    """Anti-criteria: existing surfaces unchanged after Phase 2 additions."""

    def it_keeps_first_row_toggle_working(self, table_fixture):
        assert table_fixture.first_row is True  # ---template default---
        table_fixture.first_row = False
        assert table_fixture.first_row is False
        table_fixture.first_row = True
        assert table_fixture.first_row is True

    def it_keeps_horz_banding_toggle_working(self, table_fixture):
        assert table_fixture.horz_banding is True  # ---template default---
        table_fixture.horz_banding = False
        assert table_fixture.horz_banding is False

    def it_keeps_vert_banding_toggle_working(self, table_fixture):
        assert table_fixture.vert_banding is False
        table_fixture.vert_banding = True
        assert table_fixture.vert_banding is True

    def it_keeps_default_style_unchanged_for_newly_added_tables(self, table_fixture):
        # ---no regression in the existing _tbl_tmpl baked-in default---
        assert table_fixture.style_id == DEFAULT_GUID
        assert isinstance(table_fixture, Table)
