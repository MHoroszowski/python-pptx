"""Re-generate the synthetic .pptx fixtures used by the customXml integration tests.

Run from the repo root::

    python3 tests/test_files/customxml/_generate_fixtures.py

Outputs are deterministic except for auto-assigned `datastoreItem` GUIDs;
explicit GUIDs below pin them so the resulting files round-trip byte-for-byte
when re-generated.
"""

from __future__ import annotations

import os

from pptx import Presentation

_HERE = os.path.dirname(os.path.abspath(__file__))


def _path(name: str) -> str:
    return os.path.join(_HERE, name)


def write_presentation_scoped() -> None:
    prs = Presentation()
    prs.custom_xml_parts.add(
        b'<provenance xmlns="urn:my:provenance">'
        b"<source>integration-fixture</source>"
        b"<built-at>2026-05-05T17:00:00Z</built-at>"
        b"</provenance>",
        name="provenance",
        schema_refs=["urn:my:provenance"],
        datastoreItem_id="{1A2B3C4D-5E6F-7890-ABCD-EF1234567890}",
    )
    prs.save(_path("presentation-scoped.pptx"))


def write_package_scoped() -> None:
    prs = Presentation()
    prs.custom_xml_parts.add(
        b'<vsto-config xmlns="urn:my:vsto">'
        b"<template>quarterly-report</template>"
        b"</vsto-config>",
        name="vsto",
        scope="package",
        datastoreItem_id="{ABCDEF12-3456-7890-ABCD-EF1234567890}",
    )
    prs.save(_path("package-scoped.pptx"))


def write_multipart() -> None:
    prs = Presentation()
    prs.custom_properties["Source"] = "deck-builder-cli@1.4.2"
    prs.custom_properties["BuildNumber"] = 42
    prs.custom_properties["IsDraft"] = True
    prs.custom_xml_parts.add(
        b'<provenance xmlns="urn:my:p"><source>cli</source></provenance>',
        name="provenance",
        schema_refs=["urn:my:p"],
    )
    prs.custom_xml_parts.add(b"<extra/>", name="extra", scope="package")
    prs.custom_xml_parts.add_string_blob(
        "readme",
        "# Hello\n\nThis is markdown content embedded in the .pptx.",
        mime_hint="text/markdown",
    )
    prs.save(_path("multipart.pptx"))


def write_clean() -> None:
    Presentation().save(_path("clean.pptx"))


def main() -> None:
    write_presentation_scoped()
    write_package_scoped()
    write_multipart()
    write_clean()
    print("regenerated fixtures in", _HERE)


if __name__ == "__main__":
    main()
