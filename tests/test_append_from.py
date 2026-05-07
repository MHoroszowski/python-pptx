# pyright: reportPrivateUsage=false

"""Unit-test suite for `Presentation.append_from` (slide-CRUD Phase 3).

Issue: https://github.com/MHoroszowski/python-pptx/issues/11 (Phase 3 — cross-deck copy).

Implements the slide-CRUD epic's append_from sub-feature: copy slides between
TWO different python-pptx Presentation instances, with cross-package porting
of slide-layout / slide-master / theme and image-dedup at the target package.
"""

from __future__ import annotations

import io
from pathlib import Path

import pytest

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.image import ImagePart
from pptx.parts.slide import SlideMasterPart, SlidePart
from pptx.util import Inches

_PNG_PATH = Path(__file__).parent / "test_files" / "python-powered.png"


# ---------------------------------------------------------------------------
# Helpers — keep round-trip plumbing identical to test_slide_crud /
# test_slide_duplicate.
# ---------------------------------------------------------------------------


def _seed_presentation_with(n_slides: int) -> Presentation:
    prs = Presentation()
    layout = prs.slide_layouts[6]
    for i in range(n_slides):
        slide = prs.slides.add_slide(layout)
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.text = (
            "slide %d" % i
        )
    return prs


def _round_trip(prs: Presentation) -> Presentation:
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


def _image_part_count(prs: Presentation) -> int:
    return sum(1 for p in prs.part.package.iter_parts() if isinstance(p, ImagePart))


def _slide_part_count(prs: Presentation) -> int:
    return sum(1 for p in prs.part.package.iter_parts() if isinstance(p, SlidePart))


def _master_part_count(prs: Presentation) -> int:
    return sum(1 for p in prs.part.package.iter_parts() if isinstance(p, SlideMasterPart))


# ---------------------------------------------------------------------------
# API surface.
# ---------------------------------------------------------------------------


class DescribePresentation_AppendFrom_API(object):
    """Argument validation, signature, and return shape."""

    def it_exposes_append_from_on_Presentation(self):
        prs = Presentation()
        assert callable(getattr(prs, "append_from", None))

    def it_returns_a_list_of_Slide_objects(self):
        from pptx.slide import Slide

        src = _seed_presentation_with(2)
        tgt = Presentation()

        result = tgt.append_from(src)

        assert isinstance(result, list)
        assert all(isinstance(s, Slide) for s in result)
        assert len(result) == 2

    def it_appends_all_source_slides_when_slide_indexes_is_None(self):
        src = _seed_presentation_with(3)
        tgt = Presentation()
        tgt.slides.add_slide(tgt.slide_layouts[6])

        tgt.append_from(src)

        assert len(tgt.slides) == 4

    def it_appends_only_selected_indexes_when_slide_indexes_provided(self):
        src = _seed_presentation_with(3)
        tgt = Presentation()
        tgt.slides.add_slide(tgt.slide_layouts[6])

        tgt.append_from(src, slide_indexes=[0, 2])

        assert len(tgt.slides) == 3

    def it_preserves_source_index_order_in_appended_slides(self):
        src = _seed_presentation_with(3)
        tgt = Presentation()

        new_slides = tgt.append_from(src, slide_indexes=[2, 0])

        # Verify by text content of the textbox we seeded.
        def _first_text(slide):
            return next(
                shp.text_frame.text for shp in slide.shapes if getattr(shp, "has_text_frame", False)
            )

        assert _first_text(new_slides[0]) == "slide 2"
        assert _first_text(new_slides[1]) == "slide 0"

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def it_raises_IndexError_for_out_of_range_index(self, bad_idx: int):
        src = _seed_presentation_with(2)
        tgt = Presentation()
        with pytest.raises(IndexError):
            tgt.append_from(src, slide_indexes=[0, bad_idx])

    def it_treats_empty_slide_indexes_as_a_noop(self):
        src = _seed_presentation_with(3)
        tgt = Presentation()
        tgt.slides.add_slide(tgt.slide_layouts[6])
        slide_count_before = len(tgt.slides)

        result = tgt.append_from(src, slide_indexes=[])

        assert result == []
        assert len(tgt.slides) == slide_count_before

    def it_supports_self_append_from_self(self):
        """Edge case — appending from self to self acts like multi-duplicate."""
        prs = _seed_presentation_with(2)
        slide_count_before = len(prs.slides)

        new_slides = prs.append_from(prs, slide_indexes=[0])

        assert len(prs.slides) == slide_count_before + 1
        assert len(new_slides) == 1


# ---------------------------------------------------------------------------
# Cross-package porting: per-slide.
# ---------------------------------------------------------------------------


class DescribePresentation_AppendFrom_PerSlide(object):
    """Each appended slide gets its own deep-copied part."""

    def it_grows_target_slide_part_count_by_the_number_of_appended_slides(self):
        src = _seed_presentation_with(3)
        tgt = Presentation()
        tgt.slides.add_slide(tgt.slide_layouts[6])
        n_before = _slide_part_count(tgt)

        tgt.append_from(src)

        assert _slide_part_count(tgt) == n_before + 3

    def it_does_not_mutate_source_slide_count(self):
        """Anti — source presentation must be unchanged after append."""
        src = _seed_presentation_with(2)
        tgt = Presentation()
        src_slide_count_before = len(src.slides)

        tgt.append_from(src)

        assert len(src.slides) == src_slide_count_before

    def it_assigns_unique_partnames_across_appended_slides(self):
        src = _seed_presentation_with(3)
        tgt = Presentation()

        tgt.append_from(src)

        partnames = [s.part.partname for s in tgt.slides]
        assert len(set(partnames)) == len(partnames)

    def it_assigns_unique_slide_ids_across_appended_slides(self):
        src = _seed_presentation_with(2)
        tgt = Presentation()
        tgt.slides.add_slide(tgt.slide_layouts[6])

        tgt.append_from(src)

        ids = [s.slide_id for s in tgt.slides]
        assert len(set(ids)) == len(ids)

    def it_isolates_modifications_to_appended_slides_from_source(self):
        src = _seed_presentation_with(1)
        tgt = Presentation()

        new_slide = tgt.append_from(src)[0]
        textbox = next(shp for shp in new_slide.shapes if getattr(shp, "has_text_frame", False))
        textbox.text_frame.text = "MUTATED"

        src_text = next(
            shp.text_frame.text
            for shp in src.slides[0].shapes
            if getattr(shp, "has_text_frame", False)
        )
        assert src_text == "slide 0"


# ---------------------------------------------------------------------------
# Slide-master / slide-layout / theme porting.
# ---------------------------------------------------------------------------


class DescribePresentation_AppendFrom_MasterPort(object):
    """Source's master + layouts get ported into the target package."""

    def it_adds_a_new_master_to_target_when_appending_from_a_separate_pres(self):
        src = _seed_presentation_with(1)
        tgt = Presentation()
        n_masters_before = _master_part_count(tgt)

        tgt.append_from(src)

        assert _master_part_count(tgt) == n_masters_before + 1

    def it_dedups_master_within_a_single_call_for_slides_sharing_a_master(self):
        """Two source slides on same master → one ported master in target."""
        src = _seed_presentation_with(3)  # all three share src's single master
        tgt = Presentation()

        tgt.append_from(src)

        # target had 1 master, source has 1 master → after port, target has 2
        assert _master_part_count(tgt) == 2

    def it_does_not_dedup_master_across_separate_calls(self):
        """Two consecutive append_from calls re-port the master."""
        src = _seed_presentation_with(1)
        tgt = Presentation()

        tgt.append_from(src)
        tgt.append_from(src)

        # target started at 1, +1 per call = 3
        assert _master_part_count(tgt) == 3

    def it_appends_a_new_sldMasterId_entry_to_target_presentation(self):
        src = _seed_presentation_with(1)
        tgt = Presentation()
        n_master_id_entries_before = len(
            tgt.part._element.get_or_add_sldMasterIdLst().sldMasterId_lst
        )

        tgt.append_from(src)

        n_master_id_entries_after = len(
            tgt.part._element.get_or_add_sldMasterIdLst().sldMasterId_lst
        )
        assert n_master_id_entries_after == n_master_id_entries_before + 1

    def it_assigns_a_unique_master_id_to_the_new_sldMasterId_entry(self):
        src = _seed_presentation_with(1)
        tgt = Presentation()

        tgt.append_from(src)

        ids = [sm.get("id") for sm in tgt.part._element.get_or_add_sldMasterIdLst().sldMasterId_lst]
        assert all(i is not None for i in ids)
        assert len(set(ids)) == len(ids)

    def it_ports_all_layouts_of_the_source_master(self):
        """Master's layout tree is ported intact, not just the referenced one."""
        src = _seed_presentation_with(1)
        n_src_layouts = len(src.slide_master.slide_layouts)
        tgt = Presentation()
        n_tgt_layouts_before = sum(len(m.slide_layouts) for m in tgt.slide_masters)

        tgt.append_from(src)

        n_tgt_layouts_after = sum(len(m.slide_layouts) for m in tgt.slide_masters)
        assert n_tgt_layouts_after == n_tgt_layouts_before + n_src_layouts


# ---------------------------------------------------------------------------
# Image / media dedup at target package.
# ---------------------------------------------------------------------------


class DescribePresentation_AppendFrom_ImageDedup(object):
    """Cross-package SHA1 dedup: images shared by source AND target unify."""

    def it_dedups_a_source_image_already_present_in_target(self):
        """Same PNG used in both decks → one ImagePart in target after append."""
        src = Presentation()
        s1 = src.slides.add_slide(src.slide_layouts[6])
        s1.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))

        tgt = Presentation()
        t1 = tgt.slides.add_slide(tgt.slide_layouts[6])
        t1.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))
        n_images_before = _image_part_count(tgt)

        tgt.append_from(src)

        # Image IS the same blob → SHA1 dedup → no new ImagePart in target.
        assert _image_part_count(tgt) == n_images_before

    def it_dedups_within_a_single_call_when_two_source_slides_share_an_image(self):
        src = Presentation()
        s1 = src.slides.add_slide(src.slide_layouts[6])
        s1.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))
        s2 = src.slides.add_slide(src.slide_layouts[6])
        s2.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))

        tgt = Presentation()
        n_images_before = _image_part_count(tgt)

        tgt.append_from(src)

        # Both source slides referenced the same blob; target gets ONE new ImagePart.
        assert _image_part_count(tgt) == n_images_before + 1

    def it_round_trips_dedup_through_save_and_reopen(self):
        src = Presentation()
        s1 = src.slides.add_slide(src.slide_layouts[6])
        s1.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))

        tgt = Presentation()
        t1 = tgt.slides.add_slide(tgt.slide_layouts[6])
        t1.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))

        tgt.append_from(src)
        reopened = _round_trip(tgt)

        # Both slides in reopened present the same PNG bytes via shared ImagePart.
        slides_with_pictures = [
            [shp for shp in s.shapes if shp.shape_type == 13] for s in reopened.slides
        ]
        assert all(len(pics) == 1 for pics in slides_with_pictures)
        blobs = [pics[0].image.blob for pics in slides_with_pictures]
        assert blobs[0] == blobs[1]


# ---------------------------------------------------------------------------
# Notes-slide handling.
# ---------------------------------------------------------------------------


class DescribePresentation_AppendFrom_NotesSlide(object):
    """Notes-slide gets its own copy; back-rel rewired to new slide; notes-master shared."""

    def it_gives_appended_slide_its_own_notes_slide_part(self):
        src = _seed_presentation_with(1)
        src.slides[0].notes_slide.notes_text_frame.text = "speaker notes"
        tgt = Presentation()

        new_slide = tgt.append_from(src)[0]

        assert new_slide.has_notes_slide is True
        assert new_slide.notes_slide.part is not src.slides[0].notes_slide.part

    def it_carries_notes_text_to_the_appended_slide(self):
        src = _seed_presentation_with(1)
        src.slides[0].notes_slide.notes_text_frame.text = "speaker notes"
        tgt = Presentation()

        new_slide = tgt.append_from(src)[0]

        assert new_slide.notes_slide.notes_text_frame.text == "speaker notes"

    def it_rewires_the_notes_back_rel_to_the_new_slide(self):
        """Cross-package gotcha #961: notes-slide's RT.SLIDE points at NEW slide."""
        src = _seed_presentation_with(1)
        src.slides[0].notes_slide.notes_text_frame.text = "x"
        tgt = Presentation()

        new_slide = tgt.append_from(src)[0]

        new_notes_part = new_slide.notes_slide.part
        back_ref_target = new_notes_part.part_related_by(RT.SLIDE)
        assert back_ref_target is new_slide.part

    def it_uses_targets_existing_notes_master_not_a_port_of_source(self):
        """Notes-master is presentation-singleton — target's existing one is reused."""
        src = _seed_presentation_with(1)
        src.slides[0].notes_slide.notes_text_frame.text = "x"
        tgt = Presentation()
        # ---force-create target notes-master before append---
        target_notes_master_part = tgt.part.notes_master_part

        new_slide = tgt.append_from(src)[0]

        new_notes_master = new_slide.notes_slide.part.part_related_by(RT.NOTES_MASTER)
        assert new_notes_master is target_notes_master_part

    def it_does_not_create_a_notes_slide_when_source_has_none(self):
        src = _seed_presentation_with(1)
        # no notes-slide on source
        assert src.slides[0].part.has_notes_slide is False
        tgt = Presentation()

        new_slide = tgt.append_from(src)[0]

        assert new_slide.has_notes_slide is False


# ---------------------------------------------------------------------------
# Round-trip integration.
# ---------------------------------------------------------------------------


class DescribePresentation_AppendFrom_RoundTrip(object):
    """Open → append_from → save → reopen integration coverage."""

    def it_round_trips_a_basic_append_all(self):
        src = _seed_presentation_with(2)
        tgt = Presentation()
        tgt.slides.add_slide(tgt.slide_layouts[6])

        tgt.append_from(src)
        reopened = _round_trip(tgt)

        assert len(reopened.slides) == 3

    def it_round_trips_selective_append(self):
        src = _seed_presentation_with(3)
        tgt = Presentation()

        tgt.append_from(src, slide_indexes=[2, 0])
        reopened = _round_trip(tgt)

        # Verify by text content
        texts = []
        for s in reopened.slides:
            for shp in s.shapes:
                if getattr(shp, "has_text_frame", False):
                    texts.append(shp.text_frame.text)
                    break
            else:
                texts.append(None)
        assert texts == ["slide 2", "slide 0"]

    def it_round_trips_appended_slide_with_an_image(self):
        src = Presentation()
        s = src.slides.add_slide(src.slide_layouts[6])
        s.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))

        tgt = Presentation()
        tgt.append_from(src)
        reopened = _round_trip(tgt)

        pics = [shp for shp in reopened.slides[0].shapes if shp.shape_type == 13]
        assert len(pics) == 1
        assert len(pics[0].image.blob) > 0

    def it_round_trips_appended_slide_with_notes(self):
        src = _seed_presentation_with(1)
        src.slides[0].notes_slide.notes_text_frame.text = "round-trip notes"

        tgt = Presentation()
        tgt.append_from(src)
        reopened = _round_trip(tgt)

        assert reopened.slides[0].has_notes_slide is True
        assert reopened.slides[0].notes_slide.notes_text_frame.text == "round-trip notes"

    def it_preserves_master_count_through_round_trip(self):
        src = _seed_presentation_with(1)
        tgt = Presentation()

        tgt.append_from(src)
        masters_before = _master_part_count(tgt)
        reopened = _round_trip(tgt)

        assert _master_part_count(reopened) == masters_before


# ---------------------------------------------------------------------------
# Anti-criteria.
# ---------------------------------------------------------------------------


class DescribePresentation_AppendFrom_Antis(object):
    """Anti-criteria: properties that MUST NOT change after append_from."""

    def it_does_not_mutate_source_image_part_count(self):
        src = Presentation()
        s = src.slides.add_slide(src.slide_layouts[6])
        s.shapes.add_picture(str(_PNG_PATH), Inches(1), Inches(1))
        n_src_images_before = _image_part_count(src)

        tgt = Presentation()
        tgt.append_from(src)

        assert _image_part_count(src) == n_src_images_before

    def it_does_not_mutate_source_master_count(self):
        src = _seed_presentation_with(1)
        n_src_masters_before = _master_part_count(src)

        tgt = Presentation()
        tgt.append_from(src)

        assert _master_part_count(src) == n_src_masters_before

    def it_does_not_drop_existing_target_slides(self):
        src = _seed_presentation_with(1)
        tgt = _seed_presentation_with(2)
        existing_ids = [s.slide_id for s in tgt.slides]

        tgt.append_from(src)

        ids_after = [s.slide_id for s in tgt.slides]
        for prior_id in existing_ids:
            assert prior_id in ids_after

    def it_does_not_carry_comments_through_an_append(self):
        """Phase 3 scope: comments rels are dropped, same as Phase 2."""
        src = _seed_presentation_with(1)
        tgt = Presentation()

        new_slide = tgt.append_from(src)[0]

        with pytest.raises(KeyError):
            new_slide.part.part_related_by(RT.COMMENTS)
