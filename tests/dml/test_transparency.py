"""Unit-test suite for transparency functionality in `pptx.dml` module."""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.dml.color import ColorFormat
from pptx.dml.fill import FillFormat, _GradientStop, _NoFill, _SolidFill
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches

from ..oxml.unitdata.dml import a_alpha, a_solidFill, an_srgbClr


class DescribeColorFormatTransparency(object):
    """Unit-test suite for ColorFormat transparency property."""

    def it_knows_its_transparency_value(self, transparency_get_fixture):
        color_format, expected_transparency = transparency_get_fixture
        assert color_format.transparency == expected_transparency

    def it_can_set_its_transparency_value(self, transparency_set_fixture):
        color_format, transparency, expected_xml = transparency_set_fixture
        color_format.transparency = transparency
        assert color_format._xFill.xml == expected_xml

    def it_returns_zero_transparency_for_NoneColor(self, _NoneColor_color_format):
        assert _NoneColor_color_format.transparency == 0.0

    def it_raises_on_transparency_set_for_NoneColor(self, _NoneColor_color_format):
        with pytest.raises(ValueError):
            _NoneColor_color_format.transparency = 0.5

    def it_raises_on_assign_invalid_transparency_value(self, rgb_color_format):
        with pytest.raises(ValueError):
            rgb_color_format.transparency = 1.1
        with pytest.raises(ValueError):
            rgb_color_format.transparency = -0.1

    def it_can_set_transparency_to_zero_removes_alpha_element(self, rgb_color_format):
        rgb_color_format.transparency = 0.5
        assert rgb_color_format._color._xClr.alpha is not None

        rgb_color_format.transparency = 0.0
        assert rgb_color_format._color._xClr.alpha is None
        assert rgb_color_format.transparency == 0.0

    def it_can_set_transparency_to_one_creates_zero_alpha(self, rgb_color_format):
        rgb_color_format.transparency = 1.0

        alpha_element = rgb_color_format._color._xClr.alpha
        assert alpha_element is not None
        assert alpha_element.val == 0.0
        assert rgb_color_format.transparency == 1.0

    def it_validates_transparency_value_range(self, rgb_color_format):
        rgb_color_format._validate_transparency_value(0.0)
        rgb_color_format._validate_transparency_value(0.5)
        rgb_color_format._validate_transparency_value(1.0)

        with pytest.raises(ValueError, match="transparency must be number in range 0.0 to 1.0"):
            rgb_color_format._validate_transparency_value(-0.1)
        with pytest.raises(ValueError, match="transparency must be number in range 0.0 to 1.0"):
            rgb_color_format._validate_transparency_value(1.1)

    # fixtures -------------------------------------------------------

    @pytest.fixture(
        params=[
            (lambda: an_srgbClr().with_val("FF0000"), 0.0),
            (
                lambda: an_srgbClr().with_val("FF0000").with_child(a_alpha().with_val(100000)),
                0.0,
            ),
            (
                lambda: an_srgbClr().with_val("FF0000").with_child(a_alpha().with_val(50000)),
                0.5,
            ),
            (
                lambda: an_srgbClr().with_val("FF0000").with_child(a_alpha().with_val(0)),
                1.0,
            ),
        ]
    )
    def transparency_get_fixture(self, request):
        xClr_bldr_fn, expected_transparency = request.param
        xClr_bldr = xClr_bldr_fn()
        solidFill = a_solidFill().with_nsdecls().with_child(xClr_bldr).element
        color_format = ColorFormat.from_colorchoice_parent(solidFill)
        return color_format, expected_transparency

    @pytest.fixture(
        params=[
            (
                lambda: an_srgbClr().with_val("FF0000"),
                0.0,
                lambda: an_srgbClr().with_val("FF0000"),
            ),
            (
                lambda: an_srgbClr().with_val("FF0000"),
                0.5,
                lambda: an_srgbClr().with_val("FF0000").with_child(a_alpha().with_val(50000)),
            ),
            (
                lambda: an_srgbClr().with_val("FF0000"),
                1.0,
                lambda: an_srgbClr().with_val("FF0000").with_child(a_alpha().with_val(0)),
            ),
            (
                lambda: an_srgbClr().with_val("FF0000").with_child(a_alpha().with_val(75000)),
                0.3,
                lambda: an_srgbClr().with_val("FF0000").with_child(a_alpha().with_val(70000)),
            ),
        ]
    )
    def transparency_set_fixture(self, request):
        xClr_bldr_fn, transparency, expected_xClr_bldr_fn = request.param

        xClr_bldr = xClr_bldr_fn()
        solidFill = a_solidFill().with_nsdecls().with_child(xClr_bldr).element
        color_format = ColorFormat.from_colorchoice_parent(solidFill)

        expected_xClr_bldr = expected_xClr_bldr_fn()
        expected_xml = a_solidFill().with_nsdecls().with_child(expected_xClr_bldr).xml()

        return color_format, transparency, expected_xml

    @pytest.fixture
    def rgb_color_format(self):
        solidFill = a_solidFill().with_nsdecls().with_child(an_srgbClr().with_val("FF0000")).element
        return ColorFormat.from_colorchoice_parent(solidFill)

    @pytest.fixture
    def _NoneColor_color_format(self):
        solidFill = a_solidFill().with_nsdecls().element
        return ColorFormat.from_colorchoice_parent(solidFill)


class DescribeFillFormatTransparency(object):
    """Unit-test suite for FillFormat transparency property."""

    def it_delegates_transparency_to_solid_fill_object(self):
        xClr_bldr = an_srgbClr().with_val("FF0000")
        solidFill_elm = a_solidFill().with_nsdecls().with_child(xClr_bldr).element
        solid_fill = _SolidFill(solidFill_elm)

        class _MockParent:
            def __init__(self, fill_elm):
                self.eg_fillProperties = fill_elm

        fill_format = FillFormat(_MockParent(solidFill_elm), solid_fill)

        fill_format.transparency = 0.3

        assert abs(fill_format._fill.transparency - 0.3) < 0.001
        assert abs(fill_format.transparency - 0.3) < 0.001

    def it_raises_on_transparency_access_for_non_solid_fills(self):
        class _MockParent:
            eg_fillProperties = None

        fill_format = FillFormat(_MockParent(), _NoFill(None))

        with pytest.raises(TypeError, match="fill type .* has no transparency"):
            fill_format.transparency

    def it_raises_on_transparency_set_for_non_solid_fills(self):
        class _MockParent:
            eg_fillProperties = None

        fill_format = FillFormat(_MockParent(), _NoFill(None))

        with pytest.raises(TypeError, match="fill type .* has no transparency"):
            fill_format.transparency = 0.5


class DescribeSolidFillTransparency(object):
    """Unit-test suite for _SolidFill transparency property."""

    def it_provides_access_to_transparency(self, solid_fill_transparency_fixture):
        solid_fill, expected_transparency = solid_fill_transparency_fixture
        assert abs(solid_fill.transparency - expected_transparency) < 0.001

    def it_can_set_transparency(self, solid_fill_transparency_set_fixture):
        solid_fill, transparency = solid_fill_transparency_set_fixture
        solid_fill.transparency = transparency
        assert abs(solid_fill.transparency - transparency) < 0.001

    def it_delegates_transparency_to_fore_color(self, solid_fill_obj):
        solid_fill_obj.transparency = 0.6

        assert solid_fill_obj.fore_color.transparency == 0.6
        assert solid_fill_obj.transparency == 0.6

    def it_has_correct_fill_type(self, solid_fill_obj):
        assert solid_fill_obj.type == MSO_FILL.SOLID

    # fixtures -------------------------------------------------------

    @pytest.fixture(
        params=[
            (lambda: an_srgbClr().with_val("00FF00"), 0.0),
            (
                lambda: an_srgbClr().with_val("00FF00").with_child(a_alpha().with_val(80000)),
                0.2,
            ),
            (
                lambda: an_srgbClr().with_val("00FF00").with_child(a_alpha().with_val(30000)),
                0.7,
            ),
        ]
    )
    def solid_fill_transparency_fixture(self, request):
        xClr_bldr_fn, expected_transparency = request.param
        xClr_bldr = xClr_bldr_fn()
        solidFill_elm = a_solidFill().with_nsdecls().with_child(xClr_bldr).element

        return _SolidFill(solidFill_elm), expected_transparency

    @pytest.fixture(params=[0.0, 0.1, 0.33, 0.67, 0.9, 1.0])
    def solid_fill_transparency_set_fixture(self, request):
        transparency = request.param
        xClr_bldr = an_srgbClr().with_val("00FF00")
        solidFill_elm = a_solidFill().with_nsdecls().with_child(xClr_bldr).element

        return _SolidFill(solidFill_elm), transparency

    @pytest.fixture
    def solid_fill_obj(self):
        xClr_bldr = an_srgbClr().with_val("0000FF")
        solidFill_elm = a_solidFill().with_nsdecls().with_child(xClr_bldr).element

        return _SolidFill(solidFill_elm)


class DescribeTransparencyIntegration(object):
    """Integration tests for transparency across the entire stack."""

    def it_works_end_to_end_with_solid_fill_and_color_format(self):
        xClr_bldr = an_srgbClr().with_val("FF00FF")
        solidFill_elm = a_solidFill().with_nsdecls().with_child(xClr_bldr).element
        solid_fill = _SolidFill(solidFill_elm)

        solid_fill.transparency = 0.4

        assert abs(solid_fill.transparency - 0.4) < 0.001
        assert abs(solid_fill.fore_color.transparency - 0.4) < 0.001

        alpha_elm = solid_fill.fore_color._color._xClr.alpha
        assert alpha_elm is not None
        assert abs(alpha_elm.val - 0.6) < 0.001

    def it_handles_transparency_removal_correctly(self):
        xClr_bldr = an_srgbClr().with_val("FFFF00").with_child(a_alpha().with_val(40000))
        solidFill_elm = a_solidFill().with_nsdecls().with_child(xClr_bldr).element
        solid_fill = _SolidFill(solidFill_elm)

        assert abs(solid_fill.transparency - 0.6) < 0.001
        assert solid_fill.fore_color._color._xClr.alpha is not None

        solid_fill.transparency = 0.0

        assert solid_fill.transparency == 0.0
        assert solid_fill.fore_color._color._xClr.alpha is None

    def it_handles_color_format_transparency_directly(self):
        xClr_bldr = an_srgbClr().with_val("00FFFF")
        solidFill_elm = a_solidFill().with_nsdecls().with_child(xClr_bldr).element
        color_format = ColorFormat.from_colorchoice_parent(solidFill_elm)

        for transparency in (0.0, 0.25, 0.5, 0.75, 1.0):
            color_format.transparency = transparency
            assert abs(color_format.transparency - transparency) < 0.001

            if transparency == 0.0:
                assert color_format._color._xClr.alpha is None
            else:
                alpha_elm = color_format._color._xClr.alpha
                assert alpha_elm is not None
                expected_alpha = 1.0 - transparency
                assert abs(alpha_elm.val - expected_alpha) < 0.001


class DescribeGradientStopTransparency(object):
    """Unit-test suite for transparency on a single gradient stop.

    Covers issue #17 follow-up: `_GradientStop.color` returns a `ColorFormat`,
    so the PR #30 transparency surface should flow through transitively. These
    tests pin the behavior in so a future refactor of either side can't
    silently regress it.
    """

    def it_reads_zero_transparency_when_stop_has_no_alpha(self):
        gs = self._gs_with_srgb("FF0000")
        color_format = ColorFormat.from_colorchoice_parent(gs)

        assert color_format.transparency == 0.0

    def it_reads_transparency_from_alpha_under_srgbClr(self):
        gs = self._gs_with_srgb("FF0000", alpha_val=50000)
        color_format = ColorFormat.from_colorchoice_parent(gs)

        assert abs(color_format.transparency - 0.5) < 0.001

    def it_reads_transparency_from_alpha_under_schemeClr(self):
        gs = self._gs_with_scheme("accent1", alpha_val=30000)
        color_format = ColorFormat.from_colorchoice_parent(gs)

        assert abs(color_format.transparency - 0.7) < 0.001

    def it_writes_alpha_under_srgbClr_when_setting_transparency(self):
        gs = self._gs_with_srgb("FF0000")
        color_format = ColorFormat.from_colorchoice_parent(gs)

        color_format.transparency = 0.5

        srgbClr = gs.find(qn("a:srgbClr"))
        alpha = srgbClr.find(qn("a:alpha"))
        assert alpha is not None, "alpha must be placed under <a:srgbClr>, not <a:gs>"
        assert abs(alpha.val - 0.5) < 0.001

    def it_writes_alpha_under_schemeClr_when_setting_transparency(self):
        gs = self._gs_with_scheme("accent1")
        color_format = ColorFormat.from_colorchoice_parent(gs)

        color_format.transparency = 0.6

        schemeClr = gs.find(qn("a:schemeClr"))
        alpha = schemeClr.find(qn("a:alpha"))
        assert alpha is not None, "alpha must be placed under <a:schemeClr>, not <a:gs>"
        assert abs(alpha.val - 0.4) < 0.001

    def it_clears_alpha_when_setting_transparency_to_zero(self):
        gs = self._gs_with_srgb("FF0000", alpha_val=40000)
        color_format = ColorFormat.from_colorchoice_parent(gs)

        assert abs(color_format.transparency - 0.6) < 0.001

        color_format.transparency = 0.0

        srgbClr = gs.find(qn("a:srgbClr"))
        assert srgbClr.find(qn("a:alpha")) is None
        assert color_format.transparency == 0.0

    def it_raises_on_out_of_range_value(self):
        gs = self._gs_with_srgb("FF0000")
        color_format = ColorFormat.from_colorchoice_parent(gs)

        with pytest.raises(ValueError, match="transparency must be number in range 0.0 to 1.0"):
            color_format.transparency = 1.1
        with pytest.raises(ValueError, match="transparency must be number in range 0.0 to 1.0"):
            color_format.transparency = -0.1

    def it_raises_on_set_when_stop_has_no_color_choice(self):
        gs = parse_xml('<a:gs %s pos="50000"/>' % nsdecls("a"))
        color_format = ColorFormat.from_colorchoice_parent(gs)

        with pytest.raises(ValueError, match="can't set transparency when color.type is None"):
            color_format.transparency = 0.5

    def it_exposes_transparency_through_GradientStop_color(self):
        gs = self._gs_with_srgb("FF0000")
        stop = _GradientStop(gs)

        stop.color.transparency = 0.25

        assert abs(stop.color.transparency - 0.25) < 0.001
        srgbClr = gs.find(qn("a:srgbClr"))
        alpha = srgbClr.find(qn("a:alpha"))
        assert alpha is not None
        assert abs(alpha.val - 0.75) < 0.001

    # helpers --------------------------------------------------------

    @staticmethod
    def _gs_with_srgb(val, alpha_val=None):
        alpha = f'<a:alpha val="{alpha_val}"/>' if alpha_val is not None else ""
        xml = ('<a:gs %s pos="50000"><a:srgbClr val="%s">%s</a:srgbClr></a:gs>') % (
            nsdecls("a"),
            val,
            alpha,
        )
        return parse_xml(xml)

    @staticmethod
    def _gs_with_scheme(val, alpha_val=None):
        alpha = f'<a:alpha val="{alpha_val}"/>' if alpha_val is not None else ""
        xml = ('<a:gs %s pos="50000"><a:schemeClr val="%s">%s</a:schemeClr></a:gs>') % (
            nsdecls("a"),
            val,
            alpha,
        )
        return parse_xml(xml)


class DescribeGradientStopTransparencyRoundTrip(object):
    """End-to-end round-trip for gradient stop transparency.

    The pure XML tests above prove the API and OOXML placement. This class
    proves a real Presentation survives save+reopen with transparency intact
    on its gradient stops — the closest thing to the actual user surface
    short of a UAT visual check.
    """

    def it_round_trips_transparency_on_default_gradient_stops(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        shape.fill.gradient()
        stops = shape.fill.gradient_stops
        stops[0].color.transparency = 0.0
        stops[1].color.transparency = 0.5

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        prs2 = Presentation(buf)
        stops2 = prs2.slides[0].shapes[1].fill.gradient_stops
        assert stops2[0].color.transparency == 0.0
        assert abs(stops2[1].color.transparency - 0.5) < 0.001

    def it_persists_alpha_inside_color_choice_not_on_gs(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )
        shape.fill.gradient()
        shape.fill.gradient_stops[1].color.transparency = 0.4

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        shape2 = prs2.slides[0].shapes[1]

        gs_list = shape2.fill._xPr.findall(
            ".//" + qn("a:gradFill") + "/" + qn("a:gsLst") + "/" + qn("a:gs")
        )
        assert len(gs_list) == 2
        gs_with_alpha = gs_list[1]
        # alpha must live under the color choice element, NOT directly on <a:gs>
        assert gs_with_alpha.find(qn("a:alpha")) is None
        color_choice = gs_with_alpha.find(qn("a:schemeClr"))
        if color_choice is None:
            color_choice = gs_with_alpha.find(qn("a:srgbClr"))
        assert color_choice is not None
        alpha = color_choice.find(qn("a:alpha"))
        assert alpha is not None
        assert abs(alpha.val - 0.6) < 0.001
