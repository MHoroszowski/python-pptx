"""Fork-specific behavioral tests for ChartEx (`cx:`) modern-chart support.

Covers the surface added on top of the GetThematic port: the `add_chart`
dispatch shim, the extended `XL_CHART_TYPE` members, round-trip preservation,
content-type / relationship wiring, formerly-deferred types now being
writable (Phase C), and the `WaterfallChartData` data API. See issue #14.
"""

from __future__ import annotations

import io
import zipfile

import pytest

from pptx import Presentation
from pptx.chart.data import WaterfallChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

_CHARTEX_TYPES = [
    "WATERFALL",
    "TREEMAP",
    "SUNBURST",
    "FUNNEL",
    "BOX_WHISKER",
    "HISTOGRAM",
    "PARETO",
]


def _slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[5])


def _waterfall_data():
    cd = WaterfallChartData()
    cd.categories = ["Q1", "Q2", "Q3", "Q4", "Total"]
    cd.add_series("Revenue", [100, 50, -30, 80, 200], subtotals=[4])
    return cd


def _save_reopen(prs):
    buf = io.BytesIO()
    prs.save(buf)
    blob = buf.getvalue()
    return blob, Presentation(io.BytesIO(blob))


class DescribeXlChartTypeChartExMembers:
    @pytest.mark.parametrize("name", _CHARTEX_TYPES)
    def it_exposes_each_chartex_member(self, name):
        assert hasattr(XL_CHART_TYPE, name)

    @pytest.mark.parametrize("name", _CHARTEX_TYPES)
    def it_assigns_unique_values_to_chartex_members(self, name):
        member = getattr(XL_CHART_TYPE, name)
        clashes = [m for m in XL_CHART_TYPE if m is not member and int(m) == int(member)]
        assert clashes == []

    def it_places_chartex_values_outside_the_ms_range(self):
        for name in _CHARTEX_TYPES:
            assert int(getattr(XL_CHART_TYPE, name)) >= 1000


class DescribeWaterfallChartData:
    def it_round_trips_categories(self):
        cd = WaterfallChartData()
        cd.categories = ["a", "b", "c"]
        assert cd.categories == ["a", "b", "c"]

    def it_records_series_name_values_and_subtotals(self):
        cd = WaterfallChartData()
        cd.add_series("S", [1, 2, 3], subtotals=[2])
        assert cd.series_name == "S"
        assert cd.series_values == [1, 2, 3]
        assert cd.subtotals == [2]

    def it_defaults_subtotals_to_empty_list(self):
        cd = WaterfallChartData()
        cd.add_series("S", [1, 2])
        assert cd.subtotals == []

    def it_computes_excel_refs_from_category_count(self):
        cd = WaterfallChartData()
        cd.categories = ["a", "b", "c", "d"]
        assert cd.categories_ref == "Sheet1!$A$2:$A$5"
        assert cd.values_ref == "Sheet1!$B$2:$B$5"
        assert cd.series_name_ref == "Sheet1!$B$1"

    def it_builds_an_xlsx_blob(self):
        cd = _waterfall_data()
        blob = cd.xlsx_blob
        assert blob[:2] == b"PK"  # zip magic
        assert zipfile.ZipFile(io.BytesIO(blob)).namelist()  # valid archive

    def it_raises_when_categories_and_values_mismatch(self):
        cd = WaterfallChartData()
        cd.categories = ["a", "b"]
        cd.add_series("S", [1])
        with pytest.raises(ValueError, match="must equal"):
            _ = cd.xlsx_blob


class DescribeAddChartDispatch:
    def it_dispatches_WATERFALL_to_the_chartex_path(self):
        _, slide = _slide()
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            _waterfall_data(),
        )
        assert gf.has_chartex is True
        assert gf.has_chart is False

    def it_keeps_classic_charts_on_the_c_path(self):
        from pptx.chart.data import CategoryChartData

        _, slide = _slide()
        cd = CategoryChartData()
        cd.categories = ["a", "b"]
        cd.add_series("S", (1, 2))
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(3),
            cd,
        )
        assert gf.has_chart is True
        assert gf.has_chartex is False

    @pytest.mark.parametrize(
        "name", ["TREEMAP", "SUNBURST", "FUNNEL", "BOX_WHISKER", "HISTOGRAM", "PARETO"]
    )
    def it_now_writes_formerly_deferred_types(self, name):
        # Phase C (issue #14) inverted the Phase-A/B contract: these types no
        # longer raise NotImplementedError — they are writable. Each is
        # exercised in depth in test_chartex_phasec.py; here we just assert
        # add_chart no longer raises for them.
        from pptx.chart.data import (
            BoxWhiskerChartData,
            FunnelChartData,
            HistogramChartData,
            ParetoChartData,
            SunburstChartData,
            TreemapChartData,
        )

        _, slide = _slide()
        if name in ("TREEMAP", "SUNBURST"):
            cd = {"TREEMAP": TreemapChartData, "SUNBURST": SunburstChartData}[name]()
            cd.add_level(["A", "B"])
            cd.add_level(["x", "y"])
            cd.add_series("S", [1, 2])
        elif name in ("FUNNEL", "BOX_WHISKER"):
            cd = {"FUNNEL": FunnelChartData, "BOX_WHISKER": BoxWhiskerChartData}[name]()
            cd.categories = ["a", "b"]
            cd.add_series("S", [1, 2])
        else:
            cd = {"HISTOGRAM": HistogramChartData, "PARETO": ParetoChartData}[name]()
            cd.add_series("S", [1, 2, 3, 4], bin_count=2)
        gf = slide.shapes.add_chart(
            getattr(XL_CHART_TYPE, name), Inches(1), Inches(1), Inches(5), Inches(3), cd
        )
        assert gf.has_chartex is True

    def it_can_add_via_add_chartex_directly(self):
        _, slide = _slide()
        gf = slide.shapes.add_chartex(_waterfall_data(), Inches(1), Inches(1), Inches(6), Inches(4))
        assert gf.has_chartex is True

    def it_has_an_empty_writer_deferred_set_after_phase_c(self):
        import inspect

        from pptx.shapes.shapetree import _BaseGroupShapes

        body = inspect.getsource(_BaseGroupShapes.add_chart)
        assert "_CHARTEX_WRITER_DEFERRED = ()" in body


class DescribeChartExRoundTrip:
    def it_writes_a_chartex_part_into_the_package(self):
        prs, slide = _slide()
        slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            _waterfall_data(),
        )
        blob, _ = _save_reopen(prs)
        names = zipfile.ZipFile(io.BytesIO(blob)).namelist()
        assert any("chartEx" in n and n.endswith(".xml") for n in names)

    def it_declares_the_chartex_content_type(self):
        prs, slide = _slide()
        slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            _waterfall_data(),
        )
        blob, _ = _save_reopen(prs)
        ct = zipfile.ZipFile(io.BytesIO(blob)).read("[Content_Types].xml").decode()
        assert "chartex+xml" in ct

    def it_relates_the_slide_to_the_chartex_part(self):
        prs, slide = _slide()
        slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            _waterfall_data(),
        )
        blob, _ = _save_reopen(prs)
        z = zipfile.ZipFile(io.BytesIO(blob))
        rels = z.read("ppt/slides/_rels/slide1.xml.rels").decode()
        assert "chartEx" in rels

    def it_preserves_the_chartex_part_when_unrelated_slide_changes(self):
        prs, slide = _slide()
        slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            _waterfall_data(),
        )
        blob1, prs2 = _save_reopen(prs)
        before = {n for n in zipfile.ZipFile(io.BytesIO(blob1)).namelist() if "chartEx" in n}
        prs2.slides.add_slide(prs2.slide_layouts[6])
        blob2, _ = _save_reopen(prs2)
        after = {n for n in zipfile.ZipFile(io.BytesIO(blob2)).namelist() if "chartEx" in n}
        assert before
        assert before == after

    def it_reads_back_has_chartex_on_the_reopened_shape(self):
        prs, slide = _slide()
        slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            _waterfall_data(),
        )
        _, prs2 = _save_reopen(prs)
        shapes = [s for s in prs2.slides[0].shapes if getattr(s, "has_chartex", False)]
        assert len(shapes) == 1

    def it_exposes_a_chartex_proxy_via_the_graphic_frame(self):
        prs, slide = _slide()
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            _waterfall_data(),
        )
        assert gf.chartex is not None

    def it_emits_a_waterfall_series_layout(self):
        prs, slide = _slide()
        slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            _waterfall_data(),
        )
        blob, _ = _save_reopen(prs)
        z = zipfile.ZipFile(io.BytesIO(blob))
        cx = [
            n for n in z.namelist() if "chartEx" in n and n.endswith(".xml") and "_rels" not in n
        ][0]
        assert "waterfall" in z.read(cx).decode()

    def it_keeps_a_classic_and_a_chartex_chart_together(self):
        from pptx.chart.data import CategoryChartData

        prs, slide = _slide()
        cd = CategoryChartData()
        cd.categories = ["a", "b"]
        cd.add_series("S", (1, 2))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(3),
            cd,
        )
        slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(5),
            Inches(1),
            Inches(4),
            Inches(3),
            _waterfall_data(),
        )
        blob, _ = _save_reopen(prs)
        names = zipfile.ZipFile(io.BytesIO(blob)).namelist()
        has_c = any("charts/chart1.xml" in n for n in names)
        has_cx = any("chartEx" in n and n.endswith(".xml") for n in names)
        assert has_c
        assert has_cx

    def it_does_not_inject_a_chartex_part_into_a_plain_deck(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        blob, _ = _save_reopen(prs)
        names = zipfile.ZipFile(io.BytesIO(blob)).namelist()
        assert not any("chartEx" in n for n in names)
