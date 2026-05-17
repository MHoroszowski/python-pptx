"""Phase-C tests: ChartEx writers for Treemap/Sunburst/Funnel/BoxWhisker/
Histogram/Pareto + the generalized ``ChartEx.replace_data``. See issue #14.

Structural assertions are derived from the normative in-repo schema
``spec/ISO-IEC-29500-4/xsd/dml-chartex.xsd`` (the schema PowerPoint conforms
to) — element ``layoutId`` + ``CT_Series`` child order, never hand-guessed.
"""

from __future__ import annotations

import io
import zipfile

import pytest
from lxml import etree

from pptx import Presentation
from pptx.chart.data import (
    BoxWhiskerChartData,
    FunnelChartData,
    HistogramChartData,
    ParetoChartData,
    SunburstChartData,
    TreemapChartData,
    WaterfallChartData,
)
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

CX = "http://schemas.microsoft.com/office/drawing/2014/chartex"


def _slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[5])


def _hier(cls):
    cd = cls()
    cd.add_level(["Tech", "Tech", "Retail", "Retail"])
    cd.add_level(["Phones", "Laptops", "Apparel", "Food"])
    cd.add_series("Revenue", [50, 30, 20, 15])
    return cd


def _cat(cls):
    cd = cls()
    cd.categories = ["Leads", "Qualified", "Proposals", "Won"]
    cd.add_series("Pipeline", [1000, 600, 250, 90])
    return cd


def _hist(cls):
    cd = cls()
    cd.add_series("Scores", [55, 62, 71, 73, 88, 91, 64, 78, 82, 69], bin_count=5)
    return cd


def _pareto(cls):
    # PowerPoint Pareto is categorical (aggregate by category) per ground truth.
    cd = cls()
    cd.categories = ["Defect A", "Defect B", "Defect C", "Defect D"]
    cd.add_series("Count", [45, 30, 15, 10])
    return cd


_BUILDERS = {
    XL_CHART_TYPE.TREEMAP: (lambda: _hier(TreemapChartData), "treemap"),
    XL_CHART_TYPE.SUNBURST: (lambda: _hier(SunburstChartData), "sunburst"),
    XL_CHART_TYPE.FUNNEL: (lambda: _cat(FunnelChartData), "funnel"),
    XL_CHART_TYPE.BOX_WHISKER: (lambda: _cat(BoxWhiskerChartData), "boxWhisker"),
    XL_CHART_TYPE.HISTOGRAM: (lambda: _hist(HistogramChartData), "clusteredColumn"),
    XL_CHART_TYPE.PARETO: (lambda: _pareto(ParetoChartData), "clusteredColumn"),
}
_ALL = list(_BUILDERS)


def _add(slide, ct):
    builder, _ = _BUILDERS[ct]
    return slide.shapes.add_chart(ct, Inches(1), Inches(1), Inches(6), Inches(4), builder())


def _cx_xml(prs):
    buf = io.BytesIO()
    prs.save(buf)
    z = zipfile.ZipFile(io.BytesIO(buf.getvalue()))
    name = next(
        n for n in z.namelist() if "chartEx" in n and n.endswith(".xml") and "_rels" not in n
    )
    return z.read(name), buf.getvalue()


class DescribePhaseCWriters:
    @pytest.mark.parametrize("ct", _ALL)
    def it_writes_each_type_via_add_chart(self, ct):
        _, slide = _slide()
        gf = _add(slide, ct)
        assert gf.has_chartex is True
        assert gf.has_chart is False

    @pytest.mark.parametrize("ct", _ALL)
    def it_emits_the_correct_layoutId(self, ct):
        prs, slide = _slide()
        _add(slide, ct)
        xml, _ = _cx_xml(prs)
        _, expected_layout = _BUILDERS[ct]
        assert ('layoutId="%s"' % expected_layout).encode() in xml

    @pytest.mark.parametrize("ct", _ALL)
    def it_round_trips_each_type_c14n_stable(self, ct):
        prs, slide = _slide()
        _add(slide, ct)
        xml1, blob = _cx_xml(prs)
        prs2 = Presentation(io.BytesIO(blob))
        prs2.slides.add_slide(prs2.slide_layouts[6])
        xml2, _ = _cx_xml(prs2)
        c14n = lambda b: etree.tostring(etree.fromstring(b), method="c14n2")  # noqa: E731
        assert c14n(xml1) == c14n(xml2)
        rt = [s for s in prs2.slides[0].shapes if getattr(s, "has_chartex", False)]
        assert len(rt) == 1

    @pytest.mark.parametrize("ct", _ALL)
    def it_emits_well_formed_xml(self, ct):
        prs, slide = _slide()
        _add(slide, ct)
        xml, _ = _cx_xml(prs)
        assert etree.fromstring(xml) is not None

    @pytest.mark.parametrize("ct", _ALL)
    def it_respects_CT_Series_child_order(self, ct):
        # XSD CT_Series sequence: tx, (spPr), ..., dataLabels, dataId, layoutPr.
        # Order bugs are the silent-corruption class — assert tx precedes
        # dataId precedes layoutPr on the first series.
        prs, slide = _slide()
        _add(slide, ct)
        xml, _ = _cx_xml(prs)
        root = etree.fromstring(xml)
        series = root.findall(".//{%s}series" % CX)[0]
        tags = [etree.QName(c).localname for c in series]
        assert tags.index("tx") < tags.index("dataId")
        if "layoutPr" in tags:
            assert tags.index("dataId") < tags.index("layoutPr")

    @pytest.mark.parametrize("ct", _ALL)
    def it_keeps_every_dataId_referencing_an_existing_data(self, ct):
        prs, slide = _slide()
        _add(slide, ct)
        xml, _ = _cx_xml(prs)
        root = etree.fromstring(xml)
        data_ids = {d.get("id") for d in root.findall(".//{%s}data" % CX)}
        for di in root.findall(".//{%s}dataId" % CX):
            assert di.get("val") in data_ids

    @pytest.mark.parametrize("ct", _ALL)
    def it_embeds_an_xlsx_workbook(self, ct):
        prs, slide = _slide()
        _add(slide, ct)
        buf = io.BytesIO()
        prs.save(buf)
        z = zipfile.ZipFile(io.BytesIO(buf.getvalue()))
        assert any(n.endswith(".xlsx") for n in z.namelist())

    def it_makes_pareto_emit_a_second_paretoLine_series(self):
        prs, slide = _slide()
        _add(slide, XL_CHART_TYPE.PARETO)
        xml, _ = _cx_xml(prs)
        assert b'layoutId="clusteredColumn"' in xml
        assert b'layoutId="paretoLine"' in xml

    def it_makes_histogram_emit_binning(self):
        prs, slide = _slide()
        _add(slide, XL_CHART_TYPE.HISTOGRAM)
        xml, _ = _cx_xml(prs)
        assert b"<cx:binning" in xml or b"binning" in xml

    @pytest.mark.parametrize("ct", [XL_CHART_TYPE.TREEMAP, XL_CHART_TYPE.SUNBURST])
    def it_emits_hierarchical_levels_with_correct_ptCount(self, ct):
        prs, slide = _slide()
        _add(slide, ct)
        xml, _ = _cx_xml(prs)
        root = etree.fromstring(xml)
        strDim = root.find(".//{%s}strDim" % CX)
        lvls = strDim.findall("{%s}lvl" % CX)
        assert len(lvls) == 2  # two add_level() calls
        for lvl in lvls:
            pts = lvl.findall("{%s}pt" % CX)
            assert lvl.get("ptCount") == str(len(pts))


class DescribePhaseCDataContainers:
    def it_builds_a_hierarchical_xlsx(self):
        cd = _hier(TreemapChartData)
        assert cd.xlsx_blob[:2] == b"PK"

    def it_rejects_ragged_hierarchy_levels(self):
        cd = TreemapChartData()
        cd.add_level(["a", "b"])
        cd.add_level(["x"])
        cd.add_series("S", [1, 2])
        with pytest.raises(ValueError, match="same length"):
            _ = cd.xlsx_blob

    def it_builds_a_category_xlsx(self):
        cd = _cat(FunnelChartData)
        assert cd.xlsx_blob[:2] == b"PK"

    def it_rejects_mismatched_category_value_lengths(self):
        cd = FunnelChartData()
        cd.categories = ["a", "b"]
        cd.add_series("S", [1])
        with pytest.raises(ValueError, match="must equal"):
            _ = cd.xlsx_blob

    def it_rejects_both_bin_count_and_bin_size(self):
        cd = HistogramChartData()
        with pytest.raises(ValueError, match="only one"):
            cd.add_series("S", [1, 2, 3], bin_count=3, bin_size=1.0)

    def it_exposes_cx_chart_type_discriminators(self):
        assert TreemapChartData.cx_chart_type == "treemap"
        assert SunburstChartData.cx_chart_type == "sunburst"
        assert FunnelChartData.cx_chart_type == "funnel"
        assert BoxWhiskerChartData.cx_chart_type == "boxWhisker"
        assert HistogramChartData.cx_chart_type == "histogram"
        assert ParetoChartData.cx_chart_type == "pareto"


class DescribePhaseCReplaceData:
    def it_replaces_treemap_data(self):
        prs, slide = _slide()
        gf = _add(slide, XL_CHART_TYPE.TREEMAP)
        nd = TreemapChartData()
        nd.add_level(["New", "New"])
        nd.add_level(["Alpha", "Beta"])
        nd.add_series("S2", [99, 11])
        gf.chartex.replace_data(nd)
        xml, _ = _cx_xml(prs)
        assert b"Alpha" in xml
        assert b"Phones" not in xml

    def it_replaces_funnel_data(self):
        prs, slide = _slide()
        gf = _add(slide, XL_CHART_TYPE.FUNNEL)
        nd = FunnelChartData()
        nd.categories = ["X", "Y"]
        nd.add_series("F2", [7, 3])
        gf.chartex.replace_data(nd)
        xml, _ = _cx_xml(prs)
        assert b'<cx:pt idx="0">X</cx:pt>' in xml or b">X<" in xml

    def it_replaces_histogram_data(self):
        prs, slide = _slide()
        gf = _add(slide, XL_CHART_TYPE.HISTOGRAM)
        nd = HistogramChartData()
        nd.add_series("H2", [1, 2, 3, 4, 5, 6], bin_count=3)
        gf.chartex.replace_data(nd)
        xml, _ = _cx_xml(prs)
        assert b'layoutId="clusteredColumn"' in xml

    def it_keeps_waterfall_replace_data_working(self):
        prs, slide = _slide()
        wd = WaterfallChartData()
        wd.categories = ["a", "b"]
        wd.add_series("W", [1, 2], subtotals=[1])
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.WATERFALL, Inches(1), Inches(1), Inches(6), Inches(4), wd
        )
        nw = WaterfallChartData()
        nw.categories = ["c", "d", "e"]
        nw.add_series("W", [5, 6, 7], subtotals=[2])
        gf.chartex.replace_data(nw)
        xml, _ = _cx_xml(prs)
        assert b">c<" in xml
        assert b'<cx:idx val="2"/>' in xml

    def it_does_not_change_the_part_name_or_rel_on_replace(self):
        prs, slide = _slide()
        gf = _add(slide, XL_CHART_TYPE.SUNBURST)
        part_before = gf.chartex.part.partname
        rId_before = gf._element.chartex_rId
        nd = SunburstChartData()
        nd.add_level(["P", "P"])
        nd.add_level(["m", "n"])
        nd.add_series("S", [4, 5])
        gf.chartex.replace_data(nd)
        assert gf.chartex.part.partname == part_before
        assert gf._element.chartex_rId == rId_before

    def it_raises_on_chart_type_mismatch(self):
        prs, slide = _slide()
        gf = _add(slide, XL_CHART_TYPE.FUNNEL)
        with pytest.raises(ValueError, match="cannot change chart type"):
            gf.chartex.replace_data(_hist(HistogramChartData))

    def it_round_trips_replaced_data(self):
        prs, slide = _slide()
        gf = _add(slide, XL_CHART_TYPE.TREEMAP)
        nd = TreemapChartData()
        nd.add_level(["Z", "Z"])
        nd.add_level(["q1", "q2"])
        nd.add_series("S", [8, 9])
        gf.chartex.replace_data(nd)
        _, blob = _cx_xml(prs)
        prs2 = Presentation(io.BytesIO(blob))
        xml2, _ = _cx_xml(prs2)
        assert b"q1" in xml2


class DescribePhaseCAntiCriteria:
    def it_leaves_the_writer_deferred_set_empty(self):
        from pptx.shapes.shapetree import _BaseGroupShapes

        src = _BaseGroupShapes.add_chart.__code__
        # The empty tuple literal must be present in add_chart source.
        import inspect

        body = inspect.getsource(_BaseGroupShapes.add_chart)
        assert "_CHARTEX_WRITER_DEFERRED = ()" in body
        assert src is not None

    @pytest.mark.parametrize("ct", _ALL)
    def it_does_not_raise_NotImplementedError_for_any_cx_type(self, ct):
        _, slide = _slide()
        # Must not raise — every cx: type is writable in Phase C.
        _add(slide, ct)

    @pytest.mark.parametrize("ct", _ALL)
    def it_wires_the_full_packaging_path_for_each_type(self, ct):
        # Advisor blind-spot closure: assert the graphicFrame→part packaging
        # (not just the chartEx XML) for every Phase-C type — content-type
        # declared, slide relationship present, graphicData URI is chartEx.
        prs, slide = _slide()
        gf = _add(slide, ct)
        buf = io.BytesIO()
        prs.save(buf)
        z = zipfile.ZipFile(io.BytesIO(buf.getvalue()))
        ctypes = z.read("[Content_Types].xml").decode()
        assert "chartex+xml" in ctypes
        rels = z.read("ppt/slides/_rels/slide1.xml.rels").decode()
        assert "chartEx" in rels
        uri = gf._element.graphic.graphicData.get("uri")
        assert uri == "http://schemas.microsoft.com/office/drawing/2014/chartex"

    def it_keeps_classic_c_charts_on_the_c_path(self):
        from pptx.chart.data import CategoryChartData

        _, slide = _slide()
        cd = CategoryChartData()
        cd.categories = ["a", "b"]
        cd.add_series("S", (1, 2))
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), cd
        )
        assert gf.has_chart is True
        assert gf.has_chartex is False

    def it_does_not_inject_chartex_into_a_plain_deck(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        buf = io.BytesIO()
        prs.save(buf)
        names = zipfile.ZipFile(io.BytesIO(buf.getvalue())).namelist()
        assert not any("chartEx" in n for n in names)
