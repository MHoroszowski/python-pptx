# pyright: reportPrivateUsage=false

"""Unit-test suite for Modernization Phase 1 (issue #29).

Phase 1 of the Modernization & Ergonomics epic. Bundles:

- ``pathlib.Path`` (and any ``os.PathLike``) inputs accepted everywhere
  python-pptx accepted str paths — closes upstream PR #1123.
- ``MSO_PATTERN_TYPE.PERCENT_40`` typo fix (was ``ERCENT_40``) — closes #1131.
- ``Slide.background.element`` returns the actual ``<p:bg>`` element
  rather than its parent ``<p:cSld>`` — closes upstream issue #1126.

Deferred to Phase 2: ``Font.color`` no-mutate-on-read fix (closes
#1111/#1074), ``collections.abc`` import sweep, dev-tooling
modernization (uv / pyright strict).
"""

from __future__ import annotations

import io
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_PATTERN_TYPE
from pptx.parts.image import Image
from pptx.util import Inches

# ---------------------------------------------------------------------------
# PathLike support — Presentation, save, add_picture, Image.from_file
# ---------------------------------------------------------------------------


class DescribePathLikeForPresentationOpen(object):
    def it_accepts_a_pathlib_Path(self, tmp_path):
        # ---seed a deck so we have something to read back---
        seed = Presentation()
        seed.save(str(tmp_path / "seed.pptx"))

        # ---Path goes in, Presentation comes out without TypeError---
        prs = Presentation(tmp_path / "seed.pptx")
        assert len(prs.slides) == 0  # ---fresh deck

    def it_accepts_a_str_path_unchanged(self, tmp_path):
        seed = Presentation()
        seed.save(str(tmp_path / "seed.pptx"))

        prs = Presentation(str(tmp_path / "seed.pptx"))
        assert prs is not None

    def it_accepts_a_BytesIO_stream_unchanged(self, tmp_path):
        buf = io.BytesIO()
        Presentation().save(buf)
        buf.seek(0)

        prs = Presentation(buf)
        assert prs is not None

    def it_accepts_a_subclass_of_PathLike(self, tmp_path):
        class MyPath:
            def __init__(self, p):
                self._p = p

            def __fspath__(self):
                return str(self._p)

        Presentation().save(str(tmp_path / "seed.pptx"))
        prs = Presentation(MyPath(tmp_path / "seed.pptx"))
        assert prs is not None


class DescribePathLikeForPresentationSave(object):
    def it_accepts_a_pathlib_Path(self, tmp_path):
        prs = Presentation()
        target = tmp_path / "out.pptx"

        prs.save(target)

        assert target.exists()
        assert target.stat().st_size > 0

    def it_round_trips_through_Path_save_and_Path_open(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        target = tmp_path / "out.pptx"

        prs.save(target)
        rt = Presentation(target)

        assert len(rt.slides) == 1


class DescribePathLikeForAddPicture(object):
    def it_accepts_a_pathlib_Path(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        png_path = Path(__file__).parent / "test_files" / "python-powered.png"

        picture = slide.shapes.add_picture(png_path, Inches(1), Inches(1))

        assert picture is not None


class DescribePathLikeForImageFromFile(object):
    def it_accepts_a_pathlib_Path(self):
        png = Path(__file__).parent / "test_files" / "python-powered.png"

        image = Image.from_file(png)

        assert image is not None
        assert len(image.blob) > 0

    def it_accepts_a_str_path_unchanged(self):
        png = Path(__file__).parent / "test_files" / "python-powered.png"

        image = Image.from_file(str(png))

        assert image is not None

    def it_accepts_a_file_like_object_unchanged(self):
        png = Path(__file__).parent / "test_files" / "python-powered.png"
        with png.open("rb") as f:
            image = Image.from_file(f)

        assert image is not None
        assert len(image.blob) > 0


# ---------------------------------------------------------------------------
# PERCENT_40 typo fix
# ---------------------------------------------------------------------------


class DescribePERCENT_40_TypoFix(object):
    def it_exposes_PERCENT_40_with_correct_spelling(self):
        # ---the previous broken name `ERCENT_40` should not exist anymore
        assert hasattr(MSO_PATTERN_TYPE, "PERCENT_40") is True
        assert hasattr(MSO_PATTERN_TYPE, "ERCENT_40") is False

    def it_carries_the_correct_xml_value_and_label(self):
        member = MSO_PATTERN_TYPE.PERCENT_40
        # ---values are (value, xml_value, description) tuples on this enum
        assert member.value == 6
        assert member.xml_value == "pct40"


# ---------------------------------------------------------------------------
# Slide.background._element / .element returns <p:bg>
# ---------------------------------------------------------------------------


class DescribeSlideBackgroundElement(object):
    def it_returns_the_bg_element_not_cSld(self):
        from lxml import etree

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg_proxy = slide.background

        # ---both `.element` (public) and `._element` (private) point at <p:bg>
        public_tag = etree.QName(bg_proxy.element.tag).localname
        private_tag = etree.QName(bg_proxy._element.tag).localname
        assert public_tag == "bg", f"expected 'bg', got '{public_tag}'"
        assert private_tag == "bg", f"expected 'bg', got '{private_tag}'"

    def it_still_supports_setting_a_solid_fill_via_fill(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()

        # ---accessing .fill should not have raised; the fill is now solid
        from pptx.enum.dml import MSO_FILL

        assert bg.fill.type == MSO_FILL.SOLID

    def it_round_trips_a_solid_fill_through_save_and_open(self, tmp_path):
        from pptx.dml.color import RGBColor

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0xAB, 0xCD, 0xEF)

        target = tmp_path / "bg.pptx"
        prs.save(target)
        rt = Presentation(target)

        assert rt.slides[0].background.fill.fore_color.rgb == RGBColor(0xAB, 0xCD, 0xEF)
