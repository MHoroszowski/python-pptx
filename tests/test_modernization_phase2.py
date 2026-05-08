# pyright: reportPrivateUsage=false

"""Unit-test suite for Modernization Phase 2 — bug fixes + by_name ergonomics.

Covers:

- |Font|.color getter is non-mutating (closes scanny/python-pptx#1111, #1074):
  reading `font.color.rgb` / `.type` / etc. on an unstyled run does NOT
  insert ``<a:solidFill>`` into the underlying XML. Setting
  `font.color.rgb = ...` still works (lazy materialization on first SET).
- W3CDTF datetime parser returns tz-aware datetimes when the source string
  carries timezone info, naive datetimes when it doesn't. Setter
  normalizes tz-aware inputs to UTC before serialization. Closes
  scanny/python-pptx#957.
- |_BaseShapes|.by_name(name) lookup helper returns the first shape with
  matching name or raises KeyError. Closes scanny/python-pptx#798,
  scanny/python-pptx#309, scanny/python-pptx#532.
- Anti-criteria: existing `font.color.rgb = ...` setter unchanged; existing
  Phase-1 fixes still in place.

Issue: https://github.com/MHoroszowski/python-pptx/issues/29 (Phase 2).
"""

from __future__ import annotations

import datetime as dt

import pytest
from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.oxml.coreprops import CT_CoreProperties
from pptx.text.text import _LazyFontColorFormat
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


def _make_run():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    tf = slide.shapes.title.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "x"
    return prs, run


@pytest.fixture
def run_fixture():
    _, r = _make_run()
    return r


@pytest.fixture
def slide_fixture():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    return slide


# ---------------------------------------------------------------------------
# Font.color non-mutation (scanny#1111, #1074)
# ---------------------------------------------------------------------------


class DescribeFontColor_NonMutation(object):
    """Reading `font.color` properties does NOT modify the underlying XML."""

    def it_does_not_mutate_rPr_on_color_property_access(self, run_fixture):
        rPr = run_fixture._r.get_or_add_rPr()
        before = etree.tostring(rPr)

        _ = run_fixture.font.color  # property access alone

        after = etree.tostring(rPr)
        assert before == after

    def it_does_not_mutate_rPr_on_reading_color_type(self, run_fixture):
        rPr = run_fixture._r.get_or_add_rPr()
        before = etree.tostring(rPr)

        _ = run_fixture.font.color.type

        after = etree.tostring(rPr)
        assert before == after

    def it_does_not_mutate_rPr_on_reading_color_rgb(self, run_fixture):
        rPr = run_fixture._r.get_or_add_rPr()
        before = etree.tostring(rPr)

        _ = run_fixture.font.color.rgb

        after = etree.tostring(rPr)
        assert before == after

    def it_returns_None_for_color_type_on_unstyled_run(self, run_fixture):
        assert run_fixture.font.color.type is None

    def it_returns_None_for_color_rgb_on_unstyled_run(self, run_fixture):
        assert run_fixture.font.color.rgb is None

    def it_returns_a_LazyFontColorFormat_proxy(self, run_fixture):
        assert isinstance(run_fixture.font.color, _LazyFontColorFormat)

    def it_materializes_solidFill_on_first_rgb_setter_call(self, run_fixture):
        run_fixture.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

        rPr = run_fixture._r.get_or_add_rPr()
        # ---a:solidFill should now exist as a child---
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        assert rPr.find("%ssolidFill" % ns) is not None

    def it_round_trips_rgb_after_lazy_materialization(self, run_fixture):
        run_fixture.font.color.rgb = RGBColor(0xAB, 0xCD, 0xEF)

        assert run_fixture.font.color.rgb == RGBColor(0xAB, 0xCD, 0xEF)
        assert run_fixture.font.color.type == MSO_COLOR_TYPE.RGB

    def it_materializes_solidFill_on_theme_color_setter(self, run_fixture):
        run_fixture.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

        assert run_fixture.font.color.theme_color == MSO_THEME_COLOR.ACCENT_1
        assert run_fixture.font.color.type == MSO_COLOR_TYPE.SCHEME

    def it_returns_None_for_theme_color_on_unstyled_run(self, run_fixture):
        # ---no solidFill = inherit from style; None is the right signal.
        # ---NOT_THEME_COLOR is reserved for the explicit "solidFill exists,
        # ---no schemeClr" case — conflating them would break inheritance
        # ---if a caller reads then writes back the value.
        assert run_fixture.font.color.theme_color is None

    def it_returns_None_for_brightness_on_unstyled_run(self, run_fixture):
        # ---0.0 is a real settable "no brightness adjustment" value, so
        # ---None must be the inherit signal on a fillless run.
        assert run_fixture.font.color.brightness is None

    def it_returns_None_for_transparency_on_unstyled_run(self, run_fixture):
        # ---0.0 means "fully opaque" — also a real settable value.
        assert run_fixture.font.color.transparency is None

    def it_keeps_existing_solidFill_runs_unchanged(self, run_fixture):
        # ---first set establishes a solidFill---
        run_fixture.font.color.rgb = RGBColor(0x00, 0xFF, 0x00)
        rPr = run_fixture._r.get_or_add_rPr()
        before = etree.tostring(rPr)

        # ---reading after first set should also be byte-stable---
        _ = run_fixture.font.color.rgb
        _ = run_fixture.font.color.type

        after = etree.tostring(rPr)
        assert before == after


# ---------------------------------------------------------------------------
# UTC-aware datetime parser/setter (scanny#957)
# ---------------------------------------------------------------------------


class DescribeW3CDTF_DateTime(object):
    """`_parse_W3CDTF_to_datetime` returns tz-aware datetimes when source carries tz info."""

    def it_returns_utc_for_Z_suffix(self):
        result = CT_CoreProperties._parse_W3CDTF_to_datetime("2024-01-15T10:30:00Z")
        assert result.tzinfo == dt.timezone.utc

    def it_returns_fixed_offset_for_negative_offset(self):
        result = CT_CoreProperties._parse_W3CDTF_to_datetime("2024-01-15T10:30:00-08:00")
        assert result.tzinfo == dt.timezone(dt.timedelta(hours=-8))

    def it_returns_fixed_offset_for_positive_offset(self):
        result = CT_CoreProperties._parse_W3CDTF_to_datetime("2024-01-15T10:30:00+05:30")
        assert result.tzinfo == dt.timezone(dt.timedelta(hours=5, minutes=30))

    def it_returns_naive_datetime_when_source_has_no_offset(self):
        result = CT_CoreProperties._parse_W3CDTF_to_datetime("2024-01-15T10:30:00")
        assert result.tzinfo is None

    def it_returns_naive_for_date_only_strings(self):
        result = CT_CoreProperties._parse_W3CDTF_to_datetime("2024-01-15")
        assert result.tzinfo is None
        assert result == dt.datetime(2024, 1, 15)

    def it_preserves_the_correct_instant_with_offset(self):
        # ---10:30 PST = 18:30 UTC---
        result = CT_CoreProperties._parse_W3CDTF_to_datetime("2024-01-15T10:30:00-08:00")
        as_utc = result.astimezone(dt.timezone.utc)
        assert as_utc.hour == 18
        assert as_utc.minute == 30

    def it_round_trips_through_save_and_reload(self, tmp_path):
        prs = Presentation()
        # ---tz-aware input: PDT (UTC-7); 10:00 PDT = 17:00 UTC---
        prs.core_properties.created = dt.datetime(
            2024, 7, 4, 10, 0, 0, tzinfo=dt.timezone(dt.timedelta(hours=-7))
        )
        out = tmp_path / "rt.pptx"
        prs.save(out)

        prs2 = Presentation(out)
        reloaded = prs2.core_properties.created
        # ---written as Z (UTC); reloaded value represents 17:00 UTC---
        assert reloaded.tzinfo == dt.timezone.utc
        assert reloaded.hour == 17
        assert reloaded.minute == 0

    def it_accepts_naive_datetimes_for_backwards_compat(self, tmp_path):
        prs = Presentation()
        # ---naive datetime; written as Z (treated as UTC by convention)---
        prs.core_properties.created = dt.datetime(2024, 1, 15, 10, 30, 0)
        out = tmp_path / "rt2.pptx"
        prs.save(out)

        prs2 = Presentation(out)
        reloaded = prs2.core_properties.created
        # ---reloaded value is tz-aware (Z parsed)---
        assert reloaded.tzinfo == dt.timezone.utc
        assert reloaded.hour == 10


# ---------------------------------------------------------------------------
# Shapes.by_name (scanny#798, #309, #532)
# ---------------------------------------------------------------------------


class DescribeShapes_by_name(object):
    """`shapes.by_name(name)` lookup helper."""

    def it_returns_the_shape_with_matching_name(self, slide_fixture):
        # ---the title shape on a Title+Content layout is named 'Title 1'---
        title = slide_fixture.shapes.by_name("Title 1")
        assert title.name == "Title 1"

    def it_raises_KeyError_on_no_match(self, slide_fixture):
        with pytest.raises(KeyError) as excinfo:
            slide_fixture.shapes.by_name("Bogus")
        assert "Bogus" in str(excinfo.value)

    def it_is_case_sensitive(self, slide_fixture):
        with pytest.raises(KeyError):
            slide_fixture.shapes.by_name("title 1")  # lowercase t

    def it_returns_first_match_when_multiple_share_a_name(self, slide_fixture):
        # ---add a textbox with a name colliding with a placeholder---
        tb = slide_fixture.shapes.add_textbox(Inches(2), Inches(2), Inches(2), Inches(1))
        tb.name = "Title 1"  # ---collide with the existing title's name---

        # ---first match in document order: the original title placeholder---
        result = slide_fixture.shapes.by_name("Title 1")
        # ---identity check: same _element as the placeholder---
        assert result.shape_id == slide_fixture.shapes.title.shape_id

    def it_works_on_slide_layout_shapes(self):
        # ---the layout shapes inherit by_name through _BaseShapes---
        prs = Presentation()
        layout = prs.slide_layouts[1]
        # ---layouts also have a 'Title 1' placeholder---
        title = layout.shapes.by_name("Title 1")
        assert title.name == "Title 1"

    def it_works_on_slide_master_shapes(self):
        prs = Presentation()
        master = prs.slide_master
        # ---masters typically have 'Title Placeholder 1'---
        # ---guard against fixture variation: just confirm by_name works on at least one shape---
        names = [s.name for s in master.shapes]
        if names:
            sh = master.shapes.by_name(names[0])
            assert sh.name == names[0]


# ---------------------------------------------------------------------------
# Anti / Regression
# ---------------------------------------------------------------------------


class DescribePhase2_Regression(object):
    """Anti-criteria — existing surfaces unchanged."""

    def it_keeps_font_color_setter_working_unchanged(self, run_fixture):
        run_fixture.font.color.rgb = RGBColor(0x12, 0x34, 0x56)
        assert run_fixture.font.color.rgb == RGBColor(0x12, 0x34, 0x56)

    def it_keeps_phase1_PathLike_working(self, tmp_path):
        # ---Phase 1's PathLike fix should still work---
        prs = Presentation()
        out = tmp_path / "x.pptx"
        prs.save(out)
        prs2 = Presentation(out)
        assert len(prs2.slides) == 0

    def it_keeps_phase1_PERCENT_40_typo_fix(self):
        from pptx.enum.dml import MSO_PATTERN_TYPE

        # ---Phase 1 fixed ERCENT_40 -> PERCENT_40---
        assert MSO_PATTERN_TYPE.PERCENT_40.value == 6
