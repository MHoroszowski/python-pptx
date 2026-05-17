# pyright: reportPrivateUsage=false

"""Unit + round-trip tests for issue #18.

[Epic] Shape Effects, Arrowheads & Connectors —
https://github.com/MHoroszowski/python-pptx/issues/18

Covers the genuinely-new surface (glow / reflection / soft-edge effects,
scene_3d / shape_3d, group world-space coordinates, flip_vertical /
flip_horizontal, Shape.duplicate) plus the issue-named convenience aliases
over the already-shipped arrowhead / connector API.

Layered like `tests/test_slide_duplicate.py`:
1. API-surface unit tests (build a shape, set a property, assert XML).
2. Round-trip integration tests (save → reopen → re-read) — the only
   layer that proves Office-compatible packaging and that nothing is
   silently dropped (the §6a thesis: schema-valid can still be wrong).
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_END_SIZE, MSO_LINE_END_TYPE
from pptx.util import Emu, Inches, Pt

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _blank_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _rect(slide):
    from pptx.enum.shapes import MSO_SHAPE

    return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))


def _roundtrip_first_autoshape(prs):
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    prs2 = Presentation(buf)
    for shp in prs2.slides[0].shapes:
        if shp.shape_type is not None and "AUTO_SHAPE" in str(shp.shape_type):
            return prs2, shp
    return prs2, list(prs2.slides[0].shapes)[0]


def _effectLst(shape):
    return shape._element.spPr.find(f"{A}effectLst")


# ───────────────────────── A. Glow (SF1) ─────────────────────────


class DescribeGlowEffect:
    def it_creates_a_glow_with_color_and_radius(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.glow_effect.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        sp.shadow.glow_effect.radius = Pt(20)
        glow = _effectLst(sp).find(f"{A}glow")
        assert glow is not None
        assert glow.get("rad") == str(Pt(20))
        assert glow.find(f"{A}srgbClr").get("val") == "FF0000"

    def it_reads_back_none_when_no_glow(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        assert sp.shadow.glow_effect.radius is None
        assert sp.shadow.glow_effect.visible is False

    def it_round_trips_glow_through_save_reopen(self):
        prs, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.glow_effect.color.rgb = RGBColor(0x00, 0x80, 0xFF)
        sp.shadow.glow_effect.radius = Pt(15)
        _, sp2 = _roundtrip_first_autoshape(prs)
        assert sp2.shadow.glow_effect.radius == Emu(Pt(15))
        assert sp2.shadow.glow_effect.color.rgb == RGBColor(0x00, 0x80, 0xFF)

    def it_orders_glow_before_inner_and_outer_shadow(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.visible = True  # outerShdw
        sp.shadow.glow_effect.radius = Pt(10)
        kids = [el.tag for el in _effectLst(sp)]
        assert kids.index(f"{A}glow") < kids.index(f"{A}outerShdw")


# ───────────────────────── B. Reflection (SF2) ─────────────────────────


class DescribeReflectionEffect:
    def it_creates_a_reflection_with_attrs(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.reflection_effect.blur_radius = Pt(2)
        sp.shadow.reflection_effect.distance = Pt(5)
        sp.shadow.reflection_effect.direction = 90.0
        refl = _effectLst(sp).find(f"{A}reflection")
        assert refl is not None
        assert refl.get("blurRad") == str(Pt(2))
        assert refl.get("dist") == str(Pt(5))

    def it_reads_none_when_absent(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        assert sp.shadow.reflection_effect.blur_radius is None
        assert sp.shadow.reflection_effect.visible is False

    def it_round_trips_reflection(self):
        prs, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.reflection_effect.blur_radius = Pt(3)
        sp.shadow.reflection_effect.distance = Pt(7)
        _, sp2 = _roundtrip_first_autoshape(prs)
        assert sp2.shadow.reflection_effect.blur_radius == Emu(Pt(3))
        assert sp2.shadow.reflection_effect.distance == Emu(Pt(7))

    def it_orders_reflection_after_outer_shadow(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.visible = True
        sp.shadow.reflection_effect.blur_radius = Pt(2)
        kids = [el.tag for el in _effectLst(sp)]
        assert kids.index(f"{A}outerShdw") < kids.index(f"{A}reflection")


# ───────────────────────── C. Soft edge (SF3) ─────────────────────────


class DescribeSoftEdgeEffect:
    def it_creates_a_soft_edge_with_radius(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.soft_edge_effect.radius = Pt(4)
        se = _effectLst(sp).find(f"{A}softEdge")
        assert se is not None
        assert se.get("rad") == str(Pt(4))

    def it_reads_none_when_absent(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        assert sp.shadow.soft_edge_effect.radius is None

    def it_round_trips_soft_edge(self):
        prs, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.soft_edge_effect.radius = Pt(6)
        _, sp2 = _roundtrip_first_autoshape(prs)
        assert sp2.shadow.soft_edge_effect.radius == Emu(Pt(6))

    def it_emits_soft_edge_last_in_schema_order(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.glow_effect.radius = Pt(3)
        sp.shadow.soft_edge_effect.radius = Pt(3)
        kids = [el.tag for el in _effectLst(sp)]
        assert kids[-1] == f"{A}softEdge"


# ───────────────────────── D. 3-D scene / shape (SF4) ─────────────────────────


class DescribeScene3DAndShape3D:
    def it_sets_a_preset_camera_with_required_lightrig(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.scene_3d.camera_preset = "perspectiveRelaxedModerately"
        s3d = sp._element.spPr.find(f"{A}scene3d")
        assert s3d is not None
        assert s3d.find(f"{A}camera").get("prst") == "perspectiveRelaxedModerately"
        # ---lightRig is schema-required; absence => PowerPoint repair---
        assert s3d.find(f"{A}lightRig") is not None

    def it_sets_extrusion_and_contour(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.shape_3d.extrusion_height = Pt(12)
        sp.shape_3d.contour_width = Pt(1)
        sp3d = sp._element.spPr.find(f"{A}sp3d")
        assert sp3d.get("extrusionH") == str(Pt(12))
        assert sp3d.get("contourW") == str(Pt(1))

    def it_round_trips_3d(self):
        prs, slide = _blank_slide()
        sp = _rect(slide)
        sp.scene_3d.camera_preset = "orthographicFront"
        sp.shape_3d.extrusion_height = Pt(20)
        _, sp2 = _roundtrip_first_autoshape(prs)
        assert sp2.scene_3d.camera_preset == "orthographicFront"
        assert sp2.shape_3d.extrusion_height == Emu(Pt(20))

    def it_orders_scene3d_and_sp3d_after_effectlst(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.visible = True
        sp.scene_3d.camera_preset = "orthographicFront"
        sp.shape_3d.extrusion_height = Pt(5)
        kids = [el.tag for el in sp._element.spPr]
        assert kids.index(f"{A}effectLst") < kids.index(f"{A}scene3d")
        assert kids.index(f"{A}scene3d") < kids.index(f"{A}sp3d")


# ───────────────────────── E. Arrowhead alias (SF5) ─────────────────────────


class DescribeHeadTailEndAlias:
    def it_exposes_head_end_type_width_length(self):
        _, slide = _blank_slide()
        c = slide.shapes.add_connector(2, Inches(1), Inches(1), Inches(4), Inches(1))
        c.line.head_end.type = MSO_LINE_END_TYPE.TRIANGLE
        c.line.head_end.width = MSO_LINE_END_SIZE.LARGE
        c.line.head_end.length = MSO_LINE_END_SIZE.MEDIUM
        assert c.line.begin_arrowhead_style == MSO_LINE_END_TYPE.TRIANGLE
        assert c.line.head_end.width == MSO_LINE_END_SIZE.LARGE
        assert c.line.head_end.length == MSO_LINE_END_SIZE.MEDIUM

    def it_round_trips_tail_end_arrowhead(self):
        prs, slide = _blank_slide()
        c = slide.shapes.add_connector(2, Inches(1), Inches(2), Inches(5), Inches(2))
        c.line.tail_end.type = MSO_LINE_END_TYPE.STEALTH
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        c2 = list(prs2.slides[0].shapes)[0]
        assert c2.line.tail_end.type == MSO_LINE_END_TYPE.STEALTH

    def it_does_not_change_the_shipped_arrowhead_api(self):
        # ---regression: the shipped properties must still work standalone---
        _, slide = _blank_slide()
        c = slide.shapes.add_connector(2, Inches(1), Inches(1), Inches(2), Inches(2))
        c.line.end_arrowhead_style = MSO_LINE_END_TYPE.OVAL
        assert c.line.tail_end.type == MSO_LINE_END_TYPE.OVAL


# ───────────────────────── F. Connector alias (SF6) ─────────────────────────


class DescribeStartEndConnectionAlias:
    def it_aliases_begin_and_end_connect(self):
        _, slide = _blank_slide()
        a = _rect(slide)
        b = slide.shapes.add_shape(1, Inches(5), Inches(1), Inches(2), Inches(1))
        c = slide.shapes.add_connector(2, Inches(1), Inches(1), Inches(4), Inches(4))
        c.start_connection(a, 1)
        c.end_connection(b, 3)
        cNvCxnSpPr = c._element.nvCxnSpPr.cNvCxnSpPr
        assert cNvCxnSpPr.find(f"{A}stCxn").get("id") == str(a.shape_id)
        assert cNvCxnSpPr.find(f"{A}endCxn").get("id") == str(b.shape_id)

    def it_round_trips_a_connected_connector(self):
        prs, slide = _blank_slide()
        a = _rect(slide)
        c = slide.shapes.add_connector(2, Inches(1), Inches(1), Inches(4), Inches(4))
        c.start_connection(a, 0)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        shapes2 = list(prs2.slides[0].shapes)
        conn = [s for s in shapes2 if s.shape_type is not None and "LINE" in str(s.shape_type)][0]
        assert conn._element.nvCxnSpPr.cNvCxnSpPr.find(f"{A}stCxn") is not None


# ───────────────────────── G. Group world coords (SF7) ─────────────────────────


class DescribeGroupWorldCoordinates:
    def it_returns_plain_coords_for_ungrouped_shape(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        assert sp.slide_left == sp.left
        assert sp.slide_top == sp.top
        assert sp.slide_width == sp.width

    def it_composes_one_group_transform(self):
        _, slide = _blank_slide()
        group = slide.shapes.add_group_shape()
        child = group.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
        gx = group._element.grpSpPr.get_or_add_xfrm()
        gx.get_or_add_off().x = Emu(Inches(2))
        gx.get_or_add_off().y = Emu(Inches(2))
        gx.get_or_add_ext().cx = Emu(Inches(4))
        gx.get_or_add_ext().cy = Emu(Inches(4))
        gx.get_or_add_chOff().x = Emu(0)
        gx.get_or_add_chOff().y = Emu(0)
        gx.get_or_add_chExt().cx = Emu(Inches(2))
        gx.get_or_add_chExt().cy = Emu(Inches(2))
        # child at (1",1") size 1" in a 2"→4" (×2) group offset to (2",2")
        # world_x = 2" + (1" - 0)*2 = 4" ; world_w = 1" * 2 = 2"
        assert child.slide_left == Emu(Inches(4))
        assert child.slide_top == Emu(Inches(4))
        assert child.slide_width == Emu(Inches(2))

    def it_composes_nested_groups(self):
        _, slide = _blank_slide()
        outer = slide.shapes.add_group_shape()
        inner = outer.shapes.add_group_shape()
        child = inner.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))

        def setxf(grp, ox, oy, ex, ey, cox, coy, cex, cey):
            x = grp._element.grpSpPr.get_or_add_xfrm()
            x.get_or_add_off().x = Emu(ox)
            x.get_or_add_off().y = Emu(oy)
            x.get_or_add_ext().cx = Emu(ex)
            x.get_or_add_ext().cy = Emu(ey)
            x.get_or_add_chOff().x = Emu(cox)
            x.get_or_add_chOff().y = Emu(coy)
            x.get_or_add_chExt().cx = Emu(cex)
            x.get_or_add_chExt().cy = Emu(cey)

        # inner: child-space 0..2", rendered 0..2" (×1), offset 0
        setxf(inner, 0, 0, Inches(2), Inches(2), 0, 0, Inches(2), Inches(2))
        # outer: child-space 0..2", rendered at (3",0") size 4" (×2)
        setxf(outer, Inches(3), 0, Inches(4), Inches(4), 0, 0, Inches(2), Inches(2))
        # child world_x = 3" + (1"*1 - 0)*2 = 5"
        assert child.slide_left == Emu(Inches(5))

    def it_falls_back_to_identity_on_degenerate_group(self):
        _, slide = _blank_slide()
        group = slide.shapes.add_group_shape()
        child = group.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
        gx = group._element.grpSpPr.get_or_add_xfrm()
        gx.get_or_add_chExt().cx = Emu(0)  # degenerate
        gx.get_or_add_chExt().cy = Emu(0)
        # must not raise ZeroDivisionError; identity → plain child coord
        assert child.slide_left == Emu(Inches(1))

    def it_does_not_mutate_stored_xfrm(self):
        _, slide = _blank_slide()
        group = slide.shapes.add_group_shape()
        child = group.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
        before = child._element.xpath("string(./p:spPr/a:xfrm/a:off/@x)")
        _ = child.slide_left
        after = child._element.xpath("string(./p:spPr/a:xfrm/a:off/@x)")
        assert before == after


# ───────────────────────── H. Flip (SF8) ─────────────────────────


class DescribeFlip:
    def it_sets_and_reads_flip_vertical(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        assert sp.flip_vertical is False
        sp.flip_vertical = True
        assert sp.flip_vertical is True
        assert sp._element.xpath("string(./p:spPr/a:xfrm/@flipV)") == "1"

    def it_sets_and_reads_flip_horizontal(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.flip_horizontal = True
        assert sp.flip_horizontal is True

    def it_round_trips_flip_vertical(self):
        prs, slide = _blank_slide()
        sp = _rect(slide)
        sp.flip_vertical = True
        _, sp2 = _roundtrip_first_autoshape(prs)
        assert sp2.flip_vertical is True

    def it_does_not_corrupt_sppr_order_creating_xfrm(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.flip_vertical = True
        kids = [el.tag for el in sp._element.spPr]
        assert kids[0] == f"{A}xfrm"


# ───────────────────────── I. Duplicate (SF9) ─────────────────────────


class DescribeShapeDuplicate:
    def it_returns_a_distinct_shape(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        dup = sp.duplicate()
        assert dup is not sp
        assert dup._element is not sp._element

    def it_assigns_a_unique_id_and_name(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        dup = sp.duplicate()
        assert dup.shape_id != sp.shape_id
        assert dup.name != sp.name

    def it_appends_at_top_z_by_default(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        dup = sp.duplicate()
        assert list(slide.shapes)[-1].shape_id == dup.shape_id

    def it_inserts_at_requested_z(self):
        _, slide = _blank_slide()
        a = _rect(slide)
        slide.shapes.add_shape(1, Inches(4), Inches(1), Inches(1), Inches(1))
        dup = a.duplicate(insert_at_z=0)
        assert list(slide.shapes)[0].shape_id == dup.shape_id

    def it_is_an_independent_deep_copy(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        dup = sp.duplicate()
        dup.left = Inches(6)
        assert sp.left != dup.left

    def it_preserves_an_effect_on_the_clone(self):
        _, slide = _blank_slide()
        sp = _rect(slide)
        sp.shadow.glow_effect.color.rgb = RGBColor(0xFF, 0, 0)
        sp.shadow.glow_effect.radius = Pt(10)
        dup = sp.duplicate()
        assert dup.shadow.glow_effect.radius == Emu(Pt(10))
        assert dup.shadow.glow_effect.color.rgb == RGBColor(0xFF, 0, 0)

    def it_gives_every_cNvPr_a_unique_id_when_duplicating_a_group(self):
        # ---regression: pre-append max+1 returned the same id on every
        # ---call, so a duplicated group's children all collided → repair---
        _, slide = _blank_slide()
        group = slide.shapes.add_group_shape()
        group.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
        group.shapes.add_shape(1, Inches(2), Inches(2), Inches(1), Inches(1))
        group.duplicate()
        all_ids = slide.shapes._spTree.xpath("//p:cNvPr/@id")
        assert len(all_ids) == len(set(all_ids)), f"duplicate id collision: {all_ids}"

    def it_round_trips_a_duplicated_shape(self):
        prs, slide = _blank_slide()
        sp = _rect(slide)
        sp.duplicate()
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        autoshapes = [
            s
            for s in prs2.slides[0].shapes
            if s.shape_type is not None and "AUTO_SHAPE" in str(s.shape_type)
        ]
        assert len(autoshapes) == 2
        ids = [s.shape_id for s in prs2.slides[0].shapes]
        assert len(ids) == len(set(ids))  # no id collision => no repair


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(pytest.main([__file__, "-q"]))
