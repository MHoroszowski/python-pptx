# pyright: reportPrivateUsage=false

"""Unit + round-trip tests for issue #16.

[Epic] Advanced Text, Auto-fit & Internationalization —
https://github.com/MHoroszowski/python-pptx/issues/16

Covers all 11 sub-features: sub/superscript, strike, highlight,
character-spacing+kerning, latin/east_asian/complex_script (incl.
Font.name backward-compat), columns+spacing, text_direction, paragraph
RTL, will_overflow/overflow_info, shrink_text_to_fit, and the #168
fit_text long-word crash fix.

Layered like tests/test_issue18_shape_effects.py: API-surface unit tests
+ save→reopen round-trip integration (the only layer that proves the
file is Office-compatible and nothing was silently dropped).
"""

from __future__ import annotations

import io
import os

import pytest

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_TEXT_DIRECTION, MSO_TEXT_STRIKE_TYPE
from pptx.util import Inches, Pt

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
TEST_FONT = os.path.join(os.path.dirname(__file__), "test_files", "calibriz.ttf")


def _frame():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    return prs, tb.text_frame


def _run(tf):
    r = tf.paragraphs[0].add_run()
    r.text = "Sample"
    return r


def _roundtrip(prs):
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


def _first_tf(prs2):
    return list(prs2.slides[0].shapes)[0].text_frame


# ───────── A. super/subscript (SF1) ─────────


class DescribeSuperSubscript:
    def it_sets_superscript(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.superscript = True
        assert f.superscript is True
        assert f._rPr.get("baseline") == "30000"

    def it_sets_subscript_negative_baseline(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.subscript = True
        assert f.subscript is True
        assert int(f._rPr.get("baseline")) < 0

    def it_makes_super_and_sub_mutually_exclusive(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.superscript = True
        f.subscript = True
        assert f.subscript is True
        assert f.superscript is False

    def it_clears_baseline_on_false(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.superscript = True
        f.superscript = False
        assert f.superscript is None
        assert f._rPr.get("baseline") is None

    def it_round_trips_superscript(self):
        prs, tf = _frame()
        _run(tf).font.superscript = True
        f2 = _first_tf(_roundtrip(prs)).paragraphs[0].runs[0].font
        assert f2.superscript is True


# ───────── B. strike (SF2) ─────────


class DescribeStrike:
    def it_sets_single_strike(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.strike = MSO_TEXT_STRIKE_TYPE.SINGLE
        assert f._rPr.get("strike") == "sngStrike"

    def it_sets_double_strike(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.strike = MSO_TEXT_STRIKE_TYPE.DOUBLE
        assert f._rPr.get("strike") == "dblStrike"
        assert f.strike == MSO_TEXT_STRIKE_TYPE.DOUBLE

    def it_removes_strike_on_none(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.strike = MSO_TEXT_STRIKE_TYPE.SINGLE
        f.strike = None
        assert f._rPr.get("strike") is None
        assert f.strike is None

    def it_round_trips_strike(self):
        prs, tf = _frame()
        _run(tf).font.strike = MSO_TEXT_STRIKE_TYPE.DOUBLE
        f2 = _first_tf(_roundtrip(prs)).paragraphs[0].runs[0].font
        assert f2.strike == MSO_TEXT_STRIKE_TYPE.DOUBLE


# ───────── C. highlight (SF3) ─────────


class DescribeHighlight:
    def it_sets_highlight_rgb(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.highlight.rgb = RGBColor(0xFF, 0xFF, 0x00)
        hl = f._rPr.find(f"{A}highlight")
        assert hl is not None
        assert hl.find(f"{A}srgbClr").get("val") == "FFFF00"

    def it_reads_none_without_mutating_when_absent(self):
        prs, tf = _frame()
        f = _run(tf).font
        before = len(list(f._rPr))
        assert f.highlight.rgb is None
        assert len(list(f._rPr)) == before  # no mutation on read

    def it_orders_highlight_before_typeface_trio(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.name = "Arial"
        f.highlight.rgb = RGBColor(0x00, 0xFF, 0x00)
        kids = [c.tag for c in f._rPr]
        assert kids.index(f"{A}highlight") < kids.index(f"{A}latin")

    def it_round_trips_highlight(self):
        prs, tf = _frame()
        _run(tf).font.highlight.rgb = RGBColor(0xFF, 0x00, 0xFF)
        f2 = _first_tf(_roundtrip(prs)).paragraphs[0].runs[0].font
        assert f2.highlight.rgb == RGBColor(0xFF, 0x00, 0xFF)


# ───────── D. character spacing + kerning (SF4) ─────────


class DescribeSpacingKerning:
    def it_sets_character_spacing(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.character_spacing = Pt(2)
        assert f._rPr.get("spc") == "200"
        assert f.character_spacing.pt == 2.0

    def it_supports_negative_character_spacing(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.character_spacing = Pt(-1)
        assert int(f._rPr.get("spc")) < 0

    def it_sets_kerning(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.kerning = Pt(12)
        assert f._rPr.get("kern") == "1200"

    def it_round_trips_spacing_and_kerning(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.character_spacing = Pt(3)
        f.kerning = Pt(8)
        f2 = _first_tf(_roundtrip(prs)).paragraphs[0].runs[0].font
        assert f2.character_spacing.pt == 3.0
        assert f2.kerning.pt == 8.0


# ───────── E. latin/east_asian/complex_script (SF5) ─────────


class DescribeTypefaceTrio:
    def it_sets_east_asian_without_touching_latin(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.east_asian = "MS Gothic"
        assert f._rPr.find(f"{A}ea").get("typeface") == "MS Gothic"
        assert f._rPr.find(f"{A}latin") is None  # issue acceptance

    def it_sets_complex_script(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.complex_script = "Arial"
        assert f._rPr.find(f"{A}cs").get("typeface") == "Arial"

    def it_keeps_font_name_latin_only_backward_compat(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.name = "Calibri"
        assert f._rPr.find(f"{A}latin").get("typeface") == "Calibri"
        assert f._rPr.find(f"{A}ea") is None
        assert f._rPr.find(f"{A}cs") is None
        assert f.name == "Calibri"
        assert f.latin == "Calibri"

    def it_orders_latin_ea_cs(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.complex_script = "C"
        f.east_asian = "E"
        f.name = "L"
        kids = [c.tag for c in f._rPr]
        assert kids.index(f"{A}latin") < kids.index(f"{A}ea") < kids.index(f"{A}cs")

    def it_round_trips_the_trio(self):
        prs, tf = _frame()
        f = _run(tf).font
        f.name, f.east_asian, f.complex_script = "Calibri", "MS Gothic", "Arial"
        f2 = _first_tf(_roundtrip(prs)).paragraphs[0].runs[0].font
        assert (f2.name, f2.east_asian, f2.complex_script) == (
            "Calibri",
            "MS Gothic",
            "Arial",
        )


# ───────── F. columns + spacing (SF6) ─────────


class DescribeColumns:
    def it_defaults_to_one_column(self):
        prs, tf = _frame()
        assert tf.columns == 1

    def it_sets_columns_and_spacing(self):
        prs, tf = _frame()
        tf.columns = 2
        tf.column_spacing = Pt(36)
        assert tf._bodyPr.get("numCol") == "2"
        assert tf.column_spacing.pt == 36.0

    def it_rejects_out_of_range_columns(self):
        prs, tf = _frame()
        with pytest.raises(ValueError):
            tf.columns = 0
        with pytest.raises(ValueError):
            tf.columns = 17

    def it_round_trips_columns(self):
        prs, tf = _frame()
        tf.columns = 3
        tf.column_spacing = Pt(24)
        tf2 = _first_tf(_roundtrip(prs))
        assert tf2.columns == 3
        assert tf2.column_spacing.pt == 24.0


# ───────── G. text_direction (SF7) ─────────


class DescribeTextDirection:
    def it_sets_vertical(self):
        prs, tf = _frame()
        tf.text_direction = MSO_TEXT_DIRECTION.VERTICAL
        assert tf._bodyPr.get("vert") == "vert"

    def it_sets_east_asian_vertical(self):
        prs, tf = _frame()
        tf.text_direction = MSO_TEXT_DIRECTION.EAST_ASIAN_VERTICAL
        assert tf._bodyPr.get("vert") == "eaVert"
        assert tf.text_direction == MSO_TEXT_DIRECTION.EAST_ASIAN_VERTICAL

    def it_removes_direction_on_none(self):
        prs, tf = _frame()
        tf.text_direction = MSO_TEXT_DIRECTION.VERTICAL
        tf.text_direction = None
        assert tf._bodyPr.get("vert") is None

    def it_round_trips_direction(self):
        prs, tf = _frame()
        tf.text_direction = MSO_TEXT_DIRECTION.VERTICAL_270
        assert _first_tf(_roundtrip(prs)).text_direction == MSO_TEXT_DIRECTION.VERTICAL_270


# ───────── H. paragraph RTL (SF8) ─────────


class DescribeParagraphRtl:
    def it_sets_rtl_true(self):
        prs, tf = _frame()
        p = tf.paragraphs[0]
        p.text = "مرحبا"
        p.rtl = True
        assert p._p.find(f"{A}pPr").get("rtl") == "1"
        assert p.rtl is True

    def it_sets_rtl_false(self):
        prs, tf = _frame()
        p = tf.paragraphs[0]
        p.rtl = False
        assert p._p.find(f"{A}pPr").get("rtl") == "0"

    def it_removes_rtl_on_none(self):
        prs, tf = _frame()
        p = tf.paragraphs[0]
        p.rtl = True
        p.rtl = None
        assert p._p.find(f"{A}pPr").get("rtl") is None

    def it_round_trips_arabic_rtl_paragraph(self):
        prs, tf = _frame()
        p = tf.paragraphs[0]
        p.text = "اللغة العربية"
        p.rtl = True
        p2 = _first_tf(_roundtrip(prs)).paragraphs[0]
        assert p2.rtl is True
        assert "العربية" in p2.text


# ───────── I. overflow detection (SF9) ─────────


class DescribeOverflowDetection:
    def it_reports_no_overflow_for_empty_frame(self):
        prs, tf = _frame()
        assert tf.will_overflow(font_file=TEST_FONT) is False

    def it_reports_no_overflow_when_text_fits(self):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5)).text_frame
        tf.text = "short"
        assert tf.will_overflow(font_file=TEST_FONT) is False

    def it_reports_overflow_for_oversized_content(self):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(0.3)).text_frame
        tf.text = "Supercalifragilistic " * 12
        info = tf.overflow_info(font_file=TEST_FONT)
        assert info.overflows is True
        assert info.required_height > info.available_height

    @pytest.mark.parametrize("fr_h", [0.3, 0.4, 0.5])
    def it_detects_overflow_at_110pct_across_sizes(self, fr_h):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(fr_h)).text_frame
        tf.text = "The quick brown fox jumps over the lazy dog. " * 6
        assert tf.will_overflow(font_file=TEST_FONT) is True

    def it_does_not_mutate_on_overflow_check(self):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(0.3)).text_frame
        tf.text = "overflowing " * 20
        before = tf._txBody.xml
        before_autosize = tf.auto_size
        tf.will_overflow(font_file=TEST_FONT)
        assert tf._txBody.xml == before  # no XML mutation at all
        assert tf.auto_size == before_autosize  # autofit unchanged


# ───────── J. shrink_text_to_fit (SF10) ─────────


class DescribeShrinkTextToFit:
    def it_sets_normautofit_with_reduced_scale(self):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(0.3)).text_frame
        tf.text = "Supercalifragilistic " * 14
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(18)
        tf.shrink_text_to_fit(font_file=TEST_FONT)
        from pptx.enum.text import MSO_AUTO_SIZE

        assert tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        na = tf._txBody.bodyPr.normAutofit
        assert na is not None
        assert na.fontScale < 100

    def it_does_not_rewrite_run_sizes(self):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(0.3)).text_frame
        tf.text = "overflow " * 20
        r = tf.paragraphs[0].runs[0]
        r.font.size = Pt(18)
        tf.shrink_text_to_fit(font_file=TEST_FONT)
        assert r.font.size == Pt(18)  # eager scale, not size rewrite

    def it_is_a_noop_on_empty_frame(self):
        prs, tf = _frame()
        tf.shrink_text_to_fit(font_file=TEST_FONT)  # must not raise
        assert tf._txBody.bodyPr.normAutofit is None


# ───────── K. fit_text #168 crash fix (SF11) ─────────


class DescribeFitTextLongWordFix:
    def it_does_not_crash_on_a_single_long_word(self):
        # scanny/python-pptx#168 — previously raised TypeError
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(0), Inches(0), Inches(0.5), Inches(0.5)).text_frame
        tf.text = "Supercalifragilisticexpialidocious"
        tf.fit_text(font_file=TEST_FONT)  # no TypeError
        assert tf.auto_size is not None

    def it_still_fits_normal_multiword_text(self):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tf = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2)).text_frame
        tf.text = "one two three four five"
        tf.fit_text(font_file=TEST_FONT)
        sz = tf.paragraphs[0].runs[0].font.size
        assert sz is not None
        assert sz.pt >= 1

    def it_break_line_returns_usable_line_when_nothing_fits(self):
        from pptx.text.layout import TextFitter, _LineSource

        ls = _LineSource("Supercalifragilisticexpialidocious")
        fitter = TextFitter(ls, (Inches(0.2), Inches(0.2)), TEST_FONT)
        result = fitter._break_line(ls, 18)
        assert result is not None
        assert result.text != ""


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(pytest.main([__file__, "-q"]))
