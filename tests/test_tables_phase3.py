# pyright: reportPrivateUsage=false

"""Unit-test suite for Tables 2.0 Phase 3 — merge robustness API.

Covers:

- Read-only inspection accessors on |_Cell|: ``gridSpan``, ``rowSpan``,
  ``hMerge``, ``vMerge`` — mirror the underlying ``a:tc`` attrs and let
  callers inspect any cell's merge state without `is_merge_origin`
  heuristics.
- Range-style merge: |Table|.merge_cells(row_range, col_range) — idempotent
  on already-merged-exactly-this-way regions, no-op on single-cell ranges,
  raises |ValueError| on partial overlap with a different-shape merge.
- Range-style split: |Table|.split_cells(row_range, col_range) —
  idempotent on un-merged ranges, raises |ValueError| when a merge
  extends beyond the requested range boundary.
- Range-arg shape: tuples (inclusive) and Python ``range`` objects
  (half-open) are both accepted; order within either form is irrelevant.
- Round-trip: a merge applied via merge_cells survives save/reload.
- Anti: existing |Cell|.merge / |Cell|.split unchanged; Phase-2 style API
  unaffected by Phase-3 additions.

Issue: https://github.com/MHoroszowski/python-pptx/issues/12 (Phase 3).
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.table import _Cell, _normalize_range
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


def _make_table(rows: int, cols: int):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    gf = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(2))
    return prs, gf.table


@pytest.fixture
def t3x3():
    _, t = _make_table(3, 3)
    return t


@pytest.fixture
def t4x4():
    _, t = _make_table(4, 4)
    return t


# ---------------------------------------------------------------------------
# Read-only inspection accessors — _Cell.grid_span / rowSpan / hMerge / vMerge
# ---------------------------------------------------------------------------


class DescribeCell_InspectionAccessors(object):
    """Unit-test suite for `_Cell.grid_span`, `rowSpan`, `hMerge`, `vMerge`."""

    def it_returns_default_gridSpan_1_for_unmerged_cell(self, t3x3):
        assert t3x3.cell(0, 0).grid_span == 1

    def it_returns_default_rowSpan_1_for_unmerged_cell(self, t3x3):
        assert t3x3.cell(0, 0).row_span == 1

    def it_returns_False_hMerge_for_unmerged_cell(self, t3x3):
        assert t3x3.cell(0, 0).h_merge is False

    def it_returns_False_vMerge_for_unmerged_cell(self, t3x3):
        assert t3x3.cell(0, 0).v_merge is False

    def it_reports_gridSpan_on_merge_origin_after_merge(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))  # 2 rows x 3 cols
        assert t3x3.cell(0, 0).grid_span == 3

    def it_reports_rowSpan_on_merge_origin_after_merge(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))  # 2 rows x 3 cols
        assert t3x3.cell(0, 0).row_span == 2

    def it_reports_hMerge_True_on_horizontally_spanned_cell(self, t3x3):
        t3x3.merge_cells((0, 0), (0, 2))  # 1 row x 3 cols
        assert t3x3.cell(0, 1).h_merge is True
        assert t3x3.cell(0, 2).h_merge is True

    def it_reports_vMerge_True_on_vertically_spanned_cell(self, t3x3):
        t3x3.merge_cells((0, 2), (0, 0))  # 3 rows x 1 col
        assert t3x3.cell(1, 0).v_merge is True
        assert t3x3.cell(2, 0).v_merge is True

    def it_reports_both_hMerge_and_vMerge_on_inner_spanned_cell_of_block(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))  # 2x3 block
        # ---inner spanned (1,2) is both horizontally and vertically spanned---
        c = t3x3.cell(1, 2)
        assert c.h_merge is True
        assert c.v_merge is True

    def it_resets_all_attrs_after_split(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))
        t3x3.split_cells((0, 1), (0, 2))
        for r in range(2):
            for c in range(3):
                cell = t3x3.cell(r, c)
                assert cell.grid_span == 1
                assert cell.row_span == 1
                assert cell.h_merge is False
                assert cell.v_merge is False

    @pytest.mark.parametrize("attr", ["grid_span", "row_span", "h_merge", "v_merge"])
    def it_is_read_only_no_setter(self, t3x3, attr):
        cell = t3x3.cell(0, 0)
        with pytest.raises(AttributeError):
            setattr(cell, attr, 99)


# ---------------------------------------------------------------------------
# Table.merge_cells
# ---------------------------------------------------------------------------


class DescribeTable_merge_cells(object):
    """Unit-test suite for `Table.merge_cells`."""

    def it_merges_a_2x3_block(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))

        origin = t3x3.cell(0, 0)
        assert origin.is_merge_origin is True
        assert origin.grid_span == 3
        assert origin.row_span == 2

    def it_returns_the_merge_origin_Cell(self, t3x3):
        result = t3x3.merge_cells((0, 1), (0, 2))

        assert isinstance(result, _Cell)
        assert result._tc is t3x3.cell(0, 0)._tc

    def it_is_idempotent_on_exact_re_merge(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))
        # ---second call must not raise and not change shape---
        result = t3x3.merge_cells((0, 1), (0, 2))

        assert t3x3.cell(0, 0).grid_span == 3
        assert t3x3.cell(0, 0).row_span == 2
        assert result._tc is t3x3.cell(0, 0)._tc

    def it_accepts_a_python_range_object(self, t3x3):
        # ---half-open per Python convention; range(0, 2) covers rows 0 and 1---
        t3x3.merge_cells(range(0, 2), range(0, 3))

        assert t3x3.cell(0, 0).row_span == 2
        assert t3x3.cell(0, 0).grid_span == 3

    def it_treats_a_single_cell_range_as_a_noop(self, t3x3):
        # ---single-cell merge = no merge needed; returns the cell itself---
        result = t3x3.merge_cells((0, 0), (0, 0))

        assert result._tc is t3x3.cell(0, 0)._tc
        assert t3x3.cell(0, 0).grid_span == 1
        assert t3x3.cell(0, 0).row_span == 1

    def it_is_order_agnostic_within_a_tuple(self, t3x3):
        # ---(2,0) means same as (0,2): rows 0..2 inclusive---
        t3x3.merge_cells((2, 0), (2, 0))

        assert t3x3.cell(0, 0).row_span == 3
        assert t3x3.cell(0, 0).grid_span == 3

    def it_raises_ValueError_on_partial_overlap_with_different_shape(self, t3x3):
        # ---first merge: 1 row x 2 cols at top-left---
        t3x3.merge_cells((0, 0), (0, 1))

        # ---second merge: 2 rows x 2 cols at top-left = different shape---
        with pytest.raises(ValueError) as excinfo:
            t3x3.merge_cells((0, 1), (0, 1))
        assert "partially overlaps" in str(excinfo.value)

    def it_writes_correct_hMerge_vMerge_for_block_merge(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))  # 2x3

        # ---origin: gridSpan=3, rowSpan=2, no hMerge/vMerge
        assert t3x3.cell(0, 0).h_merge is False
        assert t3x3.cell(0, 0).v_merge is False
        # ---top row (excl origin): hMerge=True, vMerge=False
        assert t3x3.cell(0, 1).h_merge is True
        assert t3x3.cell(0, 1).v_merge is False
        assert t3x3.cell(0, 2).h_merge is True
        # ---left col (excl origin): vMerge=True
        assert t3x3.cell(1, 0).v_merge is True
        assert t3x3.cell(1, 0).h_merge is False
        # ---inner: both
        assert t3x3.cell(1, 1).h_merge is True
        assert t3x3.cell(1, 1).v_merge is True

    def it_raises_on_empty_range(self, t3x3):
        with pytest.raises(ValueError):
            t3x3.merge_cells(range(0, 0), range(0, 1))

    def it_raises_on_unsupported_range_argument_type(self, t3x3):
        with pytest.raises(TypeError):
            t3x3.merge_cells([0, 1], [0, 1])  # ---list, not tuple or range


# ---------------------------------------------------------------------------
# Table.split_cells
# ---------------------------------------------------------------------------


class DescribeTable_split_cells(object):
    """Unit-test suite for `Table.split_cells`."""

    def it_splits_a_merged_block(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))

        t3x3.split_cells((0, 1), (0, 2))

        assert t3x3.cell(0, 0).grid_span == 1
        assert t3x3.cell(0, 0).row_span == 1
        assert t3x3.cell(0, 0).is_merge_origin is False
        for r in range(2):
            for c in range(3):
                assert t3x3.cell(r, c).h_merge is False
                assert t3x3.cell(r, c).v_merge is False

    def it_is_idempotent_on_unmerged_range(self, t3x3):
        # ---no merges in the range; split_cells should be a no-op---
        t3x3.split_cells((0, 2), (0, 2))

        # ---all cells still default---
        for r in range(3):
            for c in range(3):
                assert t3x3.cell(r, c).grid_span == 1

    def it_accepts_a_python_range_object(self, t3x3):
        t3x3.merge_cells(range(0, 2), range(0, 3))

        t3x3.split_cells(range(0, 2), range(0, 3))

        for r in range(2):
            for c in range(3):
                assert t3x3.cell(r, c).grid_span == 1
                assert t3x3.cell(r, c).row_span == 1

    def it_raises_when_merge_extends_beyond_split_range(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))  # 2x3 over top-left

        # ---try to split only (0,0)..(0,1) — merge extends to col 2---
        with pytest.raises(ValueError) as excinfo:
            t3x3.split_cells((0, 0), (0, 1))
        assert "outside split range" in str(excinfo.value)

    def it_raises_when_spanned_cell_origin_is_outside_range(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))  # 2x3 over top-left

        # ---try to split just (1,1)..(1,2) — origin (0,0) is outside---
        with pytest.raises(ValueError) as excinfo:
            t3x3.split_cells((1, 1), (1, 2))
        assert "starts outside split range" in str(excinfo.value)

    def it_can_split_multiple_merges_in_one_call(self, t4x4):
        t4x4.merge_cells((0, 0), (0, 1))  # 1x2 at top-left
        t4x4.merge_cells((2, 2), (2, 3))  # 1x2 at bottom-right

        t4x4.split_cells((0, 3), (0, 3))  # entire table

        for r in range(4):
            for c in range(4):
                assert t4x4.cell(r, c).grid_span == 1
                assert t4x4.cell(r, c).h_merge is False


# ---------------------------------------------------------------------------
# _normalize_range helper
# ---------------------------------------------------------------------------


class Describe_normalize_range(object):
    """Unit-test suite for the private `_normalize_range` helper."""

    @pytest.mark.parametrize(
        ("rng", "expected"),
        [
            ((0, 1), (0, 1)),
            ((1, 0), (0, 1)),  # ---order-agnostic
            ((5, 5), (5, 5)),  # ---single-cell
            (range(0, 2), (0, 1)),  # ---half-open range -> inclusive low/high
            (range(2, 5), (2, 4)),
        ],
    )
    def it_normalizes_to_inclusive_low_high(self, rng, expected):
        assert _normalize_range(rng) == expected

    @pytest.mark.parametrize(
        "rng",
        [
            [0, 1],  # ---list, not tuple
            (0, 1, 2),  # ---3-tuple
            "01",  # ---string
        ],
    )
    def it_raises_TypeError_on_unsupported_shape(self, rng):
        with pytest.raises(TypeError):
            _normalize_range(rng)

    def it_raises_ValueError_on_negative_index(self):
        with pytest.raises(ValueError):
            _normalize_range((-1, 0))

    def it_raises_ValueError_on_empty_range(self):
        with pytest.raises(ValueError):
            _normalize_range(range(2, 2))

    def it_raises_ValueError_on_non_unit_step(self):
        with pytest.raises(ValueError):
            _normalize_range(range(0, 4, 2))


# ---------------------------------------------------------------------------
# Round-trip
# ---------------------------------------------------------------------------


class DescribeMerge_RoundTrip(object):
    """Save a presentation with a merge applied via merge_cells, reload — preserved."""

    def it_preserves_a_block_merge_through_save_and_reload(self):
        prs, t = _make_table(3, 3)
        t.merge_cells((0, 1), (0, 2))

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        prs2 = Presentation(buf)
        t2 = next(shp for shp in prs2.slides[0].shapes if shp.has_table).table
        assert t2.cell(0, 0).row_span == 2
        assert t2.cell(0, 0).grid_span == 3
        assert t2.cell(1, 1).h_merge is True
        assert t2.cell(1, 1).v_merge is True


# ---------------------------------------------------------------------------
# Anti / Regression
# ---------------------------------------------------------------------------


class DescribePhase3_Regression(object):
    """Anti-criteria: existing surfaces unchanged."""

    def it_keeps_existing_Cell_merge_working(self, t3x3):
        # ---legacy 2-cell merge API stays as-is---
        t3x3.cell(0, 0).merge(t3x3.cell(0, 1))

        assert t3x3.cell(0, 0).is_merge_origin is True
        assert t3x3.cell(0, 0).grid_span == 2

    def it_keeps_existing_Cell_split_working(self, t3x3):
        t3x3.cell(0, 0).merge(t3x3.cell(0, 1))
        t3x3.cell(0, 0).split()

        assert t3x3.cell(0, 0).grid_span == 1
        assert t3x3.cell(0, 1).h_merge is False

    def it_keeps_phase2_style_api_working(self, t3x3):
        # ---no regression from Phase 2 style surface---
        assert t3x3.style_id == "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"
        t3x3.apply_style("No Style, No Grid")
        assert t3x3.style_name == "No Style, No Grid"

    def it_keeps_existing_is_merge_origin_and_is_spanned_working(self, t3x3):
        t3x3.merge_cells((0, 1), (0, 2))

        assert t3x3.cell(0, 0).is_merge_origin is True
        assert t3x3.cell(0, 1).is_spanned is True
        assert t3x3.cell(0, 0).is_spanned is False
