# pyright: reportPrivateUsage=false

"""Unit-test suite for slide-CRUD additions: remove/move/indexed add_slide.

Issue: https://github.com/MHoroszowski/python-pptx/issues/11 (Phase 1).

Mirrors the unit-test style of `tests/test_slide.py`. Round-trip integration
tests at the bottom exercise the full pptx open → mutate → save → reopen path
on a synthetic in-memory presentation.
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.oxml.presentation import CT_SlideIdList
from pptx.parts.presentation import PresentationPart
from pptx.parts.slide import SlidePart
from pptx.slide import Slide, SlideLayout, Slides

from .unitutil.cxml import element, xml
from .unitutil.mock import instance_mock, method_mock, property_mock

# ---------------------------------------------------------------------------
# OXML LAYER — `p:sldIdLst` reordering helpers.
# ---------------------------------------------------------------------------


class DescribeCT_SlideIdList(object):
    """Unit-test suite for `CT_SlideIdList` insert/move/remove helpers."""

    def it_can_insert_into_an_empty_sldIdLst(self):
        """Cover the idx==len==0 path; assert by structure, not raw XML.

        The XML-prefix shape differs when the parent does not yet declare
        the `r:` namespace — semantic equality is what matters here.
        """
        sldIdLst = element("p:sldIdLst")
        assert isinstance(sldIdLst, CT_SlideIdList)

        sldIdLst.insert_sldId_at("rId9", 0)

        assert len(sldIdLst.sldId_lst) == 1
        assert sldIdLst.sldId_lst[0].rId == "rId9"
        assert sldIdLst.sldId_lst[0].id == 256

    @pytest.mark.parametrize(
        ("sldIdLst_cxml", "rId", "idx", "expected_cxml"),
        [
            # ---insert at head of populated list---
            (
                "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})",
                "rId9",
                0,
                (
                    "p:sldIdLst/(p:sldId{r:id=rId9,id=258},p:sldId{r:id=a,id=256},"
                    "p:sldId{r:id=b,id=257})"
                ),
            ),
            # ---insert in the middle---
            (
                "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})",
                "rId9",
                1,
                (
                    "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=rId9,id=258},"
                    "p:sldId{r:id=b,id=257})"
                ),
            ),
            # ---insert at the tail (idx == len)---
            (
                "p:sldIdLst/p:sldId{r:id=a,id=256}",
                "rId9",
                1,
                "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=rId9,id=257})",
            ),
        ],
    )
    def it_can_insert_a_sldId_at_a_specific_index(
        self, sldIdLst_cxml: str, rId: str, idx: int, expected_cxml: str
    ):
        sldIdLst = element(sldIdLst_cxml)
        assert isinstance(sldIdLst, CT_SlideIdList)

        sldIdLst.insert_sldId_at(rId, idx)

        assert sldIdLst.xml == xml(expected_cxml)

    @pytest.mark.parametrize("bad_idx", [-1, 5])
    def but_it_raises_IndexError_on_insert_out_of_range(self, bad_idx: int):
        sldIdLst = element("p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})")
        with pytest.raises(IndexError):
            sldIdLst.insert_sldId_at("rId9", bad_idx)

    @pytest.mark.parametrize(
        ("sldIdLst_cxml", "moving_rId", "new_idx", "expected_cxml"),
        [
            # ---move forward (head → tail)---
            (
                (
                    "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257},"
                    "p:sldId{r:id=c,id=258})"
                ),
                "a",
                2,
                (
                    "p:sldIdLst/(p:sldId{r:id=b,id=257},p:sldId{r:id=c,id=258},"
                    "p:sldId{r:id=a,id=256})"
                ),
            ),
            # ---move backward (tail → head)---
            (
                (
                    "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257},"
                    "p:sldId{r:id=c,id=258})"
                ),
                "c",
                0,
                (
                    "p:sldIdLst/(p:sldId{r:id=c,id=258},p:sldId{r:id=a,id=256},"
                    "p:sldId{r:id=b,id=257})"
                ),
            ),
            # ---move to current position is a no-op---
            (
                "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})",
                "a",
                0,
                "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})",
            ),
        ],
    )
    def it_can_move_a_sldId_to_a_new_index(
        self,
        sldIdLst_cxml: str,
        moving_rId: str,
        new_idx: int,
        expected_cxml: str,
    ):
        sldIdLst = element(sldIdLst_cxml)
        assert isinstance(sldIdLst, CT_SlideIdList)
        target = next(s for s in sldIdLst.sldId_lst if s.rId == moving_rId)

        sldIdLst.move_sldId_to(target, new_idx)

        assert sldIdLst.xml == xml(expected_cxml)

    @pytest.mark.parametrize("bad_idx", [-1, 3])
    def but_it_raises_IndexError_on_move_out_of_range(self, bad_idx: int):
        sldIdLst = element("p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})")
        target = sldIdLst.sldId_lst[0]
        with pytest.raises(IndexError):
            sldIdLst.move_sldId_to(target, bad_idx)

    def it_can_remove_a_sldId(self):
        sldIdLst = element("p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})")
        assert isinstance(sldIdLst, CT_SlideIdList)
        target = sldIdLst.sldId_lst[0]

        sldIdLst.remove_sldId(target)

        assert sldIdLst.xml == xml("p:sldIdLst/p:sldId{r:id=b,id=257}")

    def but_it_raises_on_remove_of_foreign_element(self):
        sldIdLst = element("p:sldIdLst/p:sldId{r:id=a,id=256}")
        foreign = element("p:sldIdLst/p:sldId{r:id=z,id=999}").sldId_lst[0]

        with pytest.raises(ValueError):
            sldIdLst.remove_sldId(foreign)


# ---------------------------------------------------------------------------
# `Slides` — indexed add_slide, move(), remove().
# ---------------------------------------------------------------------------


class DescribeSlides_CRUD(object):
    """Unit-test suite for the new CRUD methods on `pptx.slide.Slides`."""

    # -- add_slide(layout, index=None) ---------------------------------------

    def it_keeps_appending_when_no_index_is_given(self, append_fixture):
        slides, slide_layout_, expected_xml, slide_ = append_fixture

        slide = slides.add_slide(slide_layout_)

        assert slides._sldIdLst.xml == expected_xml
        assert slide is slide_

    def it_can_insert_a_new_slide_at_the_head(self, insert_head_fixture):
        slides, slide_layout_, expected_xml, slide_ = insert_head_fixture

        slide = slides.add_slide(slide_layout_, index=0)

        assert slides._sldIdLst.xml == expected_xml
        assert slide is slide_

    def it_can_insert_a_new_slide_at_a_middle_index(self, insert_middle_fixture):
        slides, slide_layout_, expected_xml, slide_ = insert_middle_fixture

        slide = slides.add_slide(slide_layout_, index=1)

        assert slides._sldIdLst.xml == expected_xml
        assert slide is slide_

    def it_can_insert_a_new_slide_at_the_tail(self, insert_tail_fixture):
        slides, slide_layout_, expected_xml, slide_ = insert_tail_fixture

        slide = slides.add_slide(slide_layout_, index=2)

        assert slides._sldIdLst.xml == expected_xml
        assert slide is slide_

    @pytest.mark.parametrize("bad_index", [-1, 99])
    def but_it_raises_on_add_slide_index_out_of_range(self, bad_index, part_prop_, slide_layout_):
        slides = Slides(
            element("p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})"),
            None,
        )
        with pytest.raises(IndexError):
            slides.add_slide(slide_layout_, index=bad_index)

    # -- move() ---------------------------------------------------------------

    def it_can_move_a_slide_to_a_new_position(self, move_fixture):
        slides, target_slide_, new_index, expected_xml = move_fixture

        slides.move(target_slide_, new_index)

        assert slides._sldIdLst.xml == expected_xml

    def but_it_raises_on_move_of_slide_not_in_collection(self, missing_slide_fixture):
        slides, slide_ = missing_slide_fixture
        with pytest.raises(ValueError):
            slides.move(slide_, 0)

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_move_index_out_of_range(self, bad_idx, part_prop_, slide_):
        slides = Slides(
            element("p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})"),
            None,
        )
        with pytest.raises(IndexError):
            slides.move(slide_, bad_idx)

    # -- remove() / Slide.delete() -------------------------------------------

    def it_can_remove_a_slide(self, remove_fixture):
        slides, target_slide_, expected_xml, prs_part_ = remove_fixture

        slides.remove(target_slide_)

        assert slides._sldIdLst.xml == expected_xml
        prs_part_.drop_rel.assert_called_once_with("a")

    def but_it_raises_on_remove_of_slide_not_in_collection(self, missing_slide_fixture):
        slides, slide_ = missing_slide_fixture
        with pytest.raises(ValueError):
            slides.remove(slide_)

    # ---------- fixtures -----------------------------------------------------

    @pytest.fixture
    def append_fixture(self, slide_layout_, part_prop_, slide_):
        slides = Slides(element("p:sldIdLst/p:sldId{r:id=rId1,id=256}"), None)
        part_ = part_prop_.return_value
        part_.add_slide.return_value = "rId2", slide_
        expected_xml = xml("p:sldIdLst/(p:sldId{r:id=rId1,id=256},p:sldId{r:id=rId2,id=257})")
        return slides, slide_layout_, expected_xml, slide_

    @pytest.fixture
    def insert_head_fixture(self, slide_layout_, part_prop_, slide_):
        slides = Slides(
            element("p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})"),
            None,
        )
        part_ = part_prop_.return_value
        part_.add_slide.return_value = "rIdNew", slide_
        expected_xml = xml(
            "p:sldIdLst/(p:sldId{r:id=rIdNew,id=258},p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})"
        )
        return slides, slide_layout_, expected_xml, slide_

    @pytest.fixture
    def insert_middle_fixture(self, slide_layout_, part_prop_, slide_):
        slides = Slides(
            element("p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})"),
            None,
        )
        part_ = part_prop_.return_value
        part_.add_slide.return_value = "rIdNew", slide_
        expected_xml = xml(
            "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=rIdNew,id=258},p:sldId{r:id=b,id=257})"
        )
        return slides, slide_layout_, expected_xml, slide_

    @pytest.fixture
    def insert_tail_fixture(self, slide_layout_, part_prop_, slide_):
        slides = Slides(
            element("p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})"),
            None,
        )
        part_ = part_prop_.return_value
        part_.add_slide.return_value = "rIdNew", slide_
        expected_xml = xml(
            "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257},p:sldId{r:id=rIdNew,id=258})"
        )
        return slides, slide_layout_, expected_xml, slide_

    @pytest.fixture
    def move_fixture(self, request, part_prop_, _index_):
        # ---move slide at index 0 to index 2---
        sldIdLst = element(
            "p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257},p:sldId{r:id=c,id=258})"
        )
        slides = Slides(sldIdLst, None)
        target_slide_ = instance_mock(request, Slide)
        _index_.return_value = 0
        expected_xml = xml(
            "p:sldIdLst/(p:sldId{r:id=b,id=257},p:sldId{r:id=c,id=258},p:sldId{r:id=a,id=256})"
        )
        return slides, target_slide_, 2, expected_xml

    @pytest.fixture
    def remove_fixture(self, request, part_prop_, _index_):
        sldIdLst = element("p:sldIdLst/(p:sldId{r:id=a,id=256},p:sldId{r:id=b,id=257})")
        slides = Slides(sldIdLst, None)
        target_slide_ = instance_mock(request, Slide)
        _index_.return_value = 0
        prs_part_ = part_prop_.return_value
        expected_xml = xml("p:sldIdLst/p:sldId{r:id=b,id=257}")
        return slides, target_slide_, expected_xml, prs_part_

    @pytest.fixture
    def missing_slide_fixture(self, request, part_prop_, _index_):
        sldIdLst = element("p:sldIdLst/p:sldId{r:id=a,id=256}")
        slides = Slides(sldIdLst, None)
        slide_ = instance_mock(request, Slide)
        _index_.side_effect = ValueError("not in collection")
        return slides, slide_

    # ---------- fixture components ------------------------------------------

    @pytest.fixture
    def part_prop_(self, request, prs_part_):
        return property_mock(request, Slides, "part", return_value=prs_part_)

    @pytest.fixture
    def prs_part_(self, request):
        return instance_mock(request, PresentationPart)

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, Slide)

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def _index_(self, request):
        return method_mock(request, Slides, "index")


# ---------------------------------------------------------------------------
# `Slide.delete()` alias.
# ---------------------------------------------------------------------------


class DescribeSlide_Delete(object):
    """Unit-test suite for `pptx.slide.Slide.delete`."""

    def it_delegates_to_Slides_remove_on_the_owning_presentation(self, request):
        slide_part_ = instance_mock(request, SlidePart)
        prs_ = instance_mock(request, type("Prs", (), {"slides": None}))
        prs_.slides = instance_mock(request, Slides)
        prs_part_ = instance_mock(request, PresentationPart)
        prs_part_.presentation = prs_
        slide_part_.package.presentation_part = prs_part_

        slide = Slide(element("p:sld/p:cSld"), slide_part_)

        slide.delete()

        prs_.slides.remove.assert_called_once_with(slide)


# ---------------------------------------------------------------------------
# Round-trip integration tests — open → mutate → save → reopen.
# ---------------------------------------------------------------------------


def _seed_presentation_with(n_slides: int) -> Presentation:
    """Build an in-memory Presentation seeded with `n_slides` blank slides."""
    prs = Presentation()
    layout = prs.slide_layouts[6]  # ---blank layout---
    for _ in range(n_slides):
        prs.slides.add_slide(layout)
    return prs


def _round_trip(prs: Presentation) -> Presentation:
    """Save `prs` to a bytes buffer and reopen it."""
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


class DescribeSlideCrudRoundTrip(object):
    """Open → mutate → save → reopen integration coverage."""

    def it_round_trips_an_indexed_add(self):
        prs = _seed_presentation_with(2)
        layout = prs.slide_layouts[6]
        original_ids = [s.slide_id for s in prs.slides]

        new_slide = prs.slides.add_slide(layout, index=1)
        new_id = new_slide.slide_id

        round_tripped = _round_trip(prs)
        ids_after = [s.slide_id for s in round_tripped.slides]

        assert len(round_tripped.slides) == 3
        assert ids_after == [original_ids[0], new_id, original_ids[1]]

    def it_round_trips_a_remove(self):
        prs = _seed_presentation_with(3)
        ids_before = [s.slide_id for s in prs.slides]

        prs.slides.remove(prs.slides[1])

        round_tripped = _round_trip(prs)
        ids_after = [s.slide_id for s in round_tripped.slides]

        assert len(round_tripped.slides) == 2
        assert ids_after == [ids_before[0], ids_before[2]]

    def it_round_trips_a_Slide_delete_call(self):
        prs = _seed_presentation_with(3)
        ids_before = [s.slide_id for s in prs.slides]

        prs.slides[0].delete()

        round_tripped = _round_trip(prs)
        ids_after = [s.slide_id for s in round_tripped.slides]

        assert len(round_tripped.slides) == 2
        assert ids_after == [ids_before[1], ids_before[2]]

    def it_round_trips_a_move(self):
        prs = _seed_presentation_with(3)
        ids_before = [s.slide_id for s in prs.slides]

        prs.slides.move(prs.slides[0], 2)

        round_tripped = _round_trip(prs)
        ids_after = [s.slide_id for s in round_tripped.slides]

        assert len(round_tripped.slides) == 3
        assert ids_after == [ids_before[1], ids_before[2], ids_before[0]]

    def it_treats_index_equal_to_len_as_append(self):
        """`add_slide(layout, index=len(slides))` should be a valid append."""
        prs = _seed_presentation_with(2)
        layout = prs.slide_layouts[6]

        new_slide = prs.slides.add_slide(layout, index=len(prs.slides))

        assert len(prs.slides) == 3
        assert prs.slides[2].slide_id == new_slide.slide_id

    def it_preserves_a_shared_image_after_removing_one_referencing_slide(self):
        """Shared image-part survives remove() of one of two slides using it.

        Image dedup: python-pptx's package serializer keeps any part that
        remains reachable; if slide B references the same image as slide A
        and we remove A, the image must still be readable on B.
        """
        from pathlib import Path

        from pptx.util import Inches

        png = Path(__file__).parent / "test_files" / "python-powered.png"
        prs = Presentation()
        layout = prs.slide_layouts[6]
        slide_a = prs.slides.add_slide(layout)
        slide_a.shapes.add_picture(str(png), Inches(1), Inches(1))
        slide_b = prs.slides.add_slide(layout)
        slide_b.shapes.add_picture(str(png), Inches(1), Inches(1))

        prs.slides.remove(slide_a)
        round_tripped = _round_trip(prs)

        assert len(round_tripped.slides) == 1
        # ---surviving slide must still expose its picture (== reachable image part)---
        pictures = [shp for shp in round_tripped.slides[0].shapes if shp.shape_type == 13]
        assert len(pictures) == 1
        assert pictures[0].image.blob is not None
        assert len(pictures[0].image.blob) > 0
