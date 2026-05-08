# pyright: reportPrivateUsage=false

"""Unit-test suite for Tables 2.0 Phase 1 — row/column CRUD.

Covers:

- New oxml helpers `CT_TableGrid.{insert_gridCol_at, remove_gridCol_at}`,
  `CT_Table.insert_tr_at` / `remove_tr_at`, `CT_TableRow.{insert_tc_at,
  remove_tc_at, has_cross_row_merge}`, `CT_Table.column_has_cross_column_merge`.
- Public API `Table.rows.add(at, height)`, `Table.rows.remove(index)`,
  `Table.columns.add(at, width)`, `Table.columns.remove(index)`.
- Round-trip integration: open → mutate → save → reopen.
- Anti-criteria: removing a row/column that participates in a multi-row /
  multi-column merge raises ValueError; inserts that would split a merge
  also raise.

Issue: https://github.com/MHoroszowski/python-pptx/issues/12 (Phase 1).
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.oxml.table import CT_Table, CT_TableGrid, CT_TableRow
from pptx.table import Table, _Column, _Row
from pptx.util import Emu, Inches

from .unitutil.cxml import element

# ---------------------------------------------------------------------------
# OXML LAYER — CT_TableGrid / CT_Table / CT_TableRow new helpers
# ---------------------------------------------------------------------------


class DescribeCT_TableGrid_NewHelpers(object):
    """Unit-test suite for `CT_TableGrid.{insert_gridCol_at, remove_gridCol_at}`."""

    @pytest.mark.parametrize(
        ("idx", "expected_position"),
        [(0, 0), (1, 1), (2, 2)],  # ---head, middle, append
    )
    def it_can_insert_a_gridCol_at_a_specific_index(self, idx, expected_position):
        tblGrid = element("a:tblGrid/(a:gridCol{w=100},a:gridCol{w=200})")
        assert isinstance(tblGrid, CT_TableGrid)

        new_col = tblGrid.insert_gridCol_at(idx, Emu(999))

        assert tblGrid.gridCol_lst[expected_position] is new_col
        assert new_col.w == 999

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_insert_out_of_range(self, bad_idx):
        tblGrid = element("a:tblGrid/a:gridCol{w=100}")
        with pytest.raises(IndexError):
            tblGrid.insert_gridCol_at(bad_idx, Emu(999))

    def it_can_remove_a_gridCol_at_an_index(self):
        tblGrid = element("a:tblGrid/(a:gridCol{w=100},a:gridCol{w=200},a:gridCol{w=300})")

        tblGrid.remove_gridCol_at(1)

        assert [g.w for g in tblGrid.gridCol_lst] == [100, 300]

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_remove_out_of_range(self, bad_idx):
        tblGrid = element("a:tblGrid/a:gridCol{w=100}")
        with pytest.raises(IndexError):
            tblGrid.remove_gridCol_at(bad_idx)


class DescribeCT_Table_NewHelpers(object):
    """Unit-test suite for new `CT_Table` row insert/remove + cross-col merge helpers."""

    def it_can_insert_a_tr_at_a_specific_index(self):
        tbl = element(
            "a:tbl/(a:tblGrid,a:tr{h=100}/a:tc/a:txBody/a:p,a:tr{h=200}/a:tc/a:txBody/a:p)"
        )
        assert isinstance(tbl, CT_Table)

        new_tr = tbl.insert_tr_at(1, Emu(999))

        assert tbl.tr_lst[1] is new_tr
        assert new_tr.h == 999

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_insert_tr_out_of_range(self, bad_idx):
        tbl = element("a:tbl/(a:tblGrid,a:tr{h=100}/a:tc/a:txBody/a:p)")
        with pytest.raises(IndexError):
            tbl.insert_tr_at(bad_idx, Emu(999))

    def it_can_remove_a_tr_at_an_index(self):
        tbl = element(
            "a:tbl/(a:tblGrid,a:tr{h=100}/a:tc/a:txBody/a:p,"
            "a:tr{h=200}/a:tc/a:txBody/a:p,a:tr{h=300}/a:tc/a:txBody/a:p)"
        )

        tbl.remove_tr_at(1)

        assert [tr.h for tr in tbl.tr_lst] == [100, 300]

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_remove_tr_out_of_range(self, bad_idx):
        tbl = element("a:tbl/(a:tblGrid,a:tr{h=100}/a:tc/a:txBody/a:p)")
        with pytest.raises(IndexError):
            tbl.remove_tr_at(bad_idx)

    def it_detects_cross_column_merge_via_gridSpan(self):
        tbl = element(
            "a:tbl/(a:tblGrid/(a:gridCol{w=1},a:gridCol{w=2}),"
            "a:tr{h=10}/(a:tc{gridSpan=2}/a:txBody/a:p,a:tc{hMerge=1}/a:txBody/a:p))"
        )
        assert tbl.column_has_cross_column_merge(0) is True
        assert tbl.column_has_cross_column_merge(1) is True

    def it_does_not_flag_unmerged_columns(self):
        tbl = element(
            "a:tbl/(a:tblGrid/(a:gridCol{w=1},a:gridCol{w=2}),"
            "a:tr{h=10}/(a:tc/a:txBody/a:p,a:tc/a:txBody/a:p))"
        )
        assert tbl.column_has_cross_column_merge(0) is False
        assert tbl.column_has_cross_column_merge(1) is False


class DescribeCT_TableRow_NewHelpers(object):
    """Unit-test suite for `CT_TableRow.{insert_tc_at, remove_tc_at, has_cross_row_merge}`."""

    def it_can_insert_a_tc_at_a_specific_index(self):
        tr = element("a:tr{h=10}/(a:tc/a:txBody/a:p,a:tc/a:txBody/a:p)")
        assert isinstance(tr, CT_TableRow)

        new_tc = tr.insert_tc_at(1)

        assert tr.tc_lst[1] is new_tc
        assert len(tr.tc_lst) == 3

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_insert_tc_out_of_range(self, bad_idx):
        tr = element("a:tr{h=10}/a:tc/a:txBody/a:p")
        with pytest.raises(IndexError):
            tr.insert_tc_at(bad_idx)

    def it_can_remove_a_tc_at_an_index(self):
        tr = element("a:tr{h=10}/(a:tc/a:txBody/a:p,a:tc/a:txBody/a:p,a:tc/a:txBody/a:p)")

        tr.remove_tc_at(1)

        assert len(tr.tc_lst) == 2

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_remove_tc_out_of_range(self, bad_idx):
        tr = element("a:tr{h=10}/a:tc/a:txBody/a:p")
        with pytest.raises(IndexError):
            tr.remove_tc_at(bad_idx)

    def it_detects_cross_row_merge_via_rowSpan(self):
        tr = element("a:tr{h=10}/(a:tc{rowSpan=2}/a:txBody/a:p,a:tc/a:txBody/a:p)")
        assert tr.has_cross_row_merge is True

    def it_detects_cross_row_merge_via_vMerge(self):
        tr = element("a:tr{h=10}/(a:tc{vMerge=1}/a:txBody/a:p,a:tc/a:txBody/a:p)")
        assert tr.has_cross_row_merge is True

    def it_does_not_flag_a_row_without_vertical_merges(self):
        tr = element("a:tr{h=10}/(a:tc{gridSpan=2}/a:txBody/a:p,a:tc{hMerge=1}/a:txBody/a:p)")
        assert tr.has_cross_row_merge is False


# ---------------------------------------------------------------------------
# Integration — `Table.rows` / `Table.columns` add/remove on a real table
# ---------------------------------------------------------------------------


def _new_table(rows: int, cols: int) -> Table:
    """Return a fresh table with `rows` × `cols` empty cells, on a fresh slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(2))
    return shape.table


def _round_trip_table(table: Table) -> Table:
    """Save the table's presentation to bytes, reopen, and return the table."""
    prs = table._graphic_frame.part.package.presentation_part.presentation
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    reopened = Presentation(buf)
    # ---first table on the first slide---
    for shape in reopened.slides[0].shapes:
        if shape.has_table:
            return shape.table
    raise AssertionError("no table found in round-tripped slide")


class DescribeRowAdd(object):
    """`Table.rows.add(at, height)` — row insertion behavior."""

    def it_appends_a_row_when_at_is_None(self):
        table = _new_table(2, 3)

        new_row = table.rows.add()

        assert isinstance(new_row, _Row)
        assert len(table.rows) == 3
        # ---new row is the last row---
        assert table.rows[2]._tr is new_row._tr

    def it_can_insert_a_row_at_the_head(self):
        table = _new_table(2, 3)

        new_row = table.rows.add(at=0)

        assert len(table.rows) == 3
        assert table.rows[0]._tr is new_row._tr

    def it_can_insert_a_row_at_a_middle_index(self):
        table = _new_table(3, 3)

        new_row = table.rows.add(at=1)

        assert len(table.rows) == 4
        assert table.rows[1]._tr is new_row._tr

    def it_populates_the_new_row_with_empty_cells_matching_column_count(self):
        table = _new_table(2, 4)

        table.rows.add()

        new_row = table.rows[2]
        assert len(new_row.cells) == 4
        for cell in new_row.cells:
            assert cell.text == ""

    def it_inherits_height_from_the_first_row_when_height_is_None(self):
        table = _new_table(2, 3)
        first_row_h = table.rows[0].height

        table.rows.add()

        assert table.rows[2].height == first_row_h

    def it_uses_explicit_height_when_supplied(self):
        table = _new_table(2, 3)

        table.rows.add(height=Emu(123456))

        assert table.rows[2].height == 123456

    @pytest.mark.parametrize("bad_at", [-1, 99])
    def but_it_raises_on_at_out_of_range(self, bad_at):
        table = _new_table(2, 3)
        with pytest.raises(IndexError):
            table.rows.add(at=bad_at)


class DescribeRowRemove(object):
    """`Table.rows.remove(index)` — row removal behavior."""

    def it_can_remove_a_row(self):
        table = _new_table(3, 2)
        table.cell(1, 0).text = "row 1 cell 0"
        before_kept = table.cell(0, 0).text  # ---row 0 (will keep)
        also_kept = table.cell(2, 0).text  # ---row 2 (will keep)

        table.rows.remove(1)

        assert len(table.rows) == 2
        assert table.cell(0, 0).text == before_kept
        assert table.cell(1, 0).text == also_kept

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_index_out_of_range(self, bad_idx):
        table = _new_table(2, 2)
        with pytest.raises(IndexError):
            table.rows.remove(bad_idx)

    def but_it_raises_on_row_with_cross_row_merge_origin(self):
        """Removing a row whose cell originates a vertical merge would orphan."""
        table = _new_table(3, 2)
        # ---merge (0,0) into (1,0): origin at row 0 with rowSpan=2
        table.cell(0, 0).merge(table.cell(1, 0))

        with pytest.raises(ValueError):
            table.rows.remove(0)

    def but_it_raises_on_row_that_is_a_vertical_merge_target(self):
        table = _new_table(3, 2)
        # ---origin at row 0 with rowSpan=2; row 1 is the target (vMerge=True)
        table.cell(0, 0).merge(table.cell(1, 0))

        with pytest.raises(ValueError):
            table.rows.remove(1)


class DescribeColumnAdd(object):
    """`Table.columns.add(at, width)` — column insertion behavior."""

    def it_appends_a_column_when_at_is_None(self):
        table = _new_table(2, 3)

        new_column = table.columns.add()

        assert isinstance(new_column, _Column)
        assert len(table.columns) == 4

    def it_can_insert_a_column_at_the_head(self):
        table = _new_table(2, 3)

        new_column = table.columns.add(at=0)

        assert len(table.columns) == 4
        assert table.columns[0]._gridCol is new_column._gridCol

    def it_inserts_an_empty_cell_in_every_existing_row(self):
        table = _new_table(3, 2)
        table.cell(0, 0).text = "kept"
        table.cell(0, 1).text = "kept-too"

        table.columns.add(at=1)

        assert len(table.columns) == 3
        for r in range(3):
            assert len(table.rows[r].cells) == 3
        assert table.cell(0, 0).text == "kept"
        assert table.cell(0, 1).text == ""  # ---newly inserted, empty
        assert table.cell(0, 2).text == "kept-too"

    def it_inherits_width_from_the_first_column_when_width_is_None(self):
        table = _new_table(2, 3)
        first_col_w = table.columns[0].width

        table.columns.add()

        assert table.columns[3].width == first_col_w

    def it_uses_explicit_width_when_supplied(self):
        table = _new_table(2, 3)

        table.columns.add(width=Emu(987654))

        assert table.columns[3].width == 987654

    @pytest.mark.parametrize("bad_at", [-1, 99])
    def but_it_raises_on_at_out_of_range(self, bad_at):
        table = _new_table(2, 3)
        with pytest.raises(IndexError):
            table.columns.add(at=bad_at)


class DescribeColumnRemove(object):
    """`Table.columns.remove(index)` — column removal behavior."""

    def it_can_remove_a_column(self):
        table = _new_table(2, 3)
        table.cell(0, 0).text = "keep-A"
        table.cell(0, 2).text = "keep-C"

        table.columns.remove(1)

        assert len(table.columns) == 2
        for r in range(2):
            assert len(table.rows[r].cells) == 2
        assert table.cell(0, 0).text == "keep-A"
        assert table.cell(0, 1).text == "keep-C"

    @pytest.mark.parametrize("bad_idx", [-1, 99])
    def but_it_raises_on_index_out_of_range(self, bad_idx):
        table = _new_table(2, 2)
        with pytest.raises(IndexError):
            table.columns.remove(bad_idx)

    def but_it_raises_on_column_with_cross_column_merge_origin(self):
        """Removing column 0 of a horizontally-merged pair orphans (1,*)."""
        table = _new_table(2, 3)
        table.cell(0, 0).merge(table.cell(0, 1))

        with pytest.raises(ValueError):
            table.columns.remove(0)

    def but_it_raises_on_column_that_is_a_horizontal_merge_target(self):
        table = _new_table(2, 3)
        table.cell(0, 0).merge(table.cell(0, 1))

        with pytest.raises(ValueError):
            table.columns.remove(1)


# ---------------------------------------------------------------------------
# Round-trip integration — open → mutate → save → reopen
# ---------------------------------------------------------------------------


class DescribeTablesRoundTrip(object):
    """Save → reopen preservation of row/column CRUD operations."""

    def it_round_trips_a_row_append(self):
        table = _new_table(2, 3)
        table.cell(0, 0).text = "header-A"

        table.rows.add()
        rt = _round_trip_table(table)

        assert len(rt.rows) == 3
        assert rt.cell(0, 0).text == "header-A"

    def it_round_trips_a_column_append(self):
        table = _new_table(2, 2)
        table.cell(0, 0).text = "kept"

        table.columns.add()
        rt = _round_trip_table(table)

        assert len(rt.columns) == 3
        assert rt.cell(0, 0).text == "kept"

    def it_round_trips_a_row_removal(self):
        table = _new_table(3, 2)
        table.cell(0, 0).text = "row0"
        table.cell(2, 0).text = "row2"

        table.rows.remove(1)
        rt = _round_trip_table(table)

        assert len(rt.rows) == 2
        assert rt.cell(0, 0).text == "row0"
        assert rt.cell(1, 0).text == "row2"

    def it_round_trips_a_column_removal(self):
        table = _new_table(2, 3)
        table.cell(0, 0).text = "A"
        table.cell(0, 2).text = "C"

        table.columns.remove(1)
        rt = _round_trip_table(table)

        assert len(rt.columns) == 2
        assert rt.cell(0, 0).text == "A"
        assert rt.cell(0, 1).text == "C"

    def it_round_trips_an_indexed_row_insert(self):
        table = _new_table(2, 2)
        table.cell(0, 0).text = "first"
        table.cell(1, 0).text = "last"

        table.rows.add(at=1)
        rt = _round_trip_table(table)

        assert len(rt.rows) == 3
        assert rt.cell(0, 0).text == "first"
        assert rt.cell(1, 0).text == ""  # ---inserted empty
        assert rt.cell(2, 0).text == "last"
