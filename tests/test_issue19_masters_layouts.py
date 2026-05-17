"""Integration tests for issue #19 — Slide Masters, Layouts & .potx Templates.

These exercise the public API end-to-end (build → mutate → save → reopen)
rather than mocking, because the failure modes here (PowerPoint repair
dialogs, dangling rels, content-type mismatches) only surface on a real
round-trip. Mirrors the test discipline used for the customXml and
slide-CRUD epics in this fork.
"""

from __future__ import annotations

import io
import zipfile

import pytest

from pptx import Presentation


def _craft_potx_bytes() -> io.BytesIO:
    """Return an in-memory .potx: a default deck with the template content-type."""
    base = io.BytesIO()
    Presentation().save(base)
    base.seek(0)
    src = {}
    with zipfile.ZipFile(base) as z:
        for n in z.namelist():
            src[n] = z.read(n)
    src["[Content_Types].xml"] = src["[Content_Types].xml"].replace(
        b"presentationml.presentation.main+xml",
        b"presentationml.template.main+xml",
    )
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for n, b in src.items():
            z.writestr(n, b)
    out.seek(0)
    return out


class DescribePotxRead:
    """SF1 — Presentation() accepts .potx files."""

    def it_opens_a_potx_template_without_error(self):
        prs = Presentation(_craft_potx_bytes())
        assert prs is not None

    def it_exposes_masters_and_layouts_from_a_potx(self):
        prs = Presentation(_craft_potx_bytes())
        assert len(prs.slide_masters) >= 1
        assert len(prs.slide_layouts) >= 1

    def but_it_still_rejects_a_genuinely_non_pptx_payload(self):
        bogus = io.BytesIO(b"PK\x03\x04 not really an office file")
        with pytest.raises(Exception):
            Presentation(bogus)


class DescribePotxRoundTrip:
    """SF1 regression — a .potx survives an open→save→reopen cycle."""

    def it_round_trips_a_potx_through_save(self):
        prs = Presentation(_craft_potx_bytes())
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        reopened = Presentation(buf)
        assert len(reopened.slide_masters) >= 1


class DescribeSlideMasterGetLayout:
    """SF9 — SlideMaster.get_layout(slide_layout_id, default=None)."""

    def _master_and_first_layout_id(self):
        prs = Presentation()
        master = prs.slide_masters[0]
        sldLayoutIdLst = master._element.get_or_add_sldLayoutIdLst()
        first = sldLayoutIdLst.sldLayoutId_lst[0]
        return prs, master, first.id

    def it_returns_the_layout_matching_a_known_id(self):
        prs, master, layout_id = self._master_and_first_layout_id()
        if layout_id is None:
            pytest.skip("default template layout has no @id attribute")
        got = master.get_layout(layout_id)
        assert got is not None
        assert got == master.slide_layouts[0]

    def it_returns_default_for_an_unknown_id(self):
        prs = Presentation()
        master = prs.slide_masters[0]
        sentinel = object()
        assert master.get_layout(999999, default=sentinel) is sentinel

    def and_it_returns_None_by_default_for_an_unknown_id(self):
        prs = Presentation()
        master = prs.slide_masters[0]
        assert master.get_layout(424242) is None

    def it_does_not_raise_on_a_bad_id(self):
        prs = Presentation()
        master = prs.slide_masters[0]
        # must not raise KeyError/IndexError — returns default
        assert master.get_layout(-1) is None


def _P_NS():
    from pptx.oxml.ns import qn

    return qn


class DescribeAddLayout:
    """SF3 — SlideLayouts.add_layout(name=None) creates a new p:sldLayout part.

    Manual semantic port of upstream scanny/python-pptx#1091 onto the
    fork's ruff-formatted base (issue #19 SF3, ISA ISC-12..22).
    """

    def it_returns_a_SlideLayout_object(self):
        prs = Presentation()
        layouts = prs.slide_masters[0].slide_layouts
        new_layout = layouts.add_layout()
        from pptx.slide import SlideLayout

        assert isinstance(new_layout, SlideLayout)

    def it_increments_the_master_layout_count_by_exactly_one(self):
        prs = Presentation()
        layouts = prs.slide_masters[0].slide_layouts
        before = len(layouts)
        layouts.add_layout()
        assert len(layouts) == before + 1

    def it_sets_the_name_when_given(self):
        prs = Presentation()
        layouts = prs.slide_masters[0].slide_layouts
        new_layout = layouts.add_layout(name="My Custom Layout")
        assert new_layout.name == "My Custom Layout"

    def it_assigns_a_sensible_default_name_when_none_given(self):
        prs = Presentation()
        layouts = prs.slide_masters[0].slide_layouts
        new_layout = layouts.add_layout()
        assert new_layout.name != ""
        assert new_layout.name is not None

    def it_finds_the_new_layout_by_name(self):
        prs = Presentation()
        layouts = prs.slide_masters[0].slide_layouts
        layouts.add_layout(name="Lookup Target")
        assert layouts.get_by_name("Lookup Target") is not None

    def it_allocates_a_sldLayoutId_in_the_high_uint_range(self):
        """@id must be >= 2147483648 (0x80000000), PowerPoint's convention.

        The default template's own layout ids run 2147483649..2147483659.
        A new layout id below that floor would collide with the low
        ``p:sldId/@id`` pool (slide ids start at 256) — the regression that
        produced a "PowerPoint found a problem" repair dialog and was caught
        only by Interceptor visual verification.
        """
        from pptx.parts.slide import _OOXML_LAYOUT_ID_FLOOR

        prs = Presentation()
        master = prs.slide_masters[0]
        layouts = master.slide_layouts
        layouts.add_layout(name="Range Probe")
        qn = _P_NS()
        idLst = master._element.get_or_add_sldLayoutIdLst()
        ids = [
            int(sli.get("id"))
            for sli in idLst.findall(qn("p:sldLayoutId"))
            if sli.get("id") is not None
        ]
        assert ids, "expected at least one sldLayoutId with an @id"
        new_id = ids[-1]
        assert new_id >= _OOXML_LAYOUT_ID_FLOOR
        assert new_id <= 4294967295

    def it_does_not_collide_with_the_shared_id_pool(self):
        """Regression guard for the repair-dialog bug.

        ``p:sldMasterId/@id``, ``p:sldLayoutId/@id`` AND ``p:sldId/@id`` are
        ONE shared pool in PowerPoint's repair heuristic. The new layout id
        must be disjoint from BOTH the master ids AND the slide ids. The
        original SF3 test only checked sldMasterId and so missed the
        sldId(256) collision.
        """
        prs = Presentation()
        master = prs.slide_masters[0]
        layouts = master.slide_layouts
        new_layout = layouts.add_layout(name="Collision Probe")
        # a slide must exist so a low sldId (256) is present in the pool
        prs.slides.add_slide(new_layout)
        qn = _P_NS()

        idLst = master._element.get_or_add_sldLayoutIdLst()
        layout_ids = {
            int(sli.get("id"))
            for sli in idLst.findall(qn("p:sldLayoutId"))
            if sli.get("id") is not None
        }

        pres_el = prs.part._element
        pool = set()
        smIdLst = pres_el.find(qn("p:sldMasterIdLst"))
        if smIdLst is not None:
            for smi in smIdLst.findall(qn("p:sldMasterId")):
                raw = smi.get("id")
                if raw is not None:
                    pool.add(int(raw))
        sldIdLst = pres_el.find(qn("p:sldIdLst"))
        if sldIdLst is not None:
            for sid in sldIdLst.findall(qn("p:sldId")):
                raw = sid.get("id")
                if raw is not None:
                    pool.add(int(raw))

        assert pool, "fixture expected at least one sldMasterId + sldId"
        assert 256 in pool, "the slide's sldId(256) must be in the pool"
        assert layout_ids.isdisjoint(pool), (
            "new layout id collides with sldMasterId/sldId pool — repair-dialog bug"
        )

    def it_survives_a_save_reopen_round_trip(self):
        prs = Presentation()
        layouts = prs.slide_masters[0].slide_layouts
        layouts.add_layout(name="RoundTrip Layout")
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        layouts2 = prs2.slide_masters[0].slide_layouts
        assert layouts2.get_by_name("RoundTrip Layout") is not None

    def it_can_be_used_as_the_basis_for_a_new_slide(self):
        prs = Presentation()
        new_layout = prs.slide_masters[0].slide_layouts.add_layout(name="Slide Source")
        before = len(prs.slides)
        prs.slides.add_slide(new_layout)
        assert len(prs.slides) == before + 1

    def it_allows_adding_a_placeholder_to_the_new_layout(self):
        """LayoutShapes now inherits group-shape add_* (closes upstream #1044)."""
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        new_layout = prs.slide_masters[0].slide_layouts.add_layout(name="PH Layout")
        before = len(new_layout.placeholders)
        ph = new_layout.shapes.add_placeholder(PP_PLACEHOLDER.BODY, "horz", "full")
        assert ph is not None
        assert len(new_layout.placeholders) == before + 1

    def it_round_trips_a_placeholder_added_to_a_new_layout(self):
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        new_layout = prs.slide_masters[0].slide_layouts.add_layout(name="PH RT")
        new_layout.shapes.add_placeholder(PP_PLACEHOLDER.BODY, "horz", "full")
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        layout2 = prs2.slide_masters[0].slide_layouts.get_by_name("PH RT")
        assert layout2 is not None
        assert len(layout2.placeholders) >= 1


def _content_types_xml(buf: io.BytesIO) -> bytes:
    """Return the raw `[Content_Types].xml` bytes from a saved package buffer."""
    buf.seek(0)
    with zipfile.ZipFile(buf) as z:
        return z.read("[Content_Types].xml")


_PML_PRESENTATION_MAIN = (
    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)
_PML_TEMPLATE_MAIN = (
    b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
)


class DescribeSaveAsPotx:
    """SF2 — Presentation.save_as_potx(path) writes a template content-type.

    The output `[Content_Types].xml` carries the template main+xml override
    for the presentation part, WITHOUT mutating the in-memory package
    (ISC-6..11). Manual semantic port (issue #19 SF2).
    """

    def it_exposes_a_save_as_potx_method(self):
        prs = Presentation()
        assert hasattr(prs, "save_as_potx")
        assert callable(prs.save_as_potx)

    def it_writes_the_template_content_type_to_the_output(self):
        prs = Presentation()
        buf = io.BytesIO()
        prs.save_as_potx(buf)
        ct_xml = _content_types_xml(buf)
        assert _PML_TEMPLATE_MAIN in ct_xml
        assert _PML_PRESENTATION_MAIN not in ct_xml

    def it_produces_a_package_that_round_trips_through_Presentation(self):
        prs = Presentation()
        buf = io.BytesIO()
        prs.save_as_potx(buf)
        buf.seek(0)
        reopened = Presentation(buf)
        assert len(reopened.slide_masters) >= 1
        assert len(reopened.slide_layouts) >= 1

    def it_does_not_mutate_the_in_memory_presentation_content_type(self):
        prs = Presentation()
        before = prs.part.content_type
        buf = io.BytesIO()
        prs.save_as_potx(buf)
        after = prs.part.content_type
        assert before == after
        assert after == _PML_PRESENTATION_MAIN.decode("ascii")

    def it_accepts_a_str_path(self, tmp_path):
        prs = Presentation()
        out = tmp_path / "template_out.potx"
        prs.save_as_potx(str(out))
        with zipfile.ZipFile(str(out)) as z:
            ct_xml = z.read("[Content_Types].xml")
        assert _PML_TEMPLATE_MAIN in ct_xml

    def it_leaves_a_normal_save_unaffected_afterwards(self):
        prs = Presentation()
        potx_buf = io.BytesIO()
        prs.save_as_potx(potx_buf)
        pptx_buf = io.BytesIO()
        prs.save(pptx_buf)
        ct_xml = _content_types_xml(pptx_buf)
        assert _PML_PRESENTATION_MAIN in ct_xml
        assert _PML_TEMPLATE_MAIN not in ct_xml


class DescribeMasterShapeAuthoring:
    """SF5 — author shapes directly on a SlideMaster via master.shapes.

    `MasterShapes` reparented to `_BaseGroupShapes`, gaining
    add_textbox/add_picture/add_shape (issue #19 SF5, ISC-30..36).
    """

    def _png_bytes(self) -> io.BytesIO:
        # ---minimal 1x1 PNG---
        import base64

        data = base64.b64decode(
            b"iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAC0lEQVR4nGNgAAIAAAUAAen63NgAAAAASUVORK5CYII="
        )
        return io.BytesIO(data)

    def it_exposes_group_shape_add_methods_on_master_shapes(self):
        prs = Presentation()
        shapes = prs.slide_masters[0].shapes
        assert hasattr(shapes, "add_textbox")
        assert hasattr(shapes, "add_picture")
        assert hasattr(shapes, "add_shape")

    def it_adds_a_textbox_to_a_master_that_survives_round_trip(self):
        prs = Presentation()
        master = prs.slide_masters[0]
        tb = master.shapes.add_textbox(0, 0, 914400, 457200)
        tb.text_frame.text = "MASTER MARKER TEXT"
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        texts = [s.text_frame.text for s in prs2.slide_masters[0].shapes if s.has_text_frame]
        assert "MASTER MARKER TEXT" in texts

    def it_adds_an_autoshape_to_a_master(self):
        from pptx.enum.shapes import MSO_SHAPE

        prs = Presentation()
        master = prs.slide_masters[0]
        before = len(list(master.shapes))
        master.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0, 0, 914400, 914400)
        assert len(list(master.shapes)) == before + 1

    def it_adds_a_picture_to_a_master(self):
        prs = Presentation()
        master = prs.slide_masters[0]
        before = len(list(master.shapes))
        master.shapes.add_picture(self._png_bytes(), 0, 0, 914400, 914400)
        assert len(list(master.shapes)) == before + 1

    def it_increments_the_master_shape_count(self):
        prs = Presentation()
        master = prs.slide_masters[0]
        before = len(list(master.shapes))
        master.shapes.add_textbox(0, 0, 914400, 457200)
        master.shapes.add_textbox(0, 457200, 914400, 457200)
        assert len(list(master.shapes)) == before + 2


class DescribeCopyFromLayout:
    """SF4 — SlideLayouts.copy_from(other_layout) duplicates a layout.

    Deep-copies the source layout's spTree shapes into a fresh layout
    created via the SF3 add_layout machinery (issue #19 SF4, ISC-23..29).
    """

    def _layout_with_a_shape(self):
        prs = Presentation()
        layouts = prs.slide_masters[0].slide_layouts
        src = layouts.add_layout(name="CopySource")
        src.shapes.add_textbox(0, 0, 914400, 457200)
        return prs, layouts, src

    def it_returns_a_new_SlideLayout(self):
        from pptx.slide import SlideLayout

        prs, layouts, src = self._layout_with_a_shape()
        copy = layouts.copy_from(src)
        assert isinstance(copy, SlideLayout)
        assert copy is not src

    def it_copies_the_source_shape_count(self):
        prs, layouts, src = self._layout_with_a_shape()
        src_count = len(list(src.shapes))
        copy = layouts.copy_from(src)
        assert len(list(copy.shapes)) == src_count

    def it_preserves_placeholder_idx_and_type(self):
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        layouts = prs.slide_masters[0].slide_layouts
        src = layouts.add_layout(name="PHSource")
        src.shapes.add_placeholder(PP_PLACEHOLDER.BODY, "horz", "full")
        src_ph = [(p.element.ph_idx, p.element.ph_type) for p in src.placeholders]
        copy = layouts.copy_from(src)
        copy_ph = [(p.element.ph_idx, p.element.ph_type) for p in copy.placeholders]
        assert sorted(map(str, copy_ph)) == sorted(map(str, src_ph))

    def it_does_not_mutate_the_source_layout(self):
        prs, layouts, src = self._layout_with_a_shape()
        src_count_before = len(list(src.shapes))
        src_name_before = src.name
        layouts.copy_from(src)
        assert len(list(src.shapes)) == src_count_before
        assert src.name == src_name_before

    def it_survives_a_save_reopen_round_trip(self):
        prs, layouts, src = self._layout_with_a_shape()
        copy = layouts.copy_from(src)
        copy.name = "CopiedLayoutRT"
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        layout2 = prs2.slide_masters[0].slide_layouts.get_by_name("CopiedLayoutRT")
        assert layout2 is not None
        assert len(list(layout2.shapes)) == len(list(src.shapes))


class DescribeLayoutPlaceholdersAdd:
    """SF6 — SlideLayout.placeholders.add(idx, ph_type, ...).

    Adds a `<p:sp>` placeholder to a layout's shape tree, readable
    afterward and surviving round-trip (issue #19 SF6, ISC-37..43).
    """

    def it_adds_a_placeholder_to_the_layout(self):
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        layout = prs.slide_masters[0].slide_layouts.add_layout(name="PHAdd")
        before = len(layout.placeholders)
        layout.placeholders.add(11, PP_PLACEHOLDER.BODY)
        assert len(layout.placeholders) == before + 1

    def it_writes_the_idx_and_type(self):
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        layout = prs.slide_masters[0].slide_layouts.add_layout(name="PHAdd2")
        ph = layout.placeholders.add(12, PP_PLACEHOLDER.BODY)
        assert ph.element.ph_idx == 12
        assert ph.element.ph_type == PP_PLACEHOLDER.BODY

    def it_is_readable_through_the_collection_after_add(self):
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        layout = prs.slide_masters[0].slide_layouts.add_layout(name="PHAdd3")
        layout.placeholders.add(13, PP_PLACEHOLDER.BODY)
        idxs = [p.element.ph_idx for p in layout.placeholders]
        assert 13 in idxs

    def it_survives_a_save_reopen_round_trip(self):
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        layout = prs.slide_masters[0].slide_layouts.add_layout(name="PHAddRT")
        layout.placeholders.add(14, PP_PLACEHOLDER.BODY)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        layout2 = prs2.slide_masters[0].slide_layouts.get_by_name("PHAddRT")
        assert layout2 is not None
        idxs = [p.element.ph_idx for p in layout2.placeholders]
        assert 14 in idxs

    def it_rejects_a_duplicate_idx_on_the_same_layout(self):
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        layout = prs.slide_masters[0].slide_layouts.add_layout(name="PHDup")
        layout.placeholders.add(15, PP_PLACEHOLDER.BODY)
        with pytest.raises(ValueError):
            layout.placeholders.add(15, PP_PLACEHOLDER.BODY)


class DescribeApplyLayout:
    """SF7 — assign a slide a layout owned by a (possibly different) master.

    Re-points the slide→layout relationship via the ``Slide.slide_layout``
    setter / ``Slide.apply_layout`` method. The rel chain
    slide→layout→master(theme) must stay intact and the slide's prior
    layout's master must NOT be orphaned (issue #19 SF7; ISC-44..49).
    """

    def _deck_with_two_layouts(self):
        prs = Presentation()
        master = prs.slide_masters[0]
        base_layout = prs.slide_layouts[0]
        target_layout = master.slide_layouts.add_layout(name="SF7Target")
        slide = prs.slides.add_slide(base_layout)
        return prs, master, base_layout, target_layout, slide

    def it_repoints_the_slide_layout_via_the_setter(self):
        prs, master, base_layout, target_layout, slide = self._deck_with_two_layouts()
        assert slide.slide_layout.name != "SF7Target"
        slide.slide_layout = target_layout
        assert slide.slide_layout.name == "SF7Target"

    def it_also_exposes_an_apply_layout_method(self):
        prs, master, base_layout, target_layout, slide = self._deck_with_two_layouts()
        slide.apply_layout(target_layout)
        assert slide.slide_layout.name == "SF7Target"

    def it_keeps_the_layout_master_chain_intact(self):
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        prs, master, base_layout, target_layout, slide = self._deck_with_two_layouts()
        slide.slide_layout = target_layout
        # ---slide resolves a layout, layout resolves a master, no raise---
        resolved_layout = slide.slide_layout
        resolved_master = resolved_layout.slide_master
        assert resolved_master is not None
        # ---exactly one slide→layout rel after the repoint (no dangling)---
        layout_rels = [
            rel
            for rel in slide.part.rels.values()
            if not rel.is_external and rel.reltype == RT.SLIDE_LAYOUT
        ]
        assert len(layout_rels) == 1

    def it_does_not_orphan_the_prior_layouts_master(self):
        prs, master, base_layout, target_layout, slide = self._deck_with_two_layouts()
        # ---a second slide still uses the original base layout---
        other_slide = prs.slides.add_slide(base_layout)
        slide.slide_layout = target_layout
        # ---the other slide's layout + master are still fully reachable---
        assert other_slide.slide_layout.name == base_layout.name
        assert other_slide.slide_layout.slide_master is not None
        # ---the base layout is still in the master's collection---
        names = [lay.name for lay in master.slide_layouts]
        assert base_layout.name in names

    def it_survives_a_save_reopen_round_trip(self):
        prs, master, base_layout, target_layout, slide = self._deck_with_two_layouts()
        slide.slide_layout = target_layout
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        assert prs2.slides[0].slide_layout.name == "SF7Target"
        # ---the round-tripped slide still resolves its master---
        assert prs2.slides[0].slide_layout.slide_master is not None

    def it_is_idempotent_when_applied_twice(self):
        prs, master, base_layout, target_layout, slide = self._deck_with_two_layouts()
        slide.slide_layout = target_layout
        slide.slide_layout = target_layout
        assert slide.slide_layout.name == "SF7Target"

    def it_rejects_a_non_layout_argument(self):
        prs, master, base_layout, target_layout, slide = self._deck_with_two_layouts()
        with pytest.raises(TypeError):
            slide.slide_layout = "not a layout"


class DescribeInsertChartIntoPlaceholder:
    """SF8 — insert a chart into a chart-capable placeholder.

    A CHART placeholder is replaced by a `<p:graphicFrame>` holding the
    chart, sized/positioned from the placeholder. Non-chart placeholders
    reject ``insert_chart`` cleanly. Upstream scanny/python-pptx#199
    (issue #19 SF8; ISC-50..55).
    """

    def _slide_with_chart_placeholder(self):
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        master = prs.slide_masters[0]
        layout = master.slide_layouts.add_layout(name="SF8Chart")
        layout.placeholders.add(
            10,
            PP_PLACEHOLDER.CHART,
            left=914400,
            top=914400,
            width=4572000,
            height=2743200,
        )
        slide = prs.slides.add_slide(layout)
        chart_ph = next(
            p for p in slide.placeholders if p.placeholder_format.type == PP_PLACEHOLDER.CHART
        )
        return prs, slide, chart_ph

    def _sample_chart_data(self):
        from pptx.chart.data import CategoryChartData

        chart_data = CategoryChartData()
        chart_data.categories = ["East", "West", "Midwest"]
        chart_data.add_series("Q1 Sales", (19.2, 21.4, 16.7))
        return chart_data

    def it_replaces_the_placeholder_with_a_chart_graphic_frame(self):
        from pptx.enum.chart import XL_CHART_TYPE

        prs, slide, chart_ph = self._slide_with_chart_placeholder()
        result = chart_ph.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, self._sample_chart_data())
        assert result.has_chart is True

    def it_creates_a_chart_part_and_relationship(self):
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        prs, slide, chart_ph = self._slide_with_chart_placeholder()
        chart_rels_before = [r for r in slide.part.rels.values() if r.reltype == RT.CHART]
        chart_ph.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, self._sample_chart_data())
        chart_rels_after = [r for r in slide.part.rels.values() if r.reltype == RT.CHART]
        assert len(chart_rels_after) == len(chart_rels_before) + 1

    def it_exposes_the_chart_on_the_returned_frame(self):
        from pptx.enum.chart import XL_CHART_TYPE

        prs, slide, chart_ph = self._slide_with_chart_placeholder()
        result = chart_ph.insert_chart(XL_CHART_TYPE.PIE, self._sample_chart_data())
        assert result.chart.chart_type == XL_CHART_TYPE.PIE

    def it_survives_a_save_reopen_round_trip(self):
        from pptx.enum.chart import XL_CHART_TYPE

        prs, slide, chart_ph = self._slide_with_chart_placeholder()
        chart_ph.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, self._sample_chart_data())
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        slide2 = prs2.slides[0]
        charts = [s for s in slide2.shapes if s.has_chart]
        assert len(charts) == 1

    def it_positions_the_frame_from_the_placeholder(self):
        from pptx.enum.chart import XL_CHART_TYPE

        prs, slide, chart_ph = self._slide_with_chart_placeholder()
        left, top, width, height = (
            chart_ph.left,
            chart_ph.top,
            chart_ph.width,
            chart_ph.height,
        )
        result = chart_ph.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, self._sample_chart_data())
        assert (result.left, result.top, result.width, result.height) == (
            left,
            top,
            width,
            height,
        )

    def but_a_non_chart_placeholder_rejects_insert_chart(self):
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.enum.shapes import PP_PLACEHOLDER

        prs = Presentation()
        master = prs.slide_masters[0]
        layout = master.slide_layouts.add_layout(name="SF8NonChart")
        layout.placeholders.add(20, PP_PLACEHOLDER.BODY, left=0, top=0, width=914400, height=457200)
        slide = prs.slides.add_slide(layout)
        body_ph = next(
            p for p in slide.placeholders if p.placeholder_format.type == PP_PLACEHOLDER.BODY
        )
        with pytest.raises(TypeError):
            body_ph.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, self._sample_chart_data())
