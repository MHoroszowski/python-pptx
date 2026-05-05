# pyright: reportPrivateUsage=false

"""End-to-end test suite for `pptx.custom_properties.CustomProperties`."""

from __future__ import annotations

import datetime as dt
from io import BytesIO

import pytest

from pptx import Presentation
from pptx.custom_properties import CustomProperties
from pptx.parts.custom_properties import CustomPropertiesPart


@pytest.fixture
def empty_prs():
    """Return a fresh Presentation built from the default template."""
    return Presentation()


def _roundtrip(prs):
    """Save `prs` to BytesIO and return a freshly-reloaded Presentation."""
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


class DescribeCustomProperties_Mapping:
    def it_starts_empty_for_a_default_presentation(self, empty_prs):
        cp = empty_prs.custom_properties
        assert isinstance(cp, CustomProperties)
        assert len(cp) == 0
        assert list(cp) == []
        assert "anything" not in cp

    def it_writes_and_reads_a_string_value(self, empty_prs):
        empty_prs.custom_properties["Source"] = "cli@1.4"
        assert empty_prs.custom_properties["Source"] == "cli@1.4"
        assert "Source" in empty_prs.custom_properties

    def it_dispatches_value_types_on_assignment(self, empty_prs):
        empty_prs.custom_properties["S"] = "string"
        empty_prs.custom_properties["I"] = 42
        empty_prs.custom_properties["F"] = 3.14
        empty_prs.custom_properties["B"] = True
        empty_prs.custom_properties["D"] = dt.datetime(2026, 5, 5, 14, 0, 0)

        assert empty_prs.custom_properties["S"] == "string"
        assert empty_prs.custom_properties["I"] == 42
        assert empty_prs.custom_properties["F"] == pytest.approx(3.14)
        assert empty_prs.custom_properties["B"] is True
        assert empty_prs.custom_properties["D"] == dt.datetime(2026, 5, 5, 14, 0, 0)

    def it_replaces_an_existing_value_on_repeated_assignment(self, empty_prs):
        empty_prs.custom_properties["X"] = "old"
        empty_prs.custom_properties["X"] = "new"
        assert empty_prs.custom_properties["X"] == "new"
        assert len(empty_prs.custom_properties) == 1

    def it_replaces_value_with_a_different_type(self, empty_prs):
        empty_prs.custom_properties["X"] = "hello"
        empty_prs.custom_properties["X"] = 42
        assert empty_prs.custom_properties["X"] == 42

    def it_raises_KeyError_on_missing_lookup(self, empty_prs):
        with pytest.raises(KeyError):
            empty_prs.custom_properties["missing"]

    def it_deletes_a_property(self, empty_prs):
        empty_prs.custom_properties["X"] = "a"
        del empty_prs.custom_properties["X"]
        assert "X" not in empty_prs.custom_properties

    def it_raises_KeyError_on_delete_missing(self, empty_prs):
        with pytest.raises(KeyError):
            del empty_prs.custom_properties["missing"]

    def it_supports_iter_keys_values_items_get(self, empty_prs):
        empty_prs.custom_properties["a"] = "1"
        empty_prs.custom_properties["b"] = 2
        empty_prs.custom_properties["c"] = True

        assert list(empty_prs.custom_properties) == ["a", "b", "c"]
        assert list(empty_prs.custom_properties.keys()) == ["a", "b", "c"]
        assert dict(empty_prs.custom_properties.items()) == {"a": "1", "b": 2, "c": True}
        assert empty_prs.custom_properties.get("missing") is None
        assert empty_prs.custom_properties.get("missing", "default") == "default"

    def it_raises_TypeError_on_unsupported_value(self, empty_prs):
        with pytest.raises(TypeError):
            empty_prs.custom_properties["X"] = object()  # type: ignore[assignment]


class DescribeCustomProperties_edge_cases:
    def it_returns_False_for_non_string_contains(self, empty_prs):
        empty_prs.custom_properties["X"] = "v"
        assert (42 in empty_prs.custom_properties) is False  # type: ignore[operator]

    def it_treats_a_property_with_no_value_child_as_absent(self, empty_prs):
        # Force a malformed entry: an op:property element with no vt:* child.
        # The lookup returns None → CustomProperties surfaces it as KeyError.
        from pptx.oxml import parse_xml
        from pptx.oxml.custom_properties import DEFAULT_FMTID
        from pptx.oxml.ns import nsdecls

        cp_part = empty_prs.part.package.custom_properties_part
        # Replace _element with a malformed Properties root containing one
        # property that has no value child.
        broken = parse_xml(
            (
                "<op:Properties %s>"
                '<op:property fmtid="%s" pid="2" name="empty"/>'
                "</op:Properties>" % (nsdecls("op", "vt"), DEFAULT_FMTID)
            ).encode()
        )
        cp_part._element = broken
        with pytest.raises(KeyError):
            _ = empty_prs.custom_properties["empty"]


class DescribeCustomProperties_explicit_setters:
    def it_writes_string_with_set_string(self, empty_prs):
        # set_string("X", "42") writes vt:lpwstr, not vt:i4
        empty_prs.custom_properties.set_string("X", "42")
        assert empty_prs.custom_properties["X"] == "42"
        # confirm the underlying element is vt:lpwstr
        prop = empty_prs.part.package.custom_properties_part.get_property("X")
        assert prop is not None
        assert prop.lpwstr is not None
        assert prop.i4 is None

    def it_writes_int_with_set_int_rejecting_bool(self, empty_prs):
        empty_prs.custom_properties.set_int("X", 5)
        assert empty_prs.custom_properties["X"] == 5
        with pytest.raises(TypeError):
            empty_prs.custom_properties.set_int("X", True)  # type: ignore[arg-type]

    def it_writes_float_with_set_float(self, empty_prs):
        empty_prs.custom_properties.set_float("X", 3.14)
        prop = empty_prs.part.package.custom_properties_part.get_property("X")
        assert prop is not None
        assert prop.r8 is not None

    def it_writes_bool_with_set_bool(self, empty_prs):
        empty_prs.custom_properties.set_bool("X", False)
        assert empty_prs.custom_properties["X"] is False
        with pytest.raises(TypeError):
            empty_prs.custom_properties.set_bool("X", 0)  # type: ignore[arg-type]

    def it_writes_datetime_with_set_datetime(self, empty_prs):
        empty_prs.custom_properties.set_datetime(
            "When", dt.datetime(2026, 1, 1, 0, 0, 0)
        )
        assert empty_prs.custom_properties["When"] == dt.datetime(2026, 1, 1, 0, 0, 0)

    def it_rejects_set_string_with_non_string(self, empty_prs):
        with pytest.raises(TypeError):
            empty_prs.custom_properties.set_string("X", 42)  # type: ignore[arg-type]

    def it_rejects_set_float_with_bool_or_non_number(self, empty_prs):
        with pytest.raises(TypeError):
            empty_prs.custom_properties.set_float("X", True)  # type: ignore[arg-type]
        with pytest.raises(TypeError):
            empty_prs.custom_properties.set_float("X", "1.0")  # type: ignore[arg-type]

    def it_rejects_set_datetime_with_non_datetime(self, empty_prs):
        with pytest.raises(TypeError):
            empty_prs.custom_properties.set_datetime("X", "today")  # type: ignore[arg-type]

    def it_overwrites_an_existing_value_via_set_string(self, empty_prs):
        empty_prs.custom_properties["X"] = 42
        empty_prs.custom_properties.set_string("X", "now-a-string")
        assert empty_prs.custom_properties["X"] == "now-a-string"


class DescribeCustomProperties_lazy_creation:
    def it_creates_the_part_on_first_access(self, empty_prs):
        # default presentation has no custom_properties_part yet — the lazy
        # access path must create one.
        cp_part = empty_prs.part.package.custom_properties_part
        assert isinstance(cp_part, CustomPropertiesPart)
        # Mapping wrapper finds it
        assert isinstance(empty_prs.custom_properties, CustomProperties)

    def it_returns_the_same_wrapper_class_each_call(self, empty_prs):
        # Different instances are fine (CustomProperties is a thin facade) —
        # what matters is that they wrap the same underlying part.
        a = empty_prs.custom_properties
        b = empty_prs.custom_properties
        assert a._part is b._part


class DescribeCustomProperties_roundtrip:
    def it_round_trips_through_save_load(self, empty_prs):
        empty_prs.custom_properties["Source"] = "cli@1.4.2"
        empty_prs.custom_properties["BuildNumber"] = 42
        empty_prs.custom_properties["IsDraft"] = True
        empty_prs.custom_properties["At"] = dt.datetime(2026, 5, 5, 14, 0, 0)

        reloaded = _roundtrip(empty_prs)

        assert reloaded.custom_properties["Source"] == "cli@1.4.2"
        assert reloaded.custom_properties["BuildNumber"] == 42
        assert reloaded.custom_properties["IsDraft"] is True
        assert reloaded.custom_properties["At"] == dt.datetime(2026, 5, 5, 14, 0, 0)

    def it_preserves_core_properties_alongside_custom_ones(self, empty_prs):
        # both can coexist; custom_properties is /docProps/custom.xml,
        # core_properties is /docProps/core.xml — distinct parts
        empty_prs.core_properties.author = "Athena"
        empty_prs.custom_properties["Source"] = "cli"
        reloaded = _roundtrip(empty_prs)
        assert reloaded.core_properties.author == "Athena"
        assert reloaded.custom_properties["Source"] == "cli"

    def it_is_a_noop_when_never_touched(self, empty_prs):
        # if the API is not used, no /docProps/custom.xml is added (the part
        # is created lazily ON first call to .custom_properties_part). A bare
        # save() that never touches the API should leave the package alone.
        buf = BytesIO()
        empty_prs.save(buf)
        # Reopen and confirm no custom_properties_part rel exists yet
        buf.seek(0)
        reloaded = Presentation(buf)
        # accessing custom_properties for the first time HERE creates it,
        # but pre-access there should be no rel of CUSTOM_PROPERTIES type
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        rel_types = {r.reltype for r in reloaded.part.package._rels.values()}
        assert RT.CUSTOM_PROPERTIES not in rel_types
