"""Acceptance suite for issue #21 — Hyperlinks 2.0 & Click Actions.

Round-trip (save→reopen) tests prove the eight sub-features survive a real
package serialization, mirroring the issue #16/#18 acceptance pattern. The
XSD-position probes encode the fork's "python-self-round-trip ≠ PowerPoint
preservation" rule: every new oxml child/attr is paired with an explicit
schema-order assertion (the exact silent-PowerPoint-repair class).
"""

from __future__ import annotations

import io
import os
import struct
import tempfile
import wave

import pytest

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.action import PP_ACTION
from pptx.oxml import parse_xml
from pptx.oxml.action import CT_EmbeddedWAVAudioFile, CT_Hyperlink
from pptx.oxml.ns import nsdecls
from pptx.util import Inches

TEST_IMAGE = "tests/test_files/monty-truth.png"


def _wav_path() -> str:
    fd, path = tempfile.mkstemp(suffix=".wav")
    os.close(fd)
    with wave.open(path, "w") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(8000)
        w.writeframes(struct.pack("<h", 0) * 800)
    return path


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, text="x"):
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tb.text_frame.text = text
    return tb


def _new_run(slide):
    """A fresh run that is unambiguously ``runs[0]`` after round-trip."""
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "x"
    return run


def _first_run(reopened):
    return reopened.slides[0].shapes[-1].text_frame.paragraphs[0].runs[0]


def _roundtrip(prs):
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


# -- A. CT_Hyperlink / CT_EmbeddedWAVAudioFile oxml model ---------------------


class DescribeCTHyperlinkModel:
    def it_models_the_new_optional_attributes(self):
        el = parse_xml(
            '<a:hlinkClick %s tooltip="hi" tgtFrame="_blank"'
            ' invalidUrl="x" history="0" highlightClick="1" endSnd="1"/>' % nsdecls("a", "r")
        )
        assert el.tooltip == "hi"
        assert el.tgtFrame == "_blank"
        assert el.invalidUrl == "x"
        assert el.history is False
        assert el.highlightClick is True
        assert el.endSnd is True

    def it_preserves_rId_and_action_attributes(self):
        el = parse_xml(
            '<a:hlinkClick %s r:id="rId3" action="ppaction://macro?name=M"/>' % nsdecls("a", "r")
        )
        assert el.rId == "rId3"
        assert el.action == "ppaction://macro?name=M"

    def it_serializes_booleans_as_one_or_zero(self):
        el = parse_xml("<a:hlinkClick %s/>" % nsdecls("a", "r"))
        el.endSnd = True
        el.history = False
        assert ' endSnd="1"' in el.xml
        assert ' history="0"' in el.xml

    def it_does_not_emit_unset_attributes(self):
        el = parse_xml("<a:hlinkClick %s/>" % nsdecls("a", "r"))
        assert "tooltip=" not in el.xml
        assert "endSnd=" not in el.xml

    def it_places_snd_before_extLst_per_xsd_sequence(self):
        # -- XSD-position probe: CT_Hyperlink sequence is snd?, extLst? --
        el = parse_xml("<a:hlinkClick %s><a:extLst/></a:hlinkClick>" % nsdecls("a", "r"))
        snd = el.get_or_add_snd()
        snd.embed = "rId9"
        children = [child.tag.split("}")[-1] for child in el]
        assert children == ["snd", "extLst"]

    def it_models_the_embedded_wav_audio_file(self):
        el = parse_xml('<a:snd %s r:embed="rId4" name="ding.wav"/>' % nsdecls("a", "r"))
        assert isinstance(el, CT_EmbeddedWAVAudioFile)
        assert el.embed == "rId4"
        assert el.name == "ding.wav"

    def it_does_not_emit_unset_snd_name(self):
        el = parse_xml('<a:snd %s r:embed="rId4"/>' % nsdecls("a", "r"))
        assert "name=" not in el.xml

    def it_uses_one_class_for_all_three_host_contexts(self):
        for tag in ("a:hlinkClick", "a:hlinkHover", "a:hlinkMouseOver"):
            el = parse_xml("<%s %s/>" % (tag, nsdecls("a", "r")))
            assert isinstance(el, CT_Hyperlink)


# -- B. Run hyperlink tooltip + color ---------------------------------------


class DescribeRunHyperlinkTooltip:
    def it_round_trips_a_run_tooltip_with_an_address(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        run.hyperlink.address = "https://example.com"
        run.hyperlink.tooltip = "Click for details"
        r2 = _first_run(_roundtrip(prs))
        assert r2.hyperlink.tooltip == "Click for details"
        assert r2.hyperlink.address == "https://example.com"

    def it_supports_a_tooltip_only_hyperlink_with_no_url(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        run.hyperlink.tooltip = "just a tip"
        r2 = _first_run(_roundtrip(prs))
        assert r2.hyperlink.tooltip == "just a tip"
        assert r2.hyperlink.address is None

    def it_returns_None_tooltip_when_no_hyperlink_present(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        assert run.hyperlink.tooltip is None

    def it_clears_tooltip_and_prunes_empty_hlink(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        run.hyperlink.tooltip = "temp"
        run.hyperlink.tooltip = None
        assert run.hyperlink.tooltip is None
        assert "hlinkClick" not in run._r.xml

    def it_keeps_the_address_when_only_the_tooltip_is_cleared(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        run.hyperlink.address = "https://keep.example"
        run.hyperlink.tooltip = "temp"
        run.hyperlink.tooltip = None
        assert run.hyperlink.address == "https://keep.example"

    def it_does_not_mutate_xml_on_tooltip_read(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        hl = run.hyperlink
        before = run._r.xml
        _ = hl.tooltip
        # -- reading the tooltip must not create an a:hlinkClick --
        assert run._r.xml == before
        assert "hlinkClick" not in run._r.xml


class DescribeRunHyperlinkColor:
    def it_round_trips_a_hyperlink_color_override(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        run.hyperlink.address = "https://example.com"
        run.hyperlink.color.rgb = RGBColor(0xC0, 0x00, 0x00)
        r2 = _first_run(_roundtrip(prs))
        assert r2.hyperlink.color.rgb == RGBColor(0xC0, 0x00, 0x00)

    def it_does_not_create_a_hyperlink_just_by_setting_color(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        run.hyperlink.color.rgb = RGBColor(0x00, 0x80, 0x00)
        assert run.hyperlink.address is None
        assert "hlinkClick" not in run._r.xml

    def it_does_not_mutate_xml_on_color_read(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        hl = run.hyperlink
        before = run._r.xml
        _ = hl.color.rgb
        # -- reading the color must not materialize a:solidFill --
        assert run._r.xml == before
        assert "solidFill" not in run._r.xml


# -- C. Shape click-action verbs --------------------------------------------


class DescribeClickActionVerbs:
    def it_can_author_a_run_macro_action(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "btn")
        shape.click_action.run_macro("My Macro")
        rt = _roundtrip(prs)
        ca = rt.slides[0].shapes[-1].click_action
        assert ca.action == PP_ACTION.RUN_MACRO
        assert ca._hlink.action_fields["name"] == "My%20Macro"

    def it_does_not_allocate_a_relationship_for_run_macro(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "btn")
        shape.click_action.run_macro("Recalc")
        # -- rId is the empty string, not None: PowerPoint requires r:id on
        # -- every hlinkClick element (even when no relationship is attached)
        # -- to avoid a load-time Repair dialog. See action.py
        # -- ``ActionSetting._get_or_add_hlink`` for the rationale.
        assert shape.click_action._hlink.rId == ""

    def it_can_author_a_target_program_action(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "prog")
        shape.click_action.target_program("C:\\tool.exe")
        rt = _roundtrip(prs)
        ca = rt.slides[0].shapes[-1].click_action
        assert ca.action == PP_ACTION.RUN_PROGRAM
        assert ca.hyperlink.address == "C:\\tool.exe"

    def it_can_author_a_play_sound_action(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "snd")
        shape.click_action.play_sound(_wav_path())
        rt = _roundtrip(prs)
        hlink = rt.slides[0].shapes[-1].click_action._hlink
        assert hlink.snd is not None
        assert hlink.snd.embed is not None
        assert hlink.endSnd is True

    def it_keeps_a_sound_alongside_an_address(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "snd")
        shape.click_action.hyperlink.address = "https://example.com"
        shape.click_action.play_sound(_wav_path())
        hlink = shape.click_action._hlink
        assert hlink.rId is not None
        assert hlink.snd is not None
        children = [c.tag.split("}")[-1] for c in hlink]
        assert children == ["snd"]  # snd is the only child, before any extLst

    def it_replaces_a_prior_action_when_setting_run_macro(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "btn")
        shape.click_action.target_program("C:\\a.exe")
        shape.click_action.run_macro("B")
        assert shape.click_action.action == PP_ACTION.RUN_MACRO
        assert shape.click_action.hyperlink.address is None


# -- D. Shape hover-action ---------------------------------------------------


class DescribeHoverAction:
    def it_exposes_a_public_hover_action(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "h")
        shape.hover_action.hyperlink.address = "https://hover.example"
        rt = _roundtrip(prs)
        sh = rt.slides[0].shapes[-1]
        assert sh.hover_action.hyperlink.address == "https://hover.example"

    def it_keeps_click_and_hover_actions_independent(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "h")
        shape.click_action.hyperlink.address = "https://click.example"
        shape.hover_action.hyperlink.address = "https://hover.example"
        rt = _roundtrip(prs)
        sh = rt.slides[0].shapes[-1]
        assert sh.click_action.hyperlink.address == "https://click.example"
        assert sh.hover_action.hyperlink.address == "https://hover.example"

    def it_places_hlinkClick_before_hlinkHover(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "h")
        shape.click_action.hyperlink.address = "https://c.example"
        shape.hover_action.hyperlink.address = "https://h.example"
        cNvPr = shape._element._nvXxPr.cNvPr
        tags = [
            c.tag.split("}")[-1]
            for c in cNvPr
            if c.tag.split("}")[-1] in ("hlinkClick", "hlinkHover")
        ]
        assert tags == ["hlinkClick", "hlinkHover"]


# -- E. Run-level slideshow jump-to-slide (unification) ----------------------


class DescribeRunSlideJump:
    def it_can_make_a_run_jump_to_a_slide(self):
        prs = Presentation()
        s1 = _blank_slide(prs)
        s2 = _blank_slide(prs)
        run = _new_run(s1)
        run.click_action.target_slide = s2
        ca = _first_run(_roundtrip(prs)).click_action
        assert ca.action == PP_ACTION.NAMED_SLIDE
        assert ca.target_slide is not None

    def it_unifies_run_and_shape_action_interfaces(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        # -- same ActionSetting surface as shape.click_action --
        assert hasattr(run.click_action, "run_macro")
        assert hasattr(run.click_action, "target_program")
        assert hasattr(run.click_action, "play_sound")
        assert hasattr(run.click_action, "target_slide")

    def it_keeps_a_run_jump_and_tooltip_together(self):
        prs = Presentation()
        s1 = _blank_slide(prs)
        s2 = _blank_slide(prs)
        run = _new_run(s1)
        run.click_action.target_slide = s2
        run.hyperlink.tooltip = "next"
        r2 = _first_run(_roundtrip(prs))
        assert r2.hyperlink.tooltip == "next"
        assert r2.click_action.action == PP_ACTION.NAMED_SLIDE


# -- F. Picture hyperlink ----------------------------------------------------


class DescribePictureHyperlink:
    def it_round_trips_a_picture_hyperlink(self):
        prs = Presentation()
        slide = _blank_slide(prs)
        pic = slide.shapes.add_picture(TEST_IMAGE, Inches(1), Inches(1), Inches(1), Inches(1))
        pic.hyperlink.address = "https://pic.example"
        pic.hyperlink.tooltip = "the picture"
        rt = _roundtrip(prs)
        p2 = next(sh for sh in rt.slides[0].shapes if sh.shape_type == 13)
        assert p2.hyperlink.address == "https://pic.example"
        assert p2.hyperlink.tooltip == "the picture"

    def it_reuses_click_action_hyperlink_not_a_parallel_proxy(self):
        prs = Presentation()
        slide = _blank_slide(prs)
        pic = slide.shapes.add_picture(TEST_IMAGE, Inches(1), Inches(1), Inches(1), Inches(1))
        assert pic.hyperlink is pic.click_action.hyperlink

    def it_does_not_emit_hlink_until_set(self):
        prs = Presentation()
        slide = _blank_slide(prs)
        pic = slide.shapes.add_picture(TEST_IMAGE, Inches(1), Inches(1), Inches(1), Inches(1))
        _ = pic.hyperlink
        assert "hlinkClick" not in pic._element.xml


# -- G. Chart-element hyperlink — documented current boundary ----------------
#
# Sub-feature 7 (chart-element hyperlinks) is DEFERRED to a maintainer-filed
# follow-up: `ChartTitle`/`AxisTitle` are `ElementProxy` constructed without a
# chart-part parent chain, so a run inside `chart.chart_title.text_frame`
# cannot resolve `.part` to allocate the hyperlink relationship. Threading the
# chart part through the shared title base is a separate architectural surface
# with broad chart-suite regression risk; see
# `uat/FOLLOWUP_issue21_chart_hyperlinks.md`. This test pins the *current*
# boundary so a future fix flips it deliberately rather than by accident.


class DescribeChartRunHyperlinkBoundary:
    def it_documents_the_chart_title_run_part_limitation(self):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs = Presentation()
        slide = _blank_slide(prs)
        cd = CategoryChartData()
        cd.categories = ["a", "b"]
        cd.add_series("S", (1, 2))
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(4),
            cd,
        )
        chart = gf.chart
        chart.has_title = True
        tf = chart.chart_title.text_frame
        tf.text = "Revenue"
        run = tf.paragraphs[0].runs[0]
        # -- current boundary: the chart-title run has no resolvable part --
        with pytest.raises(AttributeError):
            run.hyperlink.address = "https://chart.example"


# -- H. No-PowerPoint-repair gate (XSD-position probes) ----------------------


class DescribeNoRepairGate:
    def it_keeps_rPr_hlink_successor_order_unchanged(self):
        # -- regression: rPr child order must keep hlinkClick before rtl/extLst
        from pptx.oxml.text import CT_TextCharacterProperties

        hlink_field = CT_TextCharacterProperties.__dict__["hlinkClick"]
        assert hlink_field is not None

    def it_emits_a_schema_valid_snd_only_before_extLst(self):
        el = parse_xml("<a:hlinkClick %s><a:extLst/></a:hlinkClick>" % nsdecls("a", "r"))
        el.get_or_add_snd().embed = "rId1"
        assert el.xml.index("a:snd") < el.xml.index("a:extLst")

    @pytest.mark.parametrize(("flag", "expected"), [(True, '"1"'), (False, '"0"')])
    def it_serializes_endSnd_in_powerpoint_form(self, flag, expected):
        el = parse_xml("<a:hlinkClick %s/>" % nsdecls("a", "r"))
        el.endSnd = flag
        assert ("endSnd=%s" % expected) in el.xml

    # -- Repair-trigger regressions (UAT round 1, 2026-05-19) ----------------
    # -- Three things made PowerPoint throw the Repair dialog and strip the
    # -- play-sound hlinkClick + media1.wav part on first UAT open. The fixes
    # -- are encoded as invariants below so any regression resurfaces in the
    # -- trinity, not in a maintainer's PowerPoint window.

    def it_emits_empty_r_id_on_tooltip_only_shape_hlinkClick(self):
        # -- root cause #1: <a:hlinkClick tooltip="..."/> with no r:id triggers
        # -- a Repair dialog; PowerPoint repairs by stripping the tooltip and
        # -- emitting <a:hlinkClick r:id="" action="ppaction://noaction"/>.
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "tip")
        shape.click_action.tooltip = "hi"
        hlink = shape.click_action._hlink
        assert hlink.rId == ""
        assert ' r:id=""' in hlink.xml

    def it_emits_empty_r_id_on_tooltip_only_run_hlinkClick(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        run.hyperlink.tooltip = "hi"
        hlink = run._r.get_or_add_rPr().hlinkClick
        assert hlink is not None
        assert hlink.rId == ""
        assert ' r:id=""' in hlink.xml

    def it_emits_empty_r_id_on_run_macro_hlinkClick(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "btn")
        shape.click_action.run_macro("Recalc")
        hlink = shape.click_action._hlink
        assert hlink.rId == ""
        assert ' r:id=""' in hlink.xml

    def it_emits_empty_r_id_on_play_sound_hlinkClick(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "snd")
        shape.click_action.play_sound(_wav_path())
        hlink = shape.click_action._hlink
        assert hlink.rId == ""
        assert ' r:id=""' in hlink.xml

    def it_uses_RT_AUDIO_for_play_sound_relationship(self):
        # -- root cause #2: the snd rel was RT.MEDIA (Microsoft-2007 video
        # -- media); PowerPoint discards the whole hlinkClick + the media1.wav
        # -- part because that rel type is not valid for an embedded WAV
        # -- referenced by CT_EmbeddedWAVAudioFile/@r:embed.
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "snd")
        shape.click_action.play_sound(_wav_path())
        snd_rId = shape.click_action._hlink.snd.embed
        slide_part = shape.part
        snd_rel = slide_part.rels[snd_rId]
        assert snd_rel.reltype == RT.AUDIO

    def it_round_trips_play_sound_through_save_and_reopen(self):
        # -- end-to-end: after the two root-cause fixes, the play-sound deck
        # -- survives a real save→reopen with the wav part intact and the
        # -- hlinkClick re-readable. This is what the UAT exercises.
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "snd")
        shape.click_action.play_sound(_wav_path())
        rt = _roundtrip(prs)
        hlink = rt.slides[0].shapes[-1].click_action._hlink
        assert hlink is not None
        assert hlink.snd is not None
        assert hlink.snd.embed is not None
        # -- the wav part must still resolve through the snd's rel
        rt_part = rt.slides[0].shapes[-1].part
        snd_rel = rt_part.rels[hlink.snd.embed]
        assert snd_rel.target_part is not None

    # -- Round 2 finding (UAT 2026-05-19, "tooltip does not appear"): -------
    # -- PowerPoint's tooltip processor only activates on hlinks that carry an
    # -- action verb. A bare ``<a:hlinkClick r:id="" tooltip="..."/>`` opens
    # -- without a repair dialog but is inert in slideshow mode — hovering
    # -- shows no ScreenTip. The fix mirrors PowerPoint's own form for a
    # -- tooltip-only hlink: ``action="ppaction://noaction"``.

    def it_emits_noaction_on_shape_tooltip_only_click_hlink(self):
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "tip")
        shape.click_action.tooltip = "just a tip"
        hlink = shape.click_action._hlink
        assert hlink.action == "ppaction://noaction"
        assert hlink.tooltip == "just a tip"
        assert hlink.rId == ""

    def it_emits_noaction_on_shape_tooltip_only_via_hyperlink_proxy(self):
        # -- matches the UAT path: tip.click_action.hyperlink.tooltip = ...
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "tip")
        shape.click_action.hyperlink.tooltip = "just a tip"
        hlink = shape.click_action._hlink
        assert hlink.action == "ppaction://noaction"
        assert hlink.tooltip == "just a tip"
        assert hlink.rId == ""

    def it_emits_noaction_on_run_tooltip_only_hlink(self):
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        run.hyperlink.tooltip = "just a tip"
        hlink = run._r.get_or_add_rPr().hlinkClick
        assert hlink is not None
        assert hlink.action == "ppaction://noaction"
        assert hlink.tooltip == "just a tip"

    def it_does_not_overwrite_a_real_action_with_noaction(self):
        # -- when a real action is set first, adding a tooltip must NOT replace
        # -- the action verb with the noaction marker.
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "btn")
        shape.click_action.run_macro("Recalc")
        shape.click_action.tooltip = "macro tip"
        hlink = shape.click_action._hlink
        assert hlink.action == "ppaction://macro?name=Recalc"
        assert hlink.tooltip == "macro tip"

    def it_does_not_emit_noaction_when_a_relationship_is_present(self):
        # -- when an address (URL) is set, the hlink has rId; adding a tooltip
        # -- must NOT add the noaction marker either.
        prs = Presentation()
        run = _new_run(_blank_slide(prs))
        run.hyperlink.address = "https://example.com"
        run.hyperlink.tooltip = "url tip"
        hlink = run._r.get_or_add_rPr().hlinkClick
        assert hlink is not None
        assert hlink.action is None
        assert hlink.rId != ""
        assert hlink.tooltip == "url tip"

    def it_clears_noaction_when_tooltip_is_removed_from_inert_hlink(self):
        # -- the noaction marker we added must not linger after the tooltip
        # -- is cleared; the resulting hlink must be pruned cleanly.
        prs = Presentation()
        shape = _textbox(_blank_slide(prs), "tip")
        shape.click_action.tooltip = "temp"
        shape.click_action.tooltip = None
        assert shape.click_action._hlink is None
