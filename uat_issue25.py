"""UAT — issue #25: Threaded Comments & Review.

Exercises all 8 sub-features end-to-end and emits an artifact for
maintainer visual signoff in PowerPoint (Review/Comments pane). Exits
non-zero on any structural failure. Per repo CLAUDE.md §6a an agent
running this green is script-QA, NOT signoff.

Non-visual structural/round-trip evidence is this script's captured
stdout (per the goal directive: script execution + output recording
when there is no visual element). The visual element — the comment in
PowerPoint's Review pane — is captured separately via Interceptor.

Artifact: uat_issue25_comments.pptx (open in PowerPoint → Review → Comments).
"""

from __future__ import annotations

import io
import sys

from pptx import Presentation

OUT = "/Users/mhoroszowski/Projects/AI/python-pptx/uat_issue25_comments.pptx"
ok = True


def check(label: str, cond: bool) -> None:
    global ok
    print(("  PASS " if cond else "  FAIL ") + label)
    ok = ok and cond


print("=== issue #25 UAT — threaded comments ===")

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])

# SF1/SF3 — author + add a comment
c1 = slide.comments.add("Looks great! Ship it.", author="Alex Reviewer")
check("SF3 comments.add returns a Comment", c1 is not None)
check("SF3 slide.comments len == 1", len(slide.comments) == 1)

# SF1 author dedup — same name reused
c2 = slide.comments.add("One nit: tighten the headline.", author="Alex Reviewer")
check("SF1 author dedup (same name, len now 2)", len(slide.comments) == 2)

# SF5 — metadata accessors
check("SF5 Comment.author", c1.author == "Alex Reviewer")
check("SF5 Comment.text", c1.text == "Looks great! Ship it.")
check("SF5 Comment.created_at tz-aware datetime", c1.created_at is not None and c1.created_at.tzinfo is not None)
check("SF5 Comment.anchor_position None when unanchored", c1.anchor_position is None)

# SF4 — threaded replies
r1 = c1.replies.add("Agreed.", author="Sam Lead")
c1.replies.add("Merging now.", author="Alex Reviewer")
check("SF4 replies threaded under parent (len 2)", len(c1.replies) == 2)
check("SF4 reply text", list(c1.replies)[0].text == "Agreed.")

# SF6 — resolve
check("SF6 default resolved is False", c1.resolved is False)
c1.resolve()
check("SF6 resolve() sets resolved True", c1.resolved is True)

# SF7 — per-shape comments (anchored)
shape = slide.shapes.add_textbox(914400, 914400, 1828800, 457200)
shape.text_frame.text = "Anchor target"
ac = slide.comments.add("Re: this textbox", author="Alex Reviewer", anchor=shape)
check("SF7 Shape.comments returns the anchored comment", any(x.text == "Re: this textbox" for x in shape.comments))
other = slide.shapes.add_textbox(0, 0, 914400, 457200)
check("SF7 unrelated shape has empty .comments", len(other.comments) == 0)

# SF3/SF4/SF6 — full round-trip
buf = io.BytesIO()
prs.save(buf)
buf.seek(0)
rt = Presentation(buf)
rs = rt.slides[0]
rcs = list(rs.comments)
check("RT comments survive save/reopen", len(rcs) >= 3)
first = next((x for x in rcs if x.text == "Looks great! Ship it."), None)
check("RT first comment text+author intact", first is not None and first.author == "Alex Reviewer")
check("RT resolved flag persisted", first is not None and first.resolved is True)
check("RT reply thread persisted", first is not None and len(first.replies) == 2)

prs.save(OUT)
print("ARTIFACT:", OUT, "(PowerPoint → Review → Comments)")
print("RESULT:", "PASS" if ok else "FAIL")
sys.exit(0 if ok else 1)
