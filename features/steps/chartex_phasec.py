"""Gherkin step implementations for ChartEx Phase-C features (issue #14):
writers for Treemap/Sunburst/Funnel/BoxWhisker/Histogram/Pareto + replace_data.
"""

from __future__ import annotations

import io
import zipfile

from behave import then, when

from pptx import Presentation
from pptx.chart.data import (
    BoxWhiskerChartData,
    FunnelChartData,
    HistogramChartData,
    ParetoChartData,
    SunburstChartData,
    TreemapChartData,
    WaterfallChartData,
)
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def _data_for(member_name):
    m = member_name.strip()
    if m == "WATERFALL":
        cd = WaterfallChartData()
        cd.categories = ["Q1", "Q2", "Total"]
        cd.add_series("R", [10, 20, 30], subtotals=[2])
        return XL_CHART_TYPE.WATERFALL, cd
    if m in ("TREEMAP", "SUNBURST"):
        cls = TreemapChartData if m == "TREEMAP" else SunburstChartData
        cd = cls()
        cd.add_level(["A", "A", "B", "B"])
        cd.add_level(["a1", "a2", "b1", "b2"])
        cd.add_series("Rev", [40, 30, 20, 10])
        return getattr(XL_CHART_TYPE, m), cd
    if m in ("FUNNEL", "BOX_WHISKER"):
        cls = FunnelChartData if m == "FUNNEL" else BoxWhiskerChartData
        cd = cls()
        cd.categories = ["Leads", "Qualified", "Won"]
        cd.add_series("Pipe", [100, 60, 25])
        return getattr(XL_CHART_TYPE, m), cd
    if m in ("HISTOGRAM", "PARETO"):
        cls = HistogramChartData if m == "HISTOGRAM" else ParetoChartData
        cd = cls()
        cd.add_series("Scores", [55, 62, 71, 73, 88, 91, 64, 78], bin_count=4)
        return getattr(XL_CHART_TYPE, m), cd
    raise KeyError(m)


def _cx_parts(blob):
    z = zipfile.ZipFile(io.BytesIO(blob))
    return [n for n in z.namelist() if "chartEx" in n and n.endswith(".xml")]


# when ====================================================


@when("I add a ChartEx {member_name} chart")
def when_i_add_a_chartex_member_chart(context, member_name):
    ct, cd = _data_for(member_name)
    context.cx_member = member_name.strip()
    context.cx_data = cd
    context.cx_frame = context.slide.shapes.add_chart(
        ct, Inches(1), Inches(1), Inches(6), Inches(4), cd
    )


@when("I replace the ChartEx {member_name} data with a smaller dataset")
def when_i_replace_chartex_data(context, member_name):
    _, new_cd = _data_for(member_name)
    # shrink it so the change is observable
    if hasattr(new_cd, "levels"):
        nd = type(new_cd)()
        nd.add_level(["Z", "Z"])
        nd.add_level(["z1", "z2"])
        nd.add_series("New", [7, 3])
    elif hasattr(new_cd, "categories"):
        nd = type(new_cd)()
        nd.categories = ["Only"]
        nd.add_series("New", [42])
    else:
        nd = type(new_cd)()
        nd.add_series("New", [1, 2, 3, 4], bin_count=2)
    context.cx_replacement = nd
    context.cx_frame.chartex.replace_data(nd)


@when("I attempt to replace a {a_type} ChartEx with {b_type} data")
def when_attempt_mismatch_replace(context, a_type, b_type):
    ct, cd = _data_for(a_type)
    frame = context.slide.shapes.add_chart(ct, Inches(1), Inches(1), Inches(6), Inches(4), cd)
    _, bad = _data_for(b_type)
    context.cx_replace_error = None
    try:
        frame.chartex.replace_data(bad)
    except ValueError as e:
        context.cx_replace_error = e


# then ====================================================


@then("the slide has a ChartEx graphic frame")
def then_slide_has_a_chartex_frame(context):
    frames = [s for s in context.slide.shapes if getattr(s, "has_chartex", False)]
    assert len(frames) >= 1, "no ChartEx graphic frame on slide"


@then("the saved package contains a ChartEx part")
def then_saved_package_contains_chartex_part(context):
    buf = io.BytesIO()
    context.prs.save(buf)
    assert _cx_parts(buf.getvalue()), "no chartEx part in saved package"


@then("the ChartEx round-trips preserving its part")
def then_chartex_round_trips(context):
    buf = io.BytesIO()
    context.prs.save(buf)
    before = sorted(_cx_parts(buf.getvalue()))
    prs2 = Presentation(io.BytesIO(buf.getvalue()))
    prs2.slides.add_slide(prs2.slide_layouts[0])  # unrelated edit (layout 0 always exists)
    buf2 = io.BytesIO()
    prs2.save(buf2)
    after = sorted(_cx_parts(buf2.getvalue()))
    assert before and before == after, f"{before!r} != {after!r}"
    rt = [s for s in prs2.slides[0].shapes if getattr(s, "has_chartex", False)]
    assert len(rt) == 1


@then("the reopened ChartEx reflects the replaced data")
def then_reopened_reflects_replaced(context):
    buf = io.BytesIO()
    context.prs.save(buf)
    prs2 = Presentation(io.BytesIO(buf.getvalue()))
    z = zipfile.ZipFile(io.BytesIO(buf.getvalue()))
    name = next(
        n for n in z.namelist() if "chartEx" in n and n.endswith(".xml") and "_rels" not in n
    )
    xml = z.read(name).decode()
    nd = context.cx_replacement
    token = "New"
    assert token in xml, "replaced series name not found after reopen"
    assert prs2 is not None


@then("a chart-type mismatch error is raised")
def then_mismatch_error_raised(context):
    assert context.cx_replace_error is not None
    assert "cannot change chart type" in str(context.cx_replace_error)
