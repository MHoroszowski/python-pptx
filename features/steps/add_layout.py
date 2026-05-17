"""Gherkin step implementations for issue #19 SF3 — SlideLayouts.add_layout()."""

from __future__ import annotations

import io

from behave import given, then, when

from pptx import Presentation

# given ===================================================


@given("a default presentation")
def given_a_default_presentation(context):
    context.prs = Presentation()
    master = context.prs.slide_masters[0]
    context.slide_layouts = master.slide_layouts
    context.layout_count_before = len(context.slide_layouts)
    context.slide_count_before = len(context.prs.slides)


# when ====================================================


@when("I call slide_layouts.add_layout()")
def when_I_call_add_layout_no_name(context):
    context.new_layout = context.slide_layouts.add_layout()


@when('I call slide_layouts.add_layout("{name}")')
def when_I_call_add_layout_with_name(context, name):
    context.new_layout = context.slide_layouts.add_layout(name)


@when("I save and reopen the presentation")
def when_I_save_and_reopen(context):
    buf = io.BytesIO()
    context.prs.save(buf)
    buf.seek(0)
    context.reopened = Presentation(buf)


@when("I add a slide based on the new layout")
def when_I_add_a_slide_based_on_the_new_layout(context):
    context.prs.slides.add_slide(context.new_layout)


# then ====================================================


@then("the slide-master layout count increased by exactly 1")
def then_layout_count_increased_by_1(context):
    assert len(context.slide_layouts) == context.layout_count_before + 1


@then("the new layout has a non-empty name")
def then_new_layout_has_non_empty_name(context):
    assert context.new_layout.name not in (None, "")


@then('slide_layouts.get_by_name("{name}") is the new layout')
def then_get_by_name_returns_new_layout(context, name):
    assert context.slide_layouts.get_by_name(name) is not None
    assert context.slide_layouts.get_by_name(name).name == context.new_layout.name


@then('the reopened presentation has a layout named "{name}"')
def then_reopened_has_layout_named(context, name):
    layouts = context.reopened.slide_masters[0].slide_layouts
    assert layouts.get_by_name(name) is not None


@then("the slide count increased by exactly 1")
def then_slide_count_increased_by_1(context):
    assert len(context.prs.slides) == context.slide_count_before + 1


# SF2 — save_as_potx ======================================

_TEMPLATE_CT = b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
_PRESENTATION_CT = (
    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)


@when("I save the presentation as a potx")
def when_I_save_as_potx(context):
    context.ct_before = context.prs.part.content_type
    context.potx_buf = io.BytesIO()
    context.prs.save_as_potx(context.potx_buf)


@then("the saved potx declares the template content-type")
def then_potx_declares_template_ct(context):
    import zipfile

    context.potx_buf.seek(0)
    with zipfile.ZipFile(context.potx_buf) as z:
        ct_xml = z.read("[Content_Types].xml")
    assert _TEMPLATE_CT in ct_xml
    assert _PRESENTATION_CT not in ct_xml


@then("the in-memory presentation content-type is unchanged")
def then_in_memory_ct_unchanged(context):
    assert context.prs.part.content_type == context.ct_before
    assert context.prs.part.content_type == _PRESENTATION_CT.decode("ascii")


@then("the saved potx reopens as a valid presentation")
def then_potx_reopens_valid(context):
    context.potx_buf.seek(0)
    reopened = Presentation(context.potx_buf)
    assert len(reopened.slide_masters) >= 1


# SF5 — master shape authoring ============================


@when("I add a textbox to the slide master")
def when_I_add_textbox_to_master(context):
    master = context.prs.slide_masters[0]
    tb = master.shapes.add_textbox(0, 0, 914400, 457200)
    tb.text_frame.text = "ACCEPTANCE MASTER TEXT"


@then("the reopened slide master has the master textbox text")
def then_reopened_master_has_textbox_text(context):
    master = context.reopened.slide_masters[0]
    texts = [s.text_frame.text for s in master.shapes if s.has_text_frame]
    assert "ACCEPTANCE MASTER TEXT" in texts


# SF4 — copy_from =========================================


@when("I add a textbox to the new layout")
def when_I_add_textbox_to_new_layout(context):
    context.new_layout.shapes.add_textbox(0, 0, 914400, 457200)


@when("I copy the new layout with copy_from")
def when_I_copy_layout_with_copy_from(context):
    context.source_shape_count = len(list(context.new_layout.shapes))
    context.copied_layout = context.slide_layouts.copy_from(context.new_layout)


@then("the copied layout has the same shape count as its source")
def then_copied_layout_same_shape_count(context):
    assert len(list(context.copied_layout.shapes)) == context.source_shape_count


@then("the source layout is unchanged after copy_from")
def then_source_layout_unchanged(context):
    assert len(list(context.new_layout.shapes)) == context.source_shape_count


# SF7 — cross-master / apply layout =======================


@when("I add a slide on the default layout")
def when_I_add_slide_on_default_layout(context):
    context.sf7_slide = context.prs.slides.add_slide(context.prs.slide_layouts[0])


@when("I apply the new layout to that slide")
def when_I_apply_new_layout_to_slide(context):
    context.sf7_slide.slide_layout = context.new_layout


@then('the reopened slide uses the layout named "{name}"')
def then_reopened_slide_uses_layout(context, name):
    assert context.reopened.slides[0].slide_layout.name == name


@then("the reopened slide still resolves its slide master")
def then_reopened_slide_resolves_master(context):
    assert context.reopened.slides[0].slide_layout.slide_master is not None


# SF8 — chart into placeholder ============================


@when("I add a layout with a chart placeholder")
def when_I_add_layout_with_chart_placeholder(context):
    from pptx.enum.shapes import PP_PLACEHOLDER

    master = context.prs.slide_masters[0]
    context.chart_layout = master.slide_layouts.add_layout("Chart PH Layout")
    context.chart_layout.placeholders.add(
        10,
        PP_PLACEHOLDER.CHART,
        left=914400,
        top=914400,
        width=4572000,
        height=2743200,
    )


@when("I add a slide on that chart-placeholder layout")
def when_I_add_slide_on_chart_layout(context):
    context.chart_slide = context.prs.slides.add_slide(context.chart_layout)


@when("I insert a chart into the slide's chart placeholder")
def when_I_insert_chart_into_placeholder(context):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import PP_PLACEHOLDER

    chart_ph = next(
        p
        for p in context.chart_slide.placeholders
        if p.placeholder_format.type == PP_PLACEHOLDER.CHART
    )
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("S1", (1.0, 2.0, 3.0))
    chart_ph.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)


@then("the reopened slide has exactly one chart")
def then_reopened_slide_has_one_chart(context):
    slide = context.reopened.slides[0]
    charts = [s for s in slide.shapes if s.has_chart]
    assert len(charts) == 1
