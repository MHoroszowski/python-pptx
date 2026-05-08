"""Gherkin step implementations for Table style API (issue #12 Phase 2)."""

from __future__ import annotations

import io

import pytest
from behave import then, when

from pptx import Presentation


# when ====================================================


@when('I call table.apply_style("{style_name_or_guid}")')
def when_apply_style(context, style_name_or_guid):
    context.table_.apply_style(style_name_or_guid)


@when("I set table.style_id to None")
def when_set_style_id_to_none(context):
    context.table_.style_id = None


@when("I save and reload the presentation via stream")
def when_save_and_reload_via_stream(context):
    buf = io.BytesIO()
    context.prs.save(buf)
    buf.seek(0)
    context.prs_reloaded = Presentation(buf)
    context.table_reloaded = next(
        shp for shp in context.prs_reloaded.slides[0].shapes if shp.has_table
    ).table


# then ====================================================


@then('table.style_id is "{guid}"')
def then_style_id_is(context, guid):
    assert context.table_.style_id == guid, (context.table_.style_id, guid)


@then("table.style_id is None")
def then_style_id_is_none(context):
    assert context.table_.style_id is None, context.table_.style_id


@then('table.style_name is "{name}"')
def then_style_name_is(context, name):
    assert context.table_.style_name == name, (context.table_.style_name, name)


@then("table.style_name is None")
def then_style_name_is_none(context):
    assert context.table_.style_name is None, context.table_.style_name


@then('calling table.apply_style("{name}") raises ValueError')
def then_apply_style_raises_ValueError(context, name):
    with pytest.raises(ValueError):
        context.table_.apply_style(name)


@then('the reloaded table has style_id "{guid}"')
def then_reloaded_style_id_is(context, guid):
    assert context.table_reloaded.style_id == guid


@then('the reloaded table has style_name "{name}"')
def then_reloaded_style_name_is(context, name):
    assert context.table_reloaded.style_name == name
