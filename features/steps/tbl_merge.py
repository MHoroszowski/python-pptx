"""Gherkin step implementations for Table merge robustness (issue #12 Phase 3)."""

from __future__ import annotations

import pytest
from behave import given, then, when

from pptx import Presentation
from pptx.util import Inches


# given ===================================================


@given("a 3x3 table on a fresh slide")
def given_a_3x3_table(context):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(6), Inches(2))
    context.prs = prs
    context.table_ = shape.table


# when ====================================================


@when("I call table.merge_cells row=({r1:d},{r2:d}) col=({c1:d},{c2:d})")
def when_merge_cells(context, r1, r2, c1, c2):
    context.table_.merge_cells((r1, r2), (c1, c2))


@when(
    "I call table.merge_cells with range({r_start:d},{r_stop:d}) and range({c_start:d},{c_stop:d})"
)
def when_merge_cells_with_range(context, r_start, r_stop, c_start, c_stop):
    context.table_.merge_cells(range(r_start, r_stop), range(c_start, c_stop))


@when("I call table.split_cells row=({r1:d},{r2:d}) col=({c1:d},{c2:d})")
def when_split_cells(context, r1, r2, c1, c2):
    context.table_.split_cells((r1, r2), (c1, c2))


# then ====================================================


@then("cell ({r:d},{c:d}) has gridSpan={gs:d} and rowSpan={rs:d}")
def then_cell_has_dimensions(context, r, c, gs, rs):
    cell = context.table_.cell(r, c)
    assert cell.grid_span == gs, (cell.grid_span, gs)
    assert cell.row_span == rs, (cell.row_span, rs)


@then("cell ({r:d},{c:d}) is_merge_origin is {expected:S}")
def then_cell_is_merge_origin(context, r, c, expected):
    actual = context.table_.cell(r, c).is_merge_origin
    want = expected == "True"
    assert actual is want, (actual, expected)


@then("cell ({r:d},{c:d}) hMerge is {expected:S}")
def then_cell_hMerge(context, r, c, expected):
    actual = context.table_.cell(r, c).h_merge
    want = expected == "True"
    assert actual is want, (actual, expected)


@then("cell ({r:d},{c:d}) vMerge is {expected:S}")
def then_cell_vMerge(context, r, c, expected):
    actual = context.table_.cell(r, c).v_merge
    want = expected == "True"
    assert actual is want, (actual, expected)


@then("calling table.merge_cells row=({r1:d},{r2:d}) col=({c1:d},{c2:d}) raises ValueError")
def then_merge_cells_raises(context, r1, r2, c1, c2):
    with pytest.raises(ValueError):
        context.table_.merge_cells((r1, r2), (c1, c2))


@then("calling table.split_cells row=({r1:d},{r2:d}) col=({c1:d},{c2:d}) raises ValueError")
def then_split_cells_raises(context, r1, r2, c1, c2):
    with pytest.raises(ValueError):
        context.table_.split_cells((r1, r2), (c1, c2))
