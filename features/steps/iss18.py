"""Step implementations for features/iss-18-shape-effects.feature (issue #18).

Self-contained: every scenario builds an in-memory blank presentation, so no
fixture .pptx files are needed.
"""

import io

from behave import given, then, when

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_END_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt


def _blank():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


@given("a blank slide with one rectangle")
def given_blank_slide_one_rectangle(context):
    context.prs, slide = _blank()
    context.shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
    )
    context.slide = slide


@given("a blank slide with one connector")
def given_blank_slide_one_connector(context):
    context.prs, slide = _blank()
    context.shape = slide.shapes.add_connector(2, Inches(1), Inches(1), Inches(4), Inches(1))
    context.slide = slide


@given("a group scaled two-to-one containing one rectangle")
def given_group_scaled_two_to_one(context):
    context.prs, slide = _blank()
    group = slide.shapes.add_group_shape()
    context.shape = group.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(1), Inches(1)
    )
    x = group._element.grpSpPr.get_or_add_xfrm()
    x.get_or_add_off().x = Emu(0)
    x.get_or_add_off().y = Emu(0)
    x.get_or_add_ext().cx = Emu(Inches(4))
    x.get_or_add_ext().cy = Emu(Inches(4))
    x.get_or_add_chOff().x = Emu(0)
    x.get_or_add_chOff().y = Emu(0)
    x.get_or_add_chExt().cx = Emu(Inches(2))
    x.get_or_add_chExt().cy = Emu(Inches(2))


@when("I set the glow color to FF0000 and radius to 20pt")
def when_set_glow(context):
    context.shape.shadow.glow_effect.color.rgb = RGBColor(0xFF, 0, 0)
    context.shape.shadow.glow_effect.radius = Pt(20)


@then("the rectangle reports glow radius 20pt and color FF0000")
def then_glow(context):
    assert context.shape.shadow.glow_effect.radius == Emu(Pt(20))
    assert context.shape.shadow.glow_effect.color.rgb == RGBColor(0xFF, 0, 0)


@when("I set the reflection blur radius to 3pt and distance to 7pt")
def when_set_reflection(context):
    context.shape.shadow.reflection_effect.blur_radius = Pt(3)
    context.shape.shadow.reflection_effect.distance = Pt(7)


@then("the rectangle reports reflection blur radius 3pt")
def then_reflection(context):
    assert context.shape.shadow.reflection_effect.blur_radius == Emu(Pt(3))


@when("I set the soft edge radius to 5pt")
def when_set_soft_edge(context):
    context.shape.shadow.soft_edge_effect.radius = Pt(5)


@then("the rectangle reports soft edge radius 5pt")
def then_soft_edge(context):
    assert context.shape.shadow.soft_edge_effect.radius == Emu(Pt(5))


@when("I set the 3-D camera preset to orthographicFront")
def when_set_camera(context):
    context.shape.scene_3d.camera_preset = "orthographicFront"


@then("the rectangle reports camera preset orthographicFront")
def then_camera(context):
    assert context.shape.scene_3d.camera_preset == "orthographicFront"


@then("the scene has a light rig")
def then_lightrig(context):
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    s3d = context.shape._element.spPr.find(f"{a}scene3d")
    assert s3d.find(f"{a}lightRig") is not None


@when("I set the extrusion height to 18pt")
def when_set_extrusion(context):
    context.shape.shape_3d.extrusion_height = Pt(18)


@then("the rectangle reports extrusion height 18pt")
def then_extrusion(context):
    assert context.shape.shape_3d.extrusion_height == Emu(Pt(18))


@when("I flip the rectangle vertically and round-trip the file")
def when_flip_v_roundtrip(context):
    context.shape.flip_vertical = True
    buf = io.BytesIO()
    context.prs.save(buf)
    buf.seek(0)
    context.prs2 = Presentation(buf)


@then("the reopened rectangle is flipped vertically")
def then_flipped_v(context):
    shp = [
        s
        for s in context.prs2.slides[0].shapes
        if s.shape_type is not None and "AUTO_SHAPE" in str(s.shape_type)
    ][0]
    assert shp.flip_vertical is True


@when("I flip the rectangle horizontally")
def when_flip_h(context):
    context.shape.flip_horizontal = True


@then("the rectangle is flipped horizontally")
def then_flipped_h(context):
    assert context.shape.flip_horizontal is True


@when("I duplicate the rectangle")
def when_duplicate(context):
    context.dup = context.shape.duplicate()


@then("the slide has two rectangles with distinct shape ids")
def then_two_rectangles(context):
    ids = [s.shape_id for s in context.slide.shapes]
    assert len(ids) == 2
    assert len(set(ids)) == 2


@when("I set the tail end arrowhead to a triangle and round-trip the file")
def when_tail_triangle_roundtrip(context):
    context.shape.line.tail_end.type = MSO_LINE_END_TYPE.TRIANGLE
    buf = io.BytesIO()
    context.prs.save(buf)
    buf.seek(0)
    context.prs2 = Presentation(buf)


@then("the reopened connector tail end is a triangle")
def then_tail_triangle(context):
    conn = list(context.prs2.slides[0].shapes)[0]
    assert conn.line.tail_end.type == MSO_LINE_END_TYPE.TRIANGLE


@then("the rectangle slide_width is double its local width")
def then_world_width_double(context):
    assert context.shape.slide_width == Emu(context.shape.width * 2)
