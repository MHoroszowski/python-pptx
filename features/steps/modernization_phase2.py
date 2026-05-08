"""Gherkin steps for Modernization Phase 2 (issue #29)."""

from __future__ import annotations

import datetime as dt
import io

import pytest
from behave import given, then, when
from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor


# given ===================================================


@given("a fresh slide with a title placeholder")
def given_a_fresh_slide_with_title(context):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    context.prs = prs
    context.slide = slide


@given("a fresh presentation for core-property datetimes")
def given_a_fresh_presentation_for_core_props(context):
    context.prs = Presentation()


# when ====================================================


@when("I read run.font.color.rgb on an unstyled run")
def when_read_font_color_unstyled(context):
    tf = context.slide.shapes.title.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "x"
    rPr = run._r.get_or_add_rPr()
    context.run = run
    context.rPr = rPr
    context.before_xml = etree.tostring(rPr)
    _ = run.font.color.rgb


@when("I set run.font.color.rgb to RGBColor(0xFF, 0x00, 0x00)")
def when_set_font_color_red(context):
    tf = context.slide.shapes.title.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "x"
    context.run = run
    context.rPr = run._r.get_or_add_rPr()
    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)


@when("I set core_properties.created to a tz-aware UTC datetime")
def when_set_created_utc(context):
    context.target_dt = dt.datetime(2024, 7, 4, 17, 0, 0, tzinfo=dt.timezone.utc)
    context.prs.core_properties.created = context.target_dt


# the "save and reload via stream" step is shared with tbl_styles.py


# then ====================================================


@then("the underlying rPr XML is unchanged from before the access")
def then_rPr_unchanged(context):
    after = etree.tostring(context.rPr)
    assert context.before_xml == after, (context.before_xml, after)


@then("the underlying rPr now contains an a:solidFill child")
def then_rPr_has_solidFill(context):
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    assert context.rPr.find("%ssolidFill" % ns) is not None


@then("run.font.color.rgb reads back as FF0000")
def then_font_color_rgb_reads_FF0000(context):
    assert context.run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)


@then("the reloaded core_properties.created is tz-aware")
def then_reloaded_created_is_tzaware(context):
    buf = io.BytesIO()
    context.prs.save(buf)
    buf.seek(0)
    prs2 = Presentation(buf)
    reloaded = prs2.core_properties.created
    assert reloaded.tzinfo is not None
    assert reloaded == context.target_dt


@then('shapes.by_name("Title 1") returns the title shape')
def then_by_name_returns_title(context):
    sh = context.slide.shapes.by_name("Title 1")
    assert sh.name == "Title 1"


@then('shapes.by_name("Bogus") raises KeyError')
def then_by_name_raises_keyerror(context):
    with pytest.raises(KeyError):
        context.slide.shapes.by_name("Bogus")
