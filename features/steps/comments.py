"""Gherkin step implementations for modern threaded comments (issue #25)."""

from __future__ import annotations

import io

from behave import given, then, when

from pptx import Presentation


# given ===================================================


@given("a blank slide with no comments")
def given_blank_slide_no_comments(context):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    context.prs = prs
    context.slide = slide
    assert len(slide.comments) == 0


# when ====================================================


@when('I add a comment "{text}" by author "{author}"')
def when_add_comment(context, text, author):
    context.comment = context.slide.comments.add(text, author)


@when('I reply "{text}" by author "{author}"')
def when_add_reply(context, text, author):
    context.comment.replies.add(text, author)


@when("I remove that comment")
def when_remove_comment(context):
    context.slide.comments.remove(context.comment)


# then ====================================================


@then("the slide has {count:d} comment")
@then("the slide has {count:d} comments")
def then_slide_has_n_comments(context, count):
    assert len(context.slide.comments) == count, "expected %d comments, got %d" % (
        count,
        len(context.slide.comments),
    )


@then('after save and reopen the first comment text is "{text}" by author "{author}"')
def then_reopen_first_comment(context, text, author):
    stream = io.BytesIO()
    context.prs.save(stream)
    stream.seek(0)
    comment = list(Presentation(stream).slides[0].comments)[0]
    assert comment.text == text, "text was %r" % comment.text
    assert comment.author == author, "author was %r" % comment.author


@then('after save and reopen the first comment has {count:d} reply with text "{text}"')
def then_reopen_reply(context, count, text):
    stream = io.BytesIO()
    context.prs.save(stream)
    stream.seek(0)
    comment = list(Presentation(stream).slides[0].comments)[0]
    replies = list(comment.replies)
    assert len(replies) == count, "expected %d replies, got %d" % (count, len(replies))
    assert replies[0].text == text, "reply text was %r" % replies[0].text


@then("after save and reopen the slide has {count:d} comments")
def then_reopen_comment_count(context, count):
    stream = io.BytesIO()
    context.prs.save(stream)
    stream.seek(0)
    reopened = Presentation(stream).slides[0]
    assert len(reopened.comments) == count, "expected %d comments after reopen, got %d" % (
        count,
        len(reopened.comments),
    )


# -- Wave 3: SF6 resolve + SF8 legacy/modern coexistence ------------------------


@when("I resolve that comment")
def when_resolve_comment(context):
    context.comment.resolve()


@then("after save and reopen the first comment is resolved")
def then_reopen_first_comment_resolved(context):
    stream = io.BytesIO()
    context.prs.save(stream)
    stream.seek(0)
    comment = list(Presentation(stream).slides[0].comments)[0]
    assert comment.resolved is True, "comment.resolved was %r" % comment.resolved


def _build_legacy_plus_modern(prs0):
    """Inject a legacy <p:cm> + author + slide rel into a saved package and
    return the reopened Presentation (modern not yet added)."""
    import zipfile

    base = io.BytesIO()
    prs0.save(base)
    base.seek(0)

    legacy_authors = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<p:cmAuthorLst xmlns:p="http://schemas.openxmlformats.org/'
        'presentationml/2006/main">'
        '<p:cmAuthor id="0" name="Legacy Larry" initials="LL" '
        'lastIdx="1" clrIdx="0"/></p:cmAuthorLst>'
    )
    legacy_comments = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<p:cmLst xmlns:p="http://schemas.openxmlformats.org/'
        'presentationml/2006/main">'
        '<p:cm authorId="0" dt="2026-05-16T00:00:00" idx="1">'
        '<p:pos x="100" y="200"/><p:text>Legacy feedback</p:text>'
        "</p:cm></p:cmLst>"
    )
    src = zipfile.ZipFile(base, "r")
    names = src.namelist()
    ct = (
        src.read("[Content_Types].xml")
        .decode("utf-8")
        .replace(
            "</Types>",
            '<Override PartName="/ppt/commentAuthors.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.'
            'presentationml.commentAuthors+xml"/>'
            '<Override PartName="/ppt/comments/comment1.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.'
            'presentationml.comments+xml"/></Types>',
        )
    )
    pres_rels = (
        src.read("ppt/_rels/presentation.xml.rels")
        .decode("utf-8")
        .replace(
            "</Relationships>",
            '<Relationship Id="rIdLegacyAuth" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
            'relationships/commentAuthors" Target="commentAuthors.xml"/>'
            "</Relationships>",
        )
    )
    slide_rels = (
        src.read("ppt/slides/_rels/slide1.xml.rels")
        .decode("utf-8")
        .replace(
            "</Relationships>",
            '<Relationship Id="rIdLegacyCmt" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
            'relationships/comments" Target="../comments/comment1.xml"/>'
            "</Relationships>",
        )
    )
    patched = {
        "[Content_Types].xml": ct,
        "ppt/_rels/presentation.xml.rels": pres_rels,
        "ppt/slides/_rels/slide1.xml.rels": slide_rels,
        "ppt/commentAuthors.xml": legacy_authors,
        "ppt/comments/comment1.xml": legacy_comments,
    }
    out = io.BytesIO()
    dst = zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED)
    for name in names:
        dst.writestr(name, patched.get(name, src.read(name)))
    for name, data in patched.items():
        if name not in names:
            dst.writestr(name, data)
    src.close()
    dst.close()
    out.seek(0)
    return Presentation(out)


@given("a slide carrying a pre-existing legacy comment")
def given_legacy_comment_slide(context):
    prs0 = Presentation()
    prs0.slides.add_slide(prs0.slide_layouts[6])
    context.prs = _build_legacy_plus_modern(prs0)
    context.slide = context.prs.slides[0]
    assert any(c.is_legacy for c in context.slide.comments)


@when('I add a modern comment "{text}" by author "{author}"')
def when_add_modern_comment(context, text, author):
    context.comment = context.slide.comments.add(text, author)


@then("after save and reopen both the legacy and modern comments are present")
def then_reopen_both_families(context):
    stream = io.BytesIO()
    context.prs.save(stream)
    stream.seek(0)
    comments = list(Presentation(stream).slides[0].comments)
    texts = sorted(c.text for c in comments)
    assert texts == ["Legacy feedback", "Modern note"], "texts were %r" % texts


@then("after save and reopen the legacy author ids are intact")
def then_reopen_legacy_author_ids(context):
    import zipfile

    stream = io.BytesIO()
    context.prs.save(stream)
    stream.seek(0)
    zf = zipfile.ZipFile(stream, "r")
    legacy_xml = zf.read("ppt/comments/comment1.xml").decode("utf-8")
    authors_xml = zf.read("ppt/commentAuthors.xml").decode("utf-8")
    zf.close()
    assert 'authorId="0"' in legacy_xml, "legacy authorId mangled: %r" % legacy_xml
    assert 'name="Legacy Larry"' in authors_xml, "legacy author lost: %r" % authors_xml
