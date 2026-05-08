"""Gherkin step implementations for core properties-related features."""

from __future__ import annotations

from datetime import datetime, timedelta, timezone

from behave import given, then, when
from helpers import no_core_props_pptx_path, saved_pptx_path

from pptx import Presentation

# given ===================================================


@given("I have a reference to the core properties of a presentation")
def step_given_ref_to_core_doc_props(context):
    context.prs = Presentation()
    context.core_properties = context.prs.core_properties


# when ====================================================


@when("I open a presentation having no core properties part")
def step_when_open_presentation_with_no_core_props_part(context):
    context.prs = Presentation(no_core_props_pptx_path)


@when("I set the core properties to valid values")
def step_when_set_core_doc_props_to_valid_values(context):
    # ---issue #29 Phase 2: datetimes round-trip as tz-aware UTC values, so
    # ---the input fixture is already tz-aware UTC for an apples-to-apples
    # ---comparison with the reloaded value.
    context.propvals = (
        ("author", "Creator"),
        ("category", "Category"),
        ("comments", "Description"),
        ("content_status", "Content Status"),
        ("created", datetime(2013, 6, 15, 12, 34, 56, tzinfo=timezone.utc)),
        ("identifier", "Identifier"),
        ("keywords", "key; word; keyword"),
        ("language", "Language"),
        ("last_modified_by", "Last Modified By"),
        ("last_printed", datetime(2013, 6, 15, 12, 34, 56, tzinfo=timezone.utc)),
        ("modified", datetime(2013, 6, 15, 12, 34, 56, tzinfo=timezone.utc)),
        ("revision", 9),
        ("subject", "Subject"),
        # --- exercise unicode-text case for Python 2.7 ---
        ("title", "åß∂Title°"),
        ("version", "Version"),
    )
    for name, value in context.propvals:
        setattr(context.prs.core_properties, name, value)


# then ====================================================


@then("a core properties part with default values is added")
def step_then_a_core_props_part_with_def_vals_is_added(context):
    core_props = context.prs.core_properties
    assert core_props.title == "PowerPoint Presentation"
    assert core_props.last_modified_by == "python-pptx"
    assert core_props.revision == 1
    # core_props.modified only stores time with seconds resolution, so
    # comparison needs to be a little loose (within two seconds). Issue #29
    # Phase 2 makes the parser return a tz-aware datetime; use a tz-aware
    # `now()` here so the subtraction is well-typed.
    modified_timedelta = datetime.now(timezone.utc) - core_props.modified
    max_expected_timedelta = timedelta(seconds=2)
    assert modified_timedelta < max_expected_timedelta


@then("the core properties of the presentation have the values I set")
def step_then_core_props_have_values_previously_set(context):
    core_props = Presentation(saved_pptx_path).core_properties
    for name, value in context.propvals:
        assert getattr(core_props, name) == value, "for core property '%s'" % name
