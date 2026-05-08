"""Gherkin step implementations for Table sizing & ergonomics (issue #12 Phase 4)."""

from __future__ import annotations

import io

from behave import given, then, when

from pptx import Presentation
from pptx.util import Emu, Inches


# given ===================================================


@given("a 3x4 table on a fresh slide")
def given_a_3x4_table(context):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(3, 4, Inches(1), Inches(1), Inches(6), Inches(2))
    context.prs = prs
    context.table_ = shape.table


# when ====================================================


@when("I add a row to the table")
def when_add_row(context):
    context.table_.rows.add()


@when("I remove column {idx:d} from the table")
def when_remove_column(context, idx):
    context.table_.columns.remove(idx)


@when("I set row {idx:d} height to {emu:d} EMU")
def when_set_row_height(context, idx, emu):
    context.table_.rows[idx].height = Emu(emu)


@when("I set column {idx:d} width to {emu:d} EMU")
def when_set_column_width(context, idx, emu):
    context.table_.columns[idx].width = Emu(emu)


# the "save and reload via stream" step is shared with tbl_styles.py — reuse


# then ====================================================


@then("table.row_count is {n:d}")
def then_row_count(context, n):
    assert context.table_.row_count == n, (context.table_.row_count, n)


@then("table.column_count is {n:d}")
def then_column_count(context, n):
    assert context.table_.column_count == n, (context.table_.column_count, n)


@then("table.dimensions is ({rows:d}, {cols:d})")
def then_dimensions(context, rows, cols):
    assert context.table_.dimensions == (rows, cols), (context.table_.dimensions, rows, cols)


@then("the reloaded row {idx:d} height is {emu:d}")
def then_reloaded_row_height(context, idx, emu):
    actual = context.table_reloaded.rows[idx].height
    assert actual == emu, (actual, emu)


@then("the reloaded column {idx:d} width is {emu:d}")
def then_reloaded_column_width(context, idx, emu):
    actual = context.table_reloaded.columns[idx].width
    assert actual == emu, (actual, emu)
