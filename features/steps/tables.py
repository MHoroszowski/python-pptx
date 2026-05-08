"""Gherkin step implementations for Table row/column CRUD (issue #12 Phase 1)."""

from __future__ import annotations

import pytest
from behave import given, then, when

from pptx import Presentation
from pptx.util import Inches


# given ===================================================


def _seed_table(context, rows: int, cols: int):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(6), Inches(2))
    context.prs = prs
    context.table_ = shape.table


@given("a 2x3 table on a fresh slide")
def given_a_2x3_table(context):
    _seed_table(context, 2, 3)


@given("a 3x2 table on a fresh slide")
def given_a_3x2_table(context):
    _seed_table(context, 3, 2)


@given("a 2x2 table on a fresh slide")
def given_a_2x2_table(context):
    _seed_table(context, 2, 2)


@given("a 3x2 table on a fresh slide with a vertical merge between rows 0 and 1")
def given_3x2_with_vmerge(context):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    table = shape.table
    table.cell(0, 0).merge(table.cell(1, 0))
    context.prs = prs
    context.table_ = table


@given("a 2x3 table on a fresh slide with a horizontal merge between columns 0 and 1")
def given_2x3_with_hmerge(context):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(6), Inches(2))
    table = shape.table
    table.cell(0, 0).merge(table.cell(0, 1))
    context.prs = prs
    context.table_ = table


# when ====================================================


@when("I call table.rows.add()")
def when_table_rows_add(context):
    context.new_row = context.table_.rows.add()


@when("I call table.rows.add(at={at:d})")
def when_table_rows_add_at(context, at):
    context.new_row = context.table_.rows.add(at=at)


@when("I call table.columns.add()")
def when_table_columns_add(context):
    context.new_column = context.table_.columns.add()


@when("I call table.rows.remove({idx:d})")
def when_table_rows_remove(context, idx):
    context.table_.rows.remove(idx)


@when("I call table.columns.remove({idx:d})")
def when_table_columns_remove(context, idx):
    context.table_.columns.remove(idx)


# then ====================================================


@then("the table has {n:d} rows")
def then_table_has_n_rows(context, n):
    assert len(context.table_.rows) == n, f"expected {n} rows, got {len(context.table_.rows)}"


@then("the table has {n:d} columns")
def then_table_has_n_columns(context, n):
    assert len(context.table_.columns) == n, (
        f"expected {n} columns, got {len(context.table_.columns)}"
    )


@then("the new row is at index {idx:d}")
def then_new_row_is_at_index(context, idx):
    assert context.table_.rows[idx]._tr is context.new_row._tr


@then("every row has {n:d} cells")
def then_every_row_has_n_cells(context, n):
    for row in context.table_.rows:
        assert len(row.cells) == n


@then("calling table.rows.remove({idx:d}) raises ValueError")
def then_rows_remove_raises_value_error(context, idx):
    with pytest.raises(ValueError):
        context.table_.rows.remove(idx)


@then("calling table.columns.remove({idx:d}) raises ValueError")
def then_columns_remove_raises_value_error(context, idx):
    with pytest.raises(ValueError):
        context.table_.columns.remove(idx)
