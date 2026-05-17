"""Gherkin step implementations for ChartEx features."""

from __future__ import annotations

import os
import zipfile

from behave import given, then, when
from helpers import saved_pptx_path

from pptx import Presentation
from pptx.chart.data import CategoryChartData, WaterfallChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml.ns import qn
from pptx.util import Inches

# given ===================================================


@given("ChartEx waterfall data case {data_case}")
def given_chartex_waterfall_data_case_data_case(context, data_case):
    chart_data = WaterfallChartData()
    series_name, categories, values, subtotals = {
        "q4-total": (
            "Revenue",
            ["Q1", "Q2", "Q3", "Q4", "Total"],
            [100, 50, -30, 80, 200],
            [4],
        ),
        "cash-bridge": (
            "Cash Flow",
            ["Start", "Sales", "Returns", "Ops", "Tax", "End"],
            [500, 120, -20, -75, -25, 500],
            [0, 5],
        ),
        "regional-rollup": (
            "Margin",
            ["East", "West", "Midwest", "Online", "Total"],
            [30, -10, 25, 15, 60],
            [4],
        ),
    }[data_case]
    chart_data.categories = categories
    chart_data.add_series(series_name, values, subtotals=subtotals)
    context.chart_data = chart_data


# when ====================================================


@when("I add the ChartEx waterfall via {add_path}")
def when_I_add_the_ChartEx_waterfall_via_add_path(context, add_path):
    shapes = context.slide.shapes
    if add_path == "add_chart":
        context.graphic_frame = shapes.add_chart(
            XL_CHART_TYPE.WATERFALL,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            context.chart_data,
        )
    else:
        context.graphic_frame = shapes.add_chartex(
            context.chart_data,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
        )
    context.graphic_frame_idx = len(shapes) - 1


@when("I round-trip the presentation for ChartEx inspection")
def when_I_round_trip_the_presentation_for_ChartEx_inspection(context):
    if os.path.isfile(saved_pptx_path):
        os.remove(saved_pptx_path)
    context.prs.save(saved_pptx_path)
    context.saved_package_names = zipfile.ZipFile(saved_pptx_path).namelist()
    context.saved_content_types_xml = (
        zipfile.ZipFile(saved_pptx_path).read("[Content_Types].xml").decode("utf-8")
    )
    context.prs = Presentation(saved_pptx_path)
    context.slide = context.prs.slides[0]
    if hasattr(context, "graphic_frame_idx"):
        context.graphic_frame = context.slide.shapes[context.graphic_frame_idx]
    if hasattr(context, "classic_graphic_frame_idx"):
        context.classic_graphic_frame = context.slide.shapes[context.classic_graphic_frame_idx]
    if hasattr(context, "chartex_graphic_frame_idx"):
        context.chartex_graphic_frame = context.slide.shapes[context.chartex_graphic_frame_idx]


@when("I add a classic chart beside a ChartEx waterfall")
def when_I_add_a_classic_chart_beside_a_ChartEx_waterfall(context):
    chart_data = CategoryChartData()
    chart_data.categories = ["North", "South", "West"]
    chart_data.add_series("Orders", (4, 7, 3))

    context.classic_graphic_frame = context.slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5),
        Inches(0.5),
        Inches(4.5),
        Inches(3.0),
        chart_data,
    )
    context.classic_graphic_frame_idx = len(context.slide.shapes) - 1

    context.chartex_graphic_frame = context.slide.shapes.add_chart(
        XL_CHART_TYPE.WATERFALL,
        Inches(5.2),
        Inches(0.5),
        Inches(4.0),
        Inches(3.0),
        context.chart_data,
    )
    context.chartex_graphic_frame_idx = len(context.slide.shapes) - 1
    context.graphic_frame = context.chartex_graphic_frame


@when("I attempt to add deferred ChartEx type {member_name}")
def when_I_attempt_to_add_deferred_ChartEx_type_member_name(context, member_name):
    try:
        context.chartex_error = None
        context.slide.shapes.add_chart(
            getattr(XL_CHART_TYPE, member_name),
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            context.chart_data,
        )
    except Exception as err:
        context.chartex_error = err


# then ====================================================


@then("the active ChartEx frame exposes ChartEx but not a classic chart")
def then_the_active_ChartEx_frame_exposes_ChartEx_but_not_a_classic_chart(context):
    assert context.graphic_frame.has_chartex is True
    assert context.graphic_frame.has_chart is False
    assert context.graphic_frame.chartex is not None


@then("the active ChartEx chart type is waterfall")
def then_the_active_ChartEx_chart_type_is_waterfall(context):
    assert context.graphic_frame.chartex.chart_type == "waterfall"


@then("the active ChartEx series is named {expected_name}")
def then_the_active_ChartEx_series_is_named_expected_name(context, expected_name):
    series = context.graphic_frame.chartex.series
    assert len(series) == 1
    assert series[0].name == expected_name


@then("the active ChartEx series values are {expected_values}")
def then_the_active_ChartEx_series_values_are_expected_values(context, expected_values):
    series = context.graphic_frame.chartex.series
    actual_values = series[0].values
    assert actual_values == _float_values(expected_values)


@then("the active ChartEx category labels are {expected_labels}")
def then_the_active_ChartEx_category_labels_are_expected_labels(context, expected_labels):
    actual_labels = _chart_labels(context.graphic_frame.chartex)
    assert actual_labels == _csv_values(expected_labels)


@then("the active ChartEx subtotal indices are {expected_indices}")
def then_the_active_ChartEx_subtotal_indices_are_expected_indices(context, expected_indices):
    actual_indices = _subtotal_indices(context.graphic_frame.chartex)
    assert actual_indices == _int_values(expected_indices)


@then("the active ChartEx part content type is the ChartEx content type")
def then_the_active_ChartEx_part_content_type_is_the_ChartEx_content_type(context):
    assert context.graphic_frame.chartex_part.content_type == CT.OFC_CHART_EX


@then("the active ChartEx partname contains chartEx")
def then_the_active_ChartEx_partname_contains_chartEx(context):
    assert "chartEx" in str(context.graphic_frame.chartex_part.partname)


@then("the slide has one classic chart frame and one ChartEx frame")
def then_the_slide_has_one_classic_chart_frame_and_one_ChartEx_frame(context):
    classic_frames = [shape for shape in context.slide.shapes if getattr(shape, "has_chart", False)]
    chartex_frames = [
        shape for shape in context.slide.shapes if getattr(shape, "has_chartex", False)
    ]
    assert len(classic_frames) == 1
    assert len(chartex_frames) == 1


@then("the classic chart frame still exposes only a classic chart")
def then_the_classic_chart_frame_still_exposes_only_a_classic_chart(context):
    assert context.classic_graphic_frame.has_chart is True
    assert context.classic_graphic_frame.has_chartex is False
    assert context.classic_graphic_frame.chart is not None


@then("the ChartEx frame still exposes only a ChartEx chart")
def then_the_ChartEx_frame_still_exposes_only_a_ChartEx_chart(context):
    assert context.chartex_graphic_frame.has_chartex is True
    assert context.chartex_graphic_frame.has_chart is False
    assert context.chartex_graphic_frame.chartex is not None


@then("the slide has no ChartEx graphic frames")
def then_the_slide_has_no_ChartEx_graphic_frames(context):
    chartex_frames = [
        shape for shape in context.slide.shapes if getattr(shape, "has_chartex", False)
    ]
    assert chartex_frames == []


@then("the saved package contains no ChartEx partnames")
def then_the_saved_package_contains_no_ChartEx_partnames(context):
    assert not any("chartEx" in name for name in context.saved_package_names)


@then("the saved package contains no ChartEx content type declaration")
def then_the_saved_package_contains_no_ChartEx_content_type_declaration(context):
    assert CT.OFC_CHART_EX not in context.saved_content_types_xml


@then("adding deferred ChartEx type {member_name} raises NotImplementedError")
def then_adding_deferred_ChartEx_type_member_name_raises_NotImplementedError(context, member_name):
    assert isinstance(context.chartex_error, NotImplementedError)
    assert member_name in str(context.chartex_error)


@then("XL_CHART_TYPE.{member_name} exists with value {expected_value}")
def then_XL_CHART_TYPE_member_name_exists_with_value_expected_value(
    context, member_name, expected_value
):
    assert hasattr(XL_CHART_TYPE, member_name)
    member = getattr(XL_CHART_TYPE, member_name)
    actual_value = member.value
    assert actual_value == int(expected_value)
    assert actual_value >= 1000


def _chart_labels(chartex):
    cat_dim = chartex._element.chartData.data_lst[0].strDim_lst[0]
    levels = cat_dim.lvl_lst
    assert len(levels) == 1
    return [pt.text for pt in levels[0]]


def _csv_values(csv_text):
    return [item.strip() for item in csv_text.split(",")]


def _float_values(csv_text):
    return [float(item.strip()) for item in csv_text.split(",")]


def _int_values(csv_text):
    if csv_text == "none":
        return []
    return [int(item.strip()) for item in csv_text.split(",")]


def _subtotal_indices(chartex):
    series = chartex._element.chart.plotArea.plotAreaRegion.series_lst[0]
    subtotal_elems = series.findall(f".//{qn('cx:subtotals')}/{qn('cx:idx')}")
    return [int(elem.get("val")) for elem in subtotal_elems]
