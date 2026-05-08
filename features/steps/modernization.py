"""Gherkin step implementations for Modernization Phase 1 (issue #29)."""

from __future__ import annotations

from pathlib import Path

from behave import given, then, when
from environment import scratch_dir  # noqa: E402

from pptx import Presentation
from pptx.enum.dml import MSO_PATTERN_TYPE


# given ===================================================


@given("a freshly-saved presentation at a Path")
def given_freshly_saved_presentation(context):
    target = Path(scratch_dir) / "modernization_seed.pptx"
    target.parent.mkdir(parents=True, exist_ok=True)
    Presentation().save(target)
    context.path = target


@given("a fresh presentation")
def given_fresh_presentation(context):
    context.prs = Presentation()


@given("a fresh slide on a fresh presentation")
def given_fresh_slide(context):
    prs = Presentation()
    context.slide = prs.slides.add_slide(prs.slide_layouts[6])


# when ====================================================


@when("I call Presentation(path) with the Path")
def when_call_presentation_with_path(context):
    context.prs = Presentation(context.path)


@when("I save it to a Path")
def when_save_to_path(context):
    target = Path(scratch_dir) / "modernization_out.pptx"
    target.parent.mkdir(parents=True, exist_ok=True)
    context.prs.save(target)
    context.saved_path = target


# then ====================================================


@then("I get a presentation back")
def then_presentation_back(context):
    assert context.prs is not None


@then("a non-empty .pptx file exists at that Path")
def then_non_empty_file_exists(context):
    assert context.saved_path.exists()
    assert context.saved_path.stat().st_size > 0


@then("MSO_PATTERN_TYPE.PERCENT_40 exists with xml_value pct40")
def then_percent_40_correct(context):
    assert MSO_PATTERN_TYPE.PERCENT_40.xml_value == "pct40"


@then("the broken name ERCENT_40 does not exist")
def then_ercent_40_absent(context):
    assert hasattr(MSO_PATTERN_TYPE, "ERCENT_40") is False


@then("slide.background.element local-name is bg")
def then_background_element_is_bg(context):
    from lxml import etree

    bg_elm = context.slide.background.element
    local = etree.QName(bg_elm.tag).localname
    assert local == "bg", f"expected 'bg', got '{local}'"
