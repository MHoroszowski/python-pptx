"""Step implementations for features/iss-16-advanced-text.feature (issue #16).

Self-contained: every scenario builds an in-memory blank presentation.
"""

import io
import os

from behave import given, then, when

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_TEXT_DIRECTION, MSO_TEXT_STRIKE_TYPE
from pptx.util import Inches, Pt

TEST_FONT = os.path.join(
    os.path.dirname(os.path.dirname(__file__)), "..", "tests", "test_files", "calibriz.ttf"
)
TEST_FONT = os.path.abspath(TEST_FONT)


def _blank_run(context):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2)).text_frame
    r = tf.paragraphs[0].add_run()
    r.text = "Sample"
    context.prs = prs
    context.tf = tf
    context.run = r


def _roundtrip(context):
    buf = io.BytesIO()
    context.prs.save(buf)
    buf.seek(0)
    context.prs2 = Presentation(buf)
    context.tf2 = list(context.prs2.slides[0].shapes)[0].text_frame
    return context.tf2


@given("a blank slide text frame with one run")
def given_blank_run(context):
    _blank_run(context)


@given("a tiny text frame stuffed with text")
def given_tiny_stuffed(context):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(0.3)).text_frame
    tf.text = "Supercalifragilistic " * 14
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(18)
    context.prs = prs
    context.tf = tf


@given("a tiny text frame with one very long word")
def given_tiny_longword(context):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(0), Inches(0), Inches(0.5), Inches(0.5)).text_frame
    tf.text = "Supercalifragilisticexpialidocious"
    context.prs = prs
    context.tf = tf


@when("I set the run superscript")
def when_superscript(context):
    context.run.font.superscript = True


@then("the run reports superscript true")
def then_superscript(context):
    assert context.run.font.superscript is True


@when("I set the run strike to double")
def when_strike_double(context):
    context.run.font.strike = MSO_TEXT_STRIKE_TYPE.DOUBLE


@then("the run reports strike double after round-trip")
def then_strike_double(context):
    f2 = _roundtrip(context).paragraphs[0].runs[0].font
    assert f2.strike == MSO_TEXT_STRIKE_TYPE.DOUBLE


@when("I set the run highlight to FFFF00")
def when_highlight(context):
    context.run.font.highlight.rgb = RGBColor(0xFF, 0xFF, 0x00)


@then("the run reports highlight FFFF00 after round-trip")
def then_highlight(context):
    f2 = _roundtrip(context).paragraphs[0].runs[0].font
    assert f2.highlight.rgb == RGBColor(0xFF, 0xFF, 0x00)


@when("I set the run character spacing to 2 points")
def when_spacing(context):
    context.run.font.character_spacing = Pt(2)


@then("the run reports character spacing 2 points")
def then_spacing(context):
    assert context.run.font.character_spacing.pt == 2.0


@when("I set east_asian to MS Gothic and name to Calibri")
def when_trio(context):
    context.run.font.east_asian = "MS Gothic"
    context.run.font.name = "Calibri"


@then("latin is Calibri and east_asian is MS Gothic and they are independent")
def then_trio(context):
    f2 = _roundtrip(context).paragraphs[0].runs[0].font
    assert f2.name == "Calibri" and f2.east_asian == "MS Gothic"
    assert f2.latin == "Calibri"


@when("I set the text frame to 2 columns spaced 36 points")
def when_columns(context):
    context.tf.columns = 2
    context.tf.column_spacing = Pt(36)


@then("the text frame reports 2 columns after round-trip")
def then_columns(context):
    assert _roundtrip(context).columns == 2


@when("I set the text direction to east asian vertical")
def when_direction(context):
    context.tf.text_direction = MSO_TEXT_DIRECTION.EAST_ASIAN_VERTICAL


@then("the text frame reports east asian vertical after round-trip")
def then_direction(context):
    assert _roundtrip(context).text_direction == MSO_TEXT_DIRECTION.EAST_ASIAN_VERTICAL


@when("I set the paragraph to Arabic right-to-left")
def when_rtl(context):
    p = context.tf.paragraphs[0]
    p.text = "اللغة العربية"
    p.rtl = True


@then("the paragraph reports rtl true after round-trip")
def then_rtl(context):
    p2 = _roundtrip(context).paragraphs[0]
    assert p2.rtl is True


@then("will_overflow reports true")
def then_will_overflow(context):
    assert context.tf.will_overflow(font_file=TEST_FONT) is True


@when("I call fit_text on it")
def when_fit_text(context):
    context.tf.fit_text(font_file=TEST_FONT)


@then("no error is raised and auto_size is set")
def then_fit_ok(context):
    assert context.tf.auto_size is not None


@when("I call shrink_text_to_fit")
def when_shrink(context):
    context.tf.shrink_text_to_fit(font_file=TEST_FONT)


@then("the normAutofit fontScale is below 100")
def then_shrink(context):
    na = context.tf._txBody.bodyPr.normAutofit
    assert na is not None and na.fontScale < 100
    assert context.tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
