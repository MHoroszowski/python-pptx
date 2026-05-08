"""Gherkin step implementations for Accessibility Phase B (issue #22)."""

from __future__ import annotations

from behave import given, then, when

from pptx import Presentation
from pptx.util import Inches


# given ===================================================


@given("a slide with one textbox")
def given_one_textbox(context):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    context.shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    context.prs = prs
    context.slide = slide


@given("a slide with three textboxes labelled A, B, C")
def given_three_textboxes(context):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    a = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.5))
    b = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(2), Inches(0.5))
    c = slide.shapes.add_textbox(Inches(5.5), Inches(0.5), Inches(2), Inches(0.5))
    a.name = "A"
    b.name = "B"
    c.name = "C"
    context.prs = prs
    context.slide = slide
    context.shape_a = a
    context.shape_b = b
    context.shape_c = c


# when ====================================================


@when("I set shape.is_hidden_from_accessibility to True")
def when_set_hidden_from_accessibility(context):
    context.shape.is_hidden_from_accessibility = True


@when("I set reading_order to (C, A, B)")
def when_set_reading_order(context):
    context.slide.shapes.reading_order = (context.shape_c, context.shape_a, context.shape_b)


@when("I tag A with alt text and B as decorative")
def when_tag_a_alt_b_deco(context):
    context.shape_a.alt_text = "Alpha"
    context.shape_b.is_decorative = True


# then ====================================================


@then("shape.is_decorative is True")
def then_shape_is_decorative_true(context):
    assert context.shape.is_decorative is True


@then("reading_order produces shapes in the order A, B, C")
def then_reading_order_a_b_c(context):
    names = tuple(s.name for s in context.slide.shapes.reading_order)
    assert names == ("A", "B", "C"), f"expected (A, B, C), got {names}"


@then("iteration produces shapes in the order C, A, B")
def then_iteration_c_a_b(context):
    names = tuple(s.name for s in context.slide.shapes)
    assert names == ("C", "A", "B"), f"expected (C, A, B), got {names}"


@then("accessibility_issues returns just C")
def then_accessibility_issues_just_c(context):
    issues = context.slide.shapes.accessibility_issues()
    assert len(issues) == 1, f"expected 1 issue, got {len(issues)}"
    assert issues[0].name == "C"
