"""Step implementations for features/iss-21-hyperlinks-clickactions.feature.

Self-contained: every scenario builds an in-memory blank presentation,
round-trips it through a BytesIO save/reopen, and asserts the reopened state.
"""

import io
import os
import struct
import tempfile
import wave

from behave import given, then, when

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.action import PP_ACTION
from pptx.util import Inches

TEST_IMAGE = os.path.abspath(
    os.path.join(
        os.path.dirname(os.path.dirname(__file__)),
        "..",
        "tests",
        "test_files",
        "monty-truth.png",
    )
)


def _wav():
    fd, path = tempfile.mkstemp(suffix=".wav")
    os.close(fd)
    w = wave.open(path, "w")
    w.setnchannels(1)
    w.setsampwidth(2)
    w.setframerate(8000)
    w.writeframes(struct.pack("<h", 0) * 800)
    w.close()
    return path


def _new_prs(context):
    context.prs = Presentation()
    context.slide = context.prs.slides.add_slide(context.prs.slide_layouts[6])


def _new_run(context):
    tb = context.slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    )
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "x"
    return run


def _reopen(context):
    buf = io.BytesIO()
    context.prs.save(buf)
    buf.seek(0)
    context.prs2 = Presentation(buf)
    context.shape2 = context.prs2.slides[0].shapes[-1]


def _first_run(context):
    return context.shape2.text_frame.paragraphs[0].runs[0]


@given("a run with an external hyperlink")
def given_run_with_hyperlink(context):
    _new_prs(context)
    context.run = _new_run(context)
    context.run.hyperlink.address = "https://example.com"


@given("a run with no hyperlink")
def given_run_no_hyperlink(context):
    _new_prs(context)
    context.run = _new_run(context)


@given("a shape")
def given_a_shape(context):
    _new_prs(context)
    sp = context.slide.shapes.add_textbox(
        Inches(1), Inches(3), Inches(3), Inches(1)
    )
    sp.text_frame.text = "btn"
    context.shape = sp


@given("a picture on a slide")
def given_a_picture(context):
    _new_prs(context)
    context.picture = context.slide.shapes.add_picture(
        TEST_IMAGE, Inches(1), Inches(1), Inches(1), Inches(1)
    )


@when('I set the run hyperlink tooltip to "{text}"')
def when_set_run_tooltip(context, text):
    context.run.hyperlink.tooltip = text


@when('I set the shape alt-text to "{text}"')
def when_set_shape_alt_text(context, text):
    context.shape.alt_text = text


@then('the reopened shape alt-text is "{text}"')
def then_shape_alt_text(context, text):
    _reopen(context)
    assert context.shape2.alt_text == text


@when("I set the run hyperlink color to {hexval}")
def when_set_run_color(context, hexval):
    context.run.hyperlink.color.rgb = RGBColor.from_string(hexval)


@when('I set the shape click action to run macro "{name}"')
def when_set_run_macro(context, name):
    context.shape.click_action.run_macro(name)


@when("I attach a click sound to the shape")
def when_attach_sound(context):
    context.shape.click_action.play_sound(_wav())


@when('I set the shape hover action address to "{url}"')
def when_set_hover(context, url):
    context.shape.hover_action.hyperlink.address = url


@when("I make the run jump to a new slide")
def when_run_jump(context):
    target = context.prs.slides.add_slide(context.prs.slide_layouts[6])
    context.run.click_action.target_slide = target


@when('I set the picture hyperlink address to "{url}"')
def when_set_picture_hyperlink(context, url):
    context.picture.hyperlink.address = url


@then('the reopened run hyperlink tooltip is "{text}"')
def then_run_tooltip(context, text):
    _reopen(context)
    assert _first_run(context).hyperlink.tooltip == text


@then("the reopened run hyperlink address is None")
def then_run_address_none(context):
    assert _first_run(context).hyperlink.address is None


@then("the reopened run hyperlink color is {hexval}")
def then_run_color(context, hexval):
    _reopen(context)
    assert _first_run(context).hyperlink.color.rgb == RGBColor.from_string(hexval)


@then("the reopened shape click action is RUN_MACRO")
def then_shape_macro(context):
    _reopen(context)
    assert context.shape2.click_action.action == PP_ACTION.RUN_MACRO


@then("the reopened shape click action has an embedded sound")
def then_shape_sound(context):
    _reopen(context)
    hlink = context.shape2.click_action._hlink
    assert hlink is not None and hlink.snd is not None
    assert hlink.snd.embed is not None


@then('the reopened shape hover action address is "{url}"')
def then_shape_hover(context, url):
    _reopen(context)
    assert context.shape2.hover_action.hyperlink.address == url


@then("the reopened run click action is NAMED_SLIDE")
def then_run_jump(context):
    _reopen(context)
    assert _first_run(context).click_action.action == PP_ACTION.NAMED_SLIDE


@then('the reopened picture hyperlink address is "{url}"')
def then_picture_hyperlink(context, url):
    _reopen(context)
    pic = next(sh for sh in context.prs2.slides[0].shapes if sh.shape_type == 13)
    assert pic.hyperlink.address == url
