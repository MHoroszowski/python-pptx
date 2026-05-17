"""Gherkin step implementations for slide collection-related features."""

from __future__ import annotations

from behave import given, then, when
from helpers import test_pptx

from pptx import Presentation

# given ===================================================


@given("a SlideLayouts object containing 2 layouts as slide_layouts")
def given_a_SlideLayouts_object_containing_2_layouts(context):
    prs = Presentation(test_pptx("mst-slide-layouts"))
    context.slide_layouts = prs.slide_master.slide_layouts


@given("a SlideMasters object containing 2 masters")
def given_a_SlideMasters_object_containing_2_masters(context):
    prs = Presentation(test_pptx("prs-slide-masters"))
    context.slide_masters = prs.slide_masters


@given("a Slides object containing 3 slides")
def given_a_Slides_object_containing_3_slides(context):
    prs = Presentation(test_pptx("sld-slides"))
    context.prs = prs
    context.slides = prs.slides
    # ---capture original slide ids for CRUD ordering assertions---
    context.original_slide_ids = [s.slide_id for s in prs.slides]


# when ====================================================


@when("I call slides.add_slide()")
def when_I_call_slides_add_slide(context):
    context.slide_layout = context.prs.slide_masters[0].slide_layouts[0]
    context.slides.add_slide(context.slide_layout)


@when("I call slide_layouts.remove(slide_layouts[1])")
def when_I_call_slide_layouts_remove(context):
    slide_layouts = context.slide_layouts
    slide_layouts.remove(slide_layouts[1])


@when("I call slides.add_slide(slide_layout, index=0)")
def when_I_call_slides_add_slide_index_0(context):
    layout = context.prs.slide_masters[0].slide_layouts[0]
    context.new_slide = context.slides.add_slide(layout, index=0)


@when("I call slides.add_slide(slide_layout, index=2)")
def when_I_call_slides_add_slide_index_2(context):
    layout = context.prs.slide_masters[0].slide_layouts[0]
    context.new_slide = context.slides.add_slide(layout, index=2)


@when("I call slides.move(slides[0], 2)")
def when_I_call_slides_move(context):
    context.slides.move(context.slides[0], 2)


@when("I call slides.remove(slides[1])")
def when_I_call_slides_remove(context):
    context.slides.remove(context.slides[1])


@when("I call slides[1].delete()")
def when_I_call_slide_delete(context):
    context.slides[1].delete()


@when("I call slides.duplicate(slides[0])")
def when_I_call_slides_duplicate_default_index(context):
    context.new_slide = context.slides.duplicate(context.slides[0])


@when("I call slides.duplicate(slides[0], index=3)")
def when_I_call_slides_duplicate_index_3(context):
    context.new_slide = context.slides.duplicate(context.slides[0], index=3)


@when("I call slides[1].duplicate()")
def when_I_call_slide_duplicate_alias(context):
    context.new_slide = context.slides[1].duplicate()


# then ====================================================


@then("iterating produces 3 NotesSlidePlaceholder objects")
def then_iterating_produces_3_NotesSlidePlaceholder_objects(context):
    idx = -1
    for idx, placeholder in enumerate(context.notes_slide.placeholders):
        typename = type(placeholder).__name__
        assert typename == "NotesSlidePlaceholder", "got %s" % typename
    assert idx == 2


@then("iterating slide_layouts produces 2 SlideLayout objects")
def then_iterating_slide_layouts_produces_2_SlideLayout_objects(context):
    slide_layouts = context.slide_layouts
    idx = -1
    for idx, slide_layout in enumerate(slide_layouts):
        assert type(slide_layout).__name__ == "SlideLayout"
    assert idx == 1


@then("iterating slide_masters produces 2 SlideMaster objects")
def then_iterating_slide_masters_produces_2_SlideMaster_objects(context):
    slide_masters = context.slide_masters
    idx = -1
    for idx, slide_master in enumerate(slide_masters):
        assert type(slide_master).__name__ == "SlideMaster"
    assert idx == 1


@then("iterating slides produces 3 Slide objects")
def then_iterating_slides_produces_3_Slide_objects(context):
    slides = context.slides
    idx = -1
    for idx, slide in enumerate(slides):
        assert type(slide).__name__ == "Slide"
    assert idx == 2


@then("len(slides) is {count}")
def then_len_slides_is_count(context, count):
    slides = context.slides
    assert len(slides) == int(count)


@then("len(slide_layouts) is {n}")
def then_len_slide_layouts_is_2(context, n):
    assert len(context.slide_layouts) == int(n)


@then("len(slide_masters) is 2")
def then_len_slide_masters_is_2(context):
    slide_masters = context.slide_masters
    assert len(slide_masters) == 2


@then("slide_layouts[1] is a SlideLayout object")
def then_slide_layouts_1_is_a_SlideLayout_object(context):
    slide_layouts = context.slide_layouts
    assert type(slide_layouts[1]).__name__ == "SlideLayout"


@then("slide_layouts.get_by_name(slide_layouts[1].name) is slide_layouts[1]")
def then_slide_layouts_get_by_name_is_slide_layout(context):
    slide_layouts = context.slide_layouts
    assert slide_layouts.get_by_name(slide_layouts[1].name) is slide_layouts[1]


@then("slide_layouts.index(slide_layouts[1]) == 1")
def then_slide_layouts_index_is_1(context):
    slide_layouts = context.slide_layouts
    assert slide_layouts.index(slide_layouts[1]) == 1


@then("slide_masters[1] is a SlideMaster object")
def then_slide_masters_1_is_a_SlideMaster_object(context):
    slide_masters = context.slide_masters
    assert type(slide_masters[1]).__name__ == "SlideMaster"


@then("slides.get(256) is slides[0]")
def then_slides_get_256_is_slides_0(context):
    slides = context.slides
    assert slides.get(256) is slides[0]


@then("slides.get(666, default=slides[2]) is slides[2]")
def then_slides_get_666_default_slides_2_is_slides_2(context):
    slides = context.slides
    assert slides.get(666, default=slides[2]) is slides[2]


@then("slides[2] is a Slide object")
def then_slides_2_is_a_Slide_object(context):
    slides = context.slides
    assert type(slides[2]).__name__ == "Slide"


@then("the new slide is at index {idx:d}")
def then_the_new_slide_is_at_index(context, idx):
    assert context.slides[idx].slide_id == context.new_slide.slide_id, (
        "expected new slide at index %d, got slide_id mismatch" % idx
    )


@then("the slide order matches the original [1, 2, 0]")
def then_slide_order_matches_1_2_0(context):
    o = context.original_slide_ids
    expected = [o[1], o[2], o[0]]
    actual = [s.slide_id for s in context.slides]
    assert actual == expected, "expected %r, got %r" % (expected, actual)


@then("the surviving slide order matches the original [0, 2]")
def then_surviving_slide_order_matches_0_2(context):
    o = context.original_slide_ids
    expected = [o[0], o[2]]
    actual = [s.slide_id for s in context.slides]
    assert actual == expected, "expected %r, got %r" % (expected, actual)


@then("the duplicate is at index {idx:d}")
def then_the_duplicate_is_at_index(context, idx):
    assert context.slides[idx].slide_id == context.new_slide.slide_id, (
        "expected duplicate at index %d, got slide_id mismatch" % idx
    )


@then("the source slide is still at index 0")
def then_source_slide_still_at_index_0(context):
    assert context.slides[0].slide_id == context.original_slide_ids[0], (
        "source slide moved off index 0"
    )


@then("the duplicate slide_id is unique")
def then_duplicate_slide_id_is_unique(context):
    assert context.new_slide.slide_id not in context.original_slide_ids, (
        "duplicate slide_id collides with an existing slide"
    )


@then("calling slides.duplicate(slides[0], index=99) raises IndexError")
def then_duplicate_index_99_raises(context):
    import pytest

    with pytest.raises(IndexError):
        context.slides.duplicate(context.slides[0], index=99)


# -- append_from steps -----------------------------------------------------


@given("two seeded presentations source and target")
def given_two_seeded_presentations(context):
    from pptx.util import Inches

    src = Presentation()
    layout = src.slide_layouts[6]
    for i in range(3):
        s = src.slides.add_slide(layout)
        s.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.text = (
            "src %d" % i
        )
    tgt = Presentation()
    tgt.slides.add_slide(tgt.slide_layouts[6])
    context.source_pres = src
    context.target_pres = tgt
    context.target_slides_before = len(tgt.slides)
    context.target_masters_before = sum(1 for _ in tgt.slide_masters)


@when("I call target.append_from(source)")
def when_target_append_from_source_all(context):
    context.new_slides = context.target_pres.append_from(context.source_pres)


@when("I call target.append_from(source, slide_indexes=[0, 1])")
def when_target_append_from_source_indexes_0_1(context):
    context.new_slides = context.target_pres.append_from(context.source_pres, slide_indexes=[0, 1])


@when("I call target.append_from(source, slide_indexes=[])")
def when_target_append_from_source_indexes_empty(context):
    context.new_slides = context.target_pres.append_from(context.source_pres, slide_indexes=[])


@then("target's slide count grew by source's slide count")
def then_target_grew_by_source_slide_count(context):
    expected = context.target_slides_before + len(context.source_pres.slides)
    actual = len(context.target_pres.slides)
    assert actual == expected, "expected %d slides, got %d" % (expected, actual)


@then("target's master count grew by 1")
def then_target_master_count_grew_by_1(context):
    actual = sum(1 for _ in context.target_pres.slide_masters)
    assert actual == context.target_masters_before + 1, "expected %d masters, got %d" % (
        context.target_masters_before + 1,
        actual,
    )


@then("target's slide count grew by 2")
def then_target_slide_count_grew_by_2(context):
    actual = len(context.target_pres.slides)
    assert actual == context.target_slides_before + 2


@then("target's slide count is unchanged")
def then_target_slide_count_unchanged(context):
    actual = len(context.target_pres.slides)
    assert actual == context.target_slides_before


@then("calling target.append_from(source, slide_indexes=[99]) raises IndexError")
def then_append_from_index_99_raises(context):
    import pytest

    with pytest.raises(IndexError):
        context.target_pres.append_from(context.source_pres, slide_indexes=[99])


# Sections (Phase 4) =======================================


@then("len(prs.sections) is {n:d}")
def then_len_prs_sections_is_n(context, n):
    assert len(context.prs.sections) == n


@when('I call prs.sections.add_section("{name}")')
def when_prs_sections_add_section(context, name):
    context.section = context.prs.sections.add_section(name)


@then('prs.sections[0].name is "{expected}"')
def then_prs_sections_0_name_is(context, expected):
    assert context.prs.sections[0].name == expected


@when("I call section.add_slide(prs.slides[0])")
def when_section_add_slide_0(context):
    context.tracked_slide_id = context.prs.slides[0].slide_id
    context.section.add_slide(context.prs.slides[0])


@then("len(section.slides) is {n:d}")
def then_len_section_slides_is_n(context, n):
    assert len(context.section.slides) == n


@then("section.slides still contains the moved slide")
def then_section_still_contains_moved_slide(context):
    section_slide_ids = [s.slide_id for s in context.section.slides]
    assert context.tracked_slide_id in section_slide_ids, (
        "expected slide_id %r in section.slides %r" % (context.tracked_slide_id, section_slide_ids)
    )


@then("the moved slide is at presentation index {idx:d}")
def then_moved_slide_at_index(context, idx):
    assert context.prs.slides[idx].slide_id == context.tracked_slide_id


@when("I call prs.sections.remove(section)")
def when_prs_sections_remove(context):
    context.prs.sections.remove(context.section)


@then("prs._element.extLst is None")
def then_prs_element_extLst_is_None(context):
    assert context.prs._element.extLst is None
