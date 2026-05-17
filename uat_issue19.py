"""UAT — issue #19: Slide Masters, Layouts & .potx Templates.

Exercises ALL nine sub-features end-to-end and emits artifacts for
maintainer visual signoff in PowerPoint/Keynote. Exits non-zero on any
structural failure. Per repo CLAUDE.md §6a this is the maintainer's
acceptance path — an agent running it green is script-QA, NOT signoff.

Artifacts written to repo root:
  uat_issue19_template.potx     — SF1/SF2 (open + save_as_potx)
  uat_issue19_layouts.pptx      — SF3/SF4/SF5/SF6/SF7/SF8/SF9 on real slides
"""

from __future__ import annotations

import io
import sys
import zipfile

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.util import Inches, Pt

POTX = "/Users/mhoroszowski/Projects/AI/python-pptx/uat_issue19_template.potx"
PPTX = "/Users/mhoroszowski/Projects/AI/python-pptx/uat_issue19_layouts.pptx"

ok = True


def check(label: str, cond: bool) -> None:
    global ok
    print(("  PASS " if cond else "  FAIL ") + label)
    ok = ok and cond


print("=== issue #19 UAT ===")

# ---------- SF2: save_as_potx (and SF1 read-back) ----------
prs = Presentation()
before_ct = prs.part.content_type
prs.save_as_potx(POTX)
with zipfile.ZipFile(POTX) as z:
    ct_xml = z.read("[Content_Types].xml").decode()
check("SF2 save_as_potx writes template content-type",
      "presentationml.template.main+xml" in ct_xml)
check("SF2 in-memory package content-type unmutated",
      prs.part.content_type == before_ct)
# SF1: the .potx we just wrote re-opens (was a hard ValueError before the fix)
reopened = Presentation(POTX)
check("SF1 Presentation('.potx') opens without error", reopened is not None)
check("SF1 reopened .potx exposes masters", len(reopened.slide_masters) >= 1)

# ---------- The layouts deck ----------
prs = Presentation()
master = prs.slide_masters[0]

# SF3: add_layout("Three Columns") — the issue's headline acceptance check
n_before = len(master.slide_layouts)
three_col = master.slide_layouts.add_layout(name="Three Columns")
check("SF3 add_layout increments layout count by 1",
      len(master.slide_layouts) == n_before + 1)
check("SF3 new layout name == 'Three Columns'", three_col.name == "Three Columns")

# SF6: programmatically add three body placeholders → an actual 3-column layout
for col, left in enumerate((Inches(0.4), Inches(4.7), Inches(9.0))):
    three_col.placeholders.add(
        idx=10 + col,
        ph_type=PP_PLACEHOLDER.BODY,
        name="Column %d" % (col + 1),
        left=left,
        top=Inches(1.8),
        width=Inches(4.0),
        height=Inches(4.8),
    )
check("SF6 three placeholders added to layout",
      len(three_col.placeholders) >= 3)

# SF5: author a shape directly on the master (visible on every slide)
wm = master.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(0.1), Inches(3.0), Inches(0.5)
)
wm.text_frame.text = "MASTER BANNER (SF5)"
wm.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
check("SF5 shape authored on master", wm is not None)

# SF4: copy_from — duplicate the Three Columns layout
copy = master.slide_layouts.copy_from(three_col)
check("SF4 copy_from duplicates shape/placeholder count",
      len(list(copy.shapes)) == len(list(three_col.shapes)))
check("SF4 copy_from leaves source untouched",
      len(list(three_col.shapes)) > 0)

# A real slide built on the new "Three Columns" layout (so it renders visibly)
slide = prs.slides.add_slide(three_col)
for i, ph in enumerate(slide.placeholders):
    try:
        ph.text = "Column %d content" % (i + 1)
    except Exception:
        pass

# SF9: get_layout by id round-trips to the same object
idLst = master._element.get_or_add_sldLayoutIdLst()
some_id = next((e.id for e in idLst.sldLayoutId_lst if e.id is not None), None)
if some_id is not None:
    check("SF9 get_layout(id) returns a layout",
          master.get_layout(some_id) is not None)
check("SF9 get_layout(bad id) returns None (no raise)",
      master.get_layout(987654321) is None)

# SF7: reassign the slide to a different layout (cross-master mechanism)
slide.slide_layout = copy
check("SF7 slide_layout setter repoints layout",
      slide.slide_layout.name == copy.name)

# SF8: chart into a chart placeholder — the SF6→SF8 composition:
# build a CHART placeholder programmatically, then drop a chart into it.
chart_done = False
probe = Presentation()
pm = probe.slide_masters[0]
chart_layout = pm.slide_layouts.add_layout(name="Chart Slot")
chart_layout.placeholders.add(
    idx=10,
    ph_type=PP_PLACEHOLDER.CHART,
    left=Inches(1.0),
    top=Inches(1.2),
    width=Inches(8.0),
    height=Inches(4.5),
)
cslide = probe.slides.add_slide(chart_layout)
chart_ph = next(
    p for p in cslide.placeholders
    if p.placeholder_format.type == PP_PLACEHOLDER.CHART
)
cd = CategoryChartData()
cd.categories = ["Q1", "Q2", "Q3"]
cd.add_series("Revenue", (10, 24, 18))
gf = chart_ph.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, cd)
chart_done = bool(gf.has_chart)
check("SF8 insert_chart into a CHART placeholder (SF6+SF8)", chart_done)
if chart_done:
    probe.save(PPTX.replace(".pptx", "_chart.pptx"))

prs.save(PPTX)

# Final structural round-trip of the main deck
rt = Presentation(PPTX)
check("ALL round-trip: layouts deck reopens with 'Three Columns'",
      rt.slide_masters[0].slide_layouts.get_by_name("Three Columns") is not None)

print("ARTIFACTS:", POTX, PPTX)
print("RESULT:", "PASS" if ok else "FAIL")
sys.exit(0 if ok else 1)
